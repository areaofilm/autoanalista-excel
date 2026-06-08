import hashlib
import gzip
import html
import hmac
import io
import json
import os
import re
import sqlite3
import unicodedata
import urllib.error
import urllib.request
import zipfile
from datetime import datetime, timedelta
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
from reportlab.graphics.charts.barcharts import HorizontalBarChart, VerticalBarChart
from reportlab.graphics.charts.piecharts import Pie
from reportlab.graphics.shapes import Drawing, String
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except Exception:  # pragma: no cover - optional until dependency is installed
    Presentation = None
    RGBColor = None
    PP_ALIGN = None
    Inches = None
    Pt = None
from sklearn.cluster import KMeans
from sklearn.compose import ColumnTransformer
from sklearn.decomposition import PCA
from sklearn.ensemble import GradientBoostingClassifier, GradientBoostingRegressor, IsolationForest, RandomForestClassifier, RandomForestRegressor
from sklearn.impute import SimpleImputer
from sklearn.inspection import permutation_importance
from sklearn.linear_model import LinearRegression, LogisticRegression
from sklearn.metrics import (
    accuracy_score,
    confusion_matrix,
    f1_score,
    mean_absolute_error,
    mean_squared_error,
    r2_score,
    silhouette_score,
)
from sklearn.model_selection import cross_validate, train_test_split
from sklearn.pipeline import Pipeline
from sklearn.preprocessing import OneHotEncoder, StandardScaler

SEED = 42
SUPPORTED_TYPES = ["xlsx", "xls", "xlsm", "xlsb", "ods", "csv", "tsv", "txt"]
MAX_ANALYSIS_ROWS = 120000
MAX_ML_ROWS = 30000
MAX_EXPORT_ROWS = 20000
MAX_ML_CATEGORICAL_LEVELS = 80
AUTH_SESSION_TIMEOUT_MINUTES = 30
AUTH_MAX_ATTEMPTS = 5
AUTH_LOCK_MINUTES = 10
DEFAULT_OUTLIER_IQR_MULTIPLIER = 1.5
DEFAULT_MAX_OUTLIER_DROP_PCT = 0.20
DEFAULT_MIN_ROWS_AFTER_TREATMENT = 40
CREATOR_SIGNATURE = "Walace.gorino"
ML_MISSING_CATEGORY = "__NULO__"
ML_OTHER_CATEGORY = "__OUTROS__"
ID_NAME_TOKENS = {
    "id",
    "ids",
    "cod",
    "codigo",
    "chamado",
    "ticket",
    "protocolo",
    "contrato",
    "cpf",
    "cnpj",
    "os",
    "ordem",
    "matricula",
    "registro",
    "numero",
    "nro",
    "serial",
}
ID_NAME_PREFIXES = ("id_", "cod_", "cod", "codigo_", "chamado", "ticket", "protocolo", "cpf", "cnpj")

ROLE_PERMISSIONS = {
    "admin": {"ml": True, "export": True, "rules": True},
    "analyst": {"ml": True, "export": True, "rules": True},
    "viewer": {"ml": False, "export": False, "rules": False},
}

APP_DIR = Path(__file__).resolve().parent
APP_DATA_DIR = APP_DIR / "app_data"
HISTORY_FILE = APP_DATA_DIR / "analysis_history.json"
RULES_FILE = APP_DATA_DIR / "rules_store.json"
LOCAL_DB_FILE = APP_DATA_DIR / "autoanalista.db"

st.set_page_config(page_title="AutoAnalista 2026", page_icon=":bar_chart:", layout="wide")


def hash_text(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def repair_text_encoding(text: str) -> str:
    markers = ("\u00c3", "\u00c2", "\u00e2", "\u00f0\u0178", "\ufffd")
    if not any(marker in text for marker in markers):
        return text
    try:
        candidate = text.encode("latin1").decode("utf-8")
    except (UnicodeEncodeError, UnicodeDecodeError):
        return text
    marker_count = sum(text.count(marker) for marker in markers)
    candidate_marker_count = sum(candidate.count(marker) for marker in markers)
    return candidate if candidate_marker_count < marker_count else text


def normalize_token_text(value: object) -> str:
    text = repair_text_encoding(str(value))
    text = unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"[^A-Za-z0-9]+", "_", text.lower()).strip("_")
    return re.sub(r"_+", "_", text)


def clean_report_text(value: object, max_len: int = 180) -> str:
    if value is None or (isinstance(value, float) and pd.isna(value)):
        return "-"
    text = repair_text_encoding(str(value))
    text = unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"[\x00-\x1f\x7f-\x9f]", " ", text)
    text = re.sub(r"^[^A-Za-z0-9_]+", "", text)
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        text = "-"
    return text[: max_len - 3] + "..." if len(text) > max_len else text


def pdf_paragraph(text: object, style) -> Paragraph:
    return Paragraph(html.escape(clean_report_text(text, max_len=2000)), style)


def sync_session_value(source_key: str, state_key: str, mirror_key: Optional[str] = None) -> None:
    """Keep equivalent controls in sidebar and tab synchronized."""
    value = st.session_state.get(source_key)
    st.session_state[state_key] = value
    if mirror_key:
        st.session_state[mirror_key] = value


DEMO_AUTH_USERS = {
    "admin": {"password_hash": hash_text("admin2026!"), "role": "admin"},
    "analista": {"password_hash": hash_text("analista2026!"), "role": "analyst"},
    "viewer": {"password_hash": hash_text("viewer2026!"), "role": "viewer"},
}


def ensure_data_storage() -> None:
    APP_DATA_DIR.mkdir(parents=True, exist_ok=True)
    if not HISTORY_FILE.exists():
        HISTORY_FILE.write_text("[]", encoding="utf-8")
    if not RULES_FILE.exists():
        RULES_FILE.write_text("{}", encoding="utf-8")


def local_database_installed() -> bool:
    return LOCAL_DB_FILE.exists()


def read_history_json() -> List[dict]:
    ensure_data_storage()
    try:
        data = json.loads(HISTORY_FILE.read_text(encoding="utf-8"))
        return data if isinstance(data, list) else []
    except Exception:
        return []


def write_history_json(entries: List[dict]) -> None:
    ensure_data_storage()
    HISTORY_FILE.write_text(json.dumps(entries, ensure_ascii=False, indent=2), encoding="utf-8")


def connect_local_database() -> sqlite3.Connection:
    ensure_data_storage()
    conn = sqlite3.connect(LOCAL_DB_FILE)
    conn.row_factory = sqlite3.Row
    return conn


def init_local_database(migrate_json: bool = True) -> int:
    conn = connect_local_database()
    try:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS analysis_history (
                id TEXT PRIMARY KEY,
                timestamp TEXT NOT NULL,
                user TEXT,
                role TEXT,
                file_name TEXT,
                sheet_name TEXT,
                rows INTEGER,
                cols INTEGER,
                quality_score REAL,
                version INTEGER
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS app_meta (
                key TEXT PRIMARY KEY,
                value TEXT
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS stored_sheets (
                id TEXT PRIMARY KEY,
                timestamp TEXT NOT NULL,
                user TEXT,
                file_name TEXT,
                sheet_name TEXT,
                rows INTEGER,
                cols INTEGER,
                quality_score REAL,
                version INTEGER,
                metadata_json TEXT,
                csv_gzip BLOB
            )
            """
        )
        conn.execute(
            "INSERT OR REPLACE INTO app_meta (key, value) VALUES (?, ?)",
            ("schema_version", "2"),
        )
        inserted = 0
        if migrate_json:
            for item in read_history_json():
                conn.execute(
                    """
                    INSERT OR IGNORE INTO analysis_history
                    (id, timestamp, user, role, file_name, sheet_name, rows, cols, quality_score, version)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                    """,
                    (
                        str(item.get("id") or hash_text(json.dumps(item, sort_keys=True, default=str))),
                        str(item.get("timestamp", "")),
                        str(item.get("user", "")),
                        str(item.get("role", "")),
                        str(item.get("file_name", "")),
                        str(item.get("sheet_name", "")),
                        int(item.get("rows", 0) or 0),
                        int(item.get("cols", 0) or 0),
                        float(item.get("quality_score", 0.0) or 0.0),
                        int(item.get("version", 1) or 1),
                    ),
                )
                inserted += 1
        conn.commit()
        return inserted
    finally:
        conn.close()


def load_history_from_db() -> List[dict]:
    init_local_database(migrate_json=False)
    conn = connect_local_database()
    try:
        rows = conn.execute(
            """
            SELECT id, timestamp, user, role, file_name, sheet_name, rows, cols, quality_score, version
            FROM analysis_history
            ORDER BY timestamp ASC
            """
        ).fetchall()
        return [dict(row) for row in rows]
    finally:
        conn.close()


def save_history_to_db(entries: List[dict]) -> None:
    init_local_database(migrate_json=False)
    conn = connect_local_database()
    try:
        conn.execute("DELETE FROM analysis_history")
        conn.executemany(
            """
            INSERT OR REPLACE INTO analysis_history
            (id, timestamp, user, role, file_name, sheet_name, rows, cols, quality_score, version)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            [
                (
                    str(item.get("id") or hash_text(json.dumps(item, sort_keys=True, default=str))),
                    str(item.get("timestamp", "")),
                    str(item.get("user", "")),
                    str(item.get("role", "")),
                    str(item.get("file_name", "")),
                    str(item.get("sheet_name", "")),
                    int(item.get("rows", 0) or 0),
                    int(item.get("cols", 0) or 0),
                    float(item.get("quality_score", 0.0) or 0.0),
                    int(item.get("version", 1) or 1),
                )
                for item in entries
            ],
        )
        conn.commit()
    finally:
        conn.close()


def clear_analysis_history() -> None:
    write_history_json([])
    if local_database_installed():
        init_local_database(migrate_json=False)
        conn = connect_local_database()
        try:
            conn.execute("DELETE FROM analysis_history")
            conn.commit()
        finally:
            conn.close()


def dataframe_to_gzip_csv(df: pd.DataFrame) -> bytes:
    csv_bytes = df.to_csv(index=False).encode("utf-8")
    return gzip.compress(csv_bytes)


def dataframe_from_gzip_csv(blob: bytes) -> pd.DataFrame:
    return pd.read_csv(io.BytesIO(gzip.decompress(blob)))


def save_sheet_to_local_db(
    df: pd.DataFrame,
    *,
    user: str,
    file_name: str,
    sheet_name: str,
    quality_score: float,
    version: int,
    metadata: Optional[dict] = None,
) -> str:
    init_local_database(migrate_json=False)
    stored_id = hash_text(f"{user}|{file_name}|{sheet_name}|{version}|{datetime.now().isoformat()}")[:24]
    payload = dataframe_to_gzip_csv(df.head(MAX_EXPORT_ROWS))
    conn = connect_local_database()
    try:
        conn.execute(
            """
            INSERT OR REPLACE INTO stored_sheets
            (id, timestamp, user, file_name, sheet_name, rows, cols, quality_score, version, metadata_json, csv_gzip)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                stored_id,
                datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                user,
                file_name,
                sheet_name,
                int(len(df)),
                int(df.shape[1]),
                float(quality_score),
                int(version),
                json.dumps(metadata or {}, ensure_ascii=False, default=str),
                payload,
            ),
        )
        conn.commit()
        return stored_id
    finally:
        conn.close()


def list_stored_sheets() -> pd.DataFrame:
    if not local_database_installed():
        return pd.DataFrame()
    init_local_database(migrate_json=False)
    conn = connect_local_database()
    try:
        rows = conn.execute(
            """
            SELECT id, timestamp, user, file_name, sheet_name, rows, cols, quality_score, version
            FROM stored_sheets
            ORDER BY timestamp DESC
            """
        ).fetchall()
        return pd.DataFrame([dict(row) for row in rows])
    finally:
        conn.close()


def load_stored_sheet(stored_id: str) -> Optional[pd.DataFrame]:
    if not local_database_installed():
        return None
    init_local_database(migrate_json=False)
    conn = connect_local_database()
    try:
        row = conn.execute("SELECT csv_gzip FROM stored_sheets WHERE id = ?", (stored_id,)).fetchone()
        if row is None:
            return None
        return dataframe_from_gzip_csv(row["csv_gzip"])
    finally:
        conn.close()


def clear_stored_sheets() -> None:
    if not local_database_installed():
        return
    init_local_database(migrate_json=False)
    conn = connect_local_database()
    try:
        conn.execute("DELETE FROM stored_sheets")
        conn.commit()
    finally:
        conn.close()


def database_status_label() -> str:
    if not local_database_installed():
        return "Nao instalado"
    try:
        total = len(load_history_from_db())
        stored_total = len(list_stored_sheets())
        return f"Instalado ({total} historico, {stored_total} planilhas)"
    except Exception:
        return "Instalado com pendencia de inicializacao"


def load_rules_store() -> dict:
    ensure_data_storage()
    try:
        data = json.loads(RULES_FILE.read_text(encoding="utf-8"))
        return data if isinstance(data, dict) else {}
    except Exception:
        return {}


def save_rules_store(data: dict) -> None:
    ensure_data_storage()
    RULES_FILE.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def auth_key_for_rules(username: str, area: str) -> str:
    return f"{username.lower()}::{area.lower()}"


def load_auth_users() -> Tuple[dict, bool]:
    env_raw = os.getenv("AUTOANALISTA_USERS_JSON")
    if env_raw:
        try:
            parsed = json.loads(env_raw)
            if isinstance(parsed, dict) and parsed:
                return parsed, False
        except Exception:
            pass

    try:
        if "auth_users" in st.secrets:
            sec = st.secrets["auth_users"]
            parsed = json.loads(sec) if isinstance(sec, str) else dict(sec)
            if isinstance(parsed, dict) and parsed:
                return parsed, False
    except Exception:
        pass

    return DEMO_AUTH_USERS, True


def get_domain_template_rules(area: str, columns: List[str]) -> List[dict]:
    def add_rule(out: List[dict], column: str, rule: str, priority: str, suggestion: str, params: Optional[dict] = None) -> None:
        out.append(
            {
                "column": column,
                "rule": rule,
                "priority": priority,
                "suggestion": suggestion,
                "params": params or {},
            }
        )

    rules: List[dict] = []
    normalized = {col: col.lower() for col in columns}

    if area == "Financeiro":
        for col, low in normalized.items():
            if any(t in low for t in ["valor", "preco", "total", "saldo", "receita", "despesa"]):
                add_rule(rules, col, "non_negative", "high", "Validar valores negativos inesperados.")
            if any(t in low for t in ["data", "venc", "pagamento", "emissao"]):
                add_rule(rules, col, "no_future_date", "medium", "Revisar datas futuras indevidas.")

    elif area == "Vendas":
        for col, low in normalized.items():
            if any(t in low for t in ["valor", "preco", "ticket", "total", "quantidade", "qtd"]):
                add_rule(rules, col, "non_negative", "high", "Garantir integridade dos indicadores de venda.")
            if "email" in low:
                add_rule(rules, col, "email", "medium", "Padronizar emails dos clientes.")

    elif area == "RH":
        for col, low in normalized.items():
            if "cpf" in low:
                add_rule(rules, col, "cpf", "high", "Corrigir CPFs invalidos no cadastro.")
            if "email" in low:
                add_rule(rules, col, "email", "medium", "Padronizar emails corporativos.")
            if any(t in low for t in ["data", "admiss", "demiss", "nascimento"]):
                add_rule(rules, col, "no_future_date", "high", "Revisar datas futuras em dados de pessoas.")

    elif area == "Operacoes":
        for col, low in normalized.items():
            if any(t in low for t in ["tempo", "sla", "duracao", "idade"]):
                add_rule(rules, col, "non_negative", "high", "Tempo/idade nao devem ser negativos.")
            if any(t in low for t in ["data", "abertura", "fechamento"]):
                add_rule(rules, col, "no_future_date", "medium", "Verificar coerencia temporal.")

    else:
        for col, low in normalized.items():
            if "email" in low:
                add_rule(rules, col, "email", "medium", "Padronizar formato de email.")
            if "cpf" in low:
                add_rule(rules, col, "cpf", "high", "Corrigir CPF invalido.")

    dedup = []
    seen = set()
    for item in rules:
        key = (item["column"], item["rule"], item["priority"])
        if key not in seen:
            dedup.append(item)
            seen.add(key)
    return dedup


def normalize_rules_for_columns(rules: List[dict], columns: List[str]) -> List[dict]:
    col_set = set(columns)
    out = []
    for rule in rules:
        if not isinstance(rule, dict):
            continue
        col = rule.get("column")
        if col not in col_set:
            continue
        out.append(
            {
                "column": col,
                "rule": str(rule.get("rule", "")),
                "priority": str(rule.get("priority", "medium")),
                "suggestion": str(rule.get("suggestion", "Revisar inconsistencias.")),
                "params": dict(rule.get("params", {})) if isinstance(rule.get("params", {}), dict) else {},
            }
        )
    return out


def load_history() -> List[dict]:
    if local_database_installed():
        try:
            return load_history_from_db()
        except Exception:
            return read_history_json()
    return read_history_json()


def save_history(entries: List[dict]) -> None:
    write_history_json(entries)
    if local_database_installed():
        save_history_to_db(entries)


def next_report_version(history: List[dict], file_name: str, sheet_name: str) -> int:
    versions = [
        int(item.get("version", 0))
        for item in history
        if item.get("file_name") == file_name and item.get("sheet_name") == sheet_name
    ]
    return (max(versions) + 1) if versions else 1


def register_history(
    history: List[dict],
    *,
    user: str,
    role: str,
    file_name: str,
    sheet_name: str,
    rows: int,
    cols: int,
    score: float,
) -> dict:
    version = next_report_version(history, file_name, sheet_name)
    entry = {
        "id": datetime.now().strftime("%Y%m%d%H%M%S%f"),
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "user": user,
        "role": role,
        "file_name": file_name,
        "sheet_name": sheet_name,
        "rows": int(rows),
        "cols": int(cols),
        "quality_score": round(float(score), 2),
        "version": int(version),
    }
    history.append(entry)
    save_history(history)
    return entry


def apply_theme(theme_mode: str) -> None:
    if theme_mode == "Claro":
        st.markdown(
            """
            <style>
            .stApp { background: linear-gradient(180deg, #F4F7FB 0%, #FFFFFF 80%); color: #0E1A2B; }
            [data-testid="stSidebar"] { background: #ECF2FA; }
            </style>
            """,
            unsafe_allow_html=True,
        )
    else:
        st.markdown(
            """
            <style>
            .stApp { background: radial-gradient(circle at top, #0E2038 0%, #081220 70%); }
            [data-testid="stSidebar"] { background: #091A30; }
            </style>
            """,
            unsafe_allow_html=True,
        )


def get_secret_value(name: str) -> str:
    try:
        value = st.secrets.get(name, "")
        return str(value) if value else ""
    except Exception:
        return os.getenv(name, "")


def render_turbo_ai_controls() -> dict:
    secret_key_available = bool(get_secret_value("OPENAI_API_KEY"))
    st.sidebar.markdown(
        """
        <style>
        .turbo-api-card {
            margin: 10px 0 12px 0;
            padding: 13px 14px;
            border: 1px solid rgba(56, 189, 248, 0.45);
            border-radius: 14px;
            background:
                radial-gradient(circle at top left, rgba(14, 165, 233, 0.24), transparent 42%),
                linear-gradient(135deg, rgba(15, 23, 42, 0.92), rgba(8, 47, 73, 0.86));
            box-shadow: 0 0 0 1px rgba(125, 211, 252, 0.08), 0 14px 30px rgba(2, 8, 23, 0.28);
        }
        .turbo-api-title {
            font-size: 18px;
            font-weight: 900;
            letter-spacing: .3px;
            color: #E0F2FE;
        }
        .turbo-api-cta {
            margin-top: 8px;
            font-size: 15px;
            font-weight: 900;
            color: #FACC15;
            text-transform: uppercase;
            letter-spacing: .7px;
            text-shadow: 0 0 12px rgba(250, 204, 21, .35);
        }
        .turbo-api-note {
            margin-top: 6px;
            font-size: 12px;
            color: #BAE6FD;
        }
        </style>
        <div class="turbo-api-card">
            <div class="turbo-api-title">Modo Turbo IA</div>
            <div class="turbo-api-note">Opcional: o app continua funcionando sem API.</div>
            <div class="turbo-api-cta">Ativar modo turbo com API</div>
        </div>
        """,
        unsafe_allow_html=True,
    )
    api_key = st.sidebar.text_input(
        "Chave API",
        type="password",
        value="",
        placeholder="Cole sua API key aqui",
        key="turbo_api_key",
        help="A chave fica apenas na sessao atual e nao e salva no banco, historico ou relatorios.",
    )
    use_secret = False
    if secret_key_available and not api_key:
        use_secret = st.sidebar.checkbox("Usar chave configurada no Streamlit Secrets", value=True, key="turbo_use_secret")
    enabled = st.sidebar.checkbox("Ativar Modo Turbo", value=False, key="turbo_enabled")
    with st.sidebar.expander("Configuracao avancada do Turbo", expanded=False):
        model = st.text_input("Modelo", value="gpt-4o-mini", key="turbo_model")
        base_url = st.text_input("Endpoint", value="https://api.openai.com/v1", key="turbo_base_url")
        st.caption("Use endpoint compativel com `/chat/completions`.")

    selected_key = api_key.strip() or (get_secret_value("OPENAI_API_KEY") if use_secret else "")
    active = bool(enabled and selected_key)
    if enabled and not selected_key:
        st.sidebar.warning("Modo Turbo marcado, mas nenhuma chave API foi informada.")
    elif active:
        st.sidebar.success("Modo Turbo ativo. A IA sera usada sob demanda.")
    else:
        st.sidebar.caption("Modo Turbo desligado. Analises locais continuam ativas.")

    return {
        "enabled": active,
        "api_key": selected_key,
        "model": model.strip() or "gpt-4o-mini",
        "base_url": base_url.strip().rstrip("/") or "https://api.openai.com/v1",
    }


def turbo_ai_available(turbo_config: Optional[dict]) -> bool:
    return bool(turbo_config and turbo_config.get("enabled") and turbo_config.get("api_key"))


def friendly_turbo_api_error(status_code: int, detail: str) -> str:
    try:
        payload = json.loads(detail)
    except Exception:
        payload = {}
    error = payload.get("error", {}) if isinstance(payload, dict) else {}
    if isinstance(error, dict):
        code = str(error.get("code", "") or "")
        err_type = str(error.get("type", "") or "")
        message = str(error.get("message", "") or detail)
    elif isinstance(error, list):
        code = ""
        err_type = ""
        message = " | ".join([str(item) for item in error[:3]]) or detail
    elif error:
        code = ""
        err_type = ""
        message = str(error)
    elif isinstance(payload, dict):
        code = str(payload.get("code", "") or "")
        err_type = str(payload.get("type", "") or "")
        message = str(payload.get("message", "") or detail)
    else:
        code = ""
        err_type = ""
        message = detail

    if status_code == 429 and (code == "insufficient_quota" or "quota" in message.lower()):
        return (
            "**Modo Turbo indisponivel por cota da API.**\n\n"
            "A chave foi aceita, mas a conta/projeto da API esta sem creditos, sem billing ativo ou excedeu a cota atual.\n\n"
            "**Como resolver:**\n"
            "- Verifique saldo, billing e limites do projeto no painel da API.\n"
            "- Use outra chave API com cota disponivel.\n"
            "- Ou desligue o Modo Turbo e continue usando a analise local normalmente."
        )
    if status_code in [401, 403] or code in ["invalid_api_key", "invalid_api_key_format"]:
        return (
            "**Chave API recusada.**\n\n"
            "Confira se a chave foi copiada corretamente, se pertence ao projeto certo e se possui permissao para usar o modelo selecionado."
        )
    if status_code == 429 or "rate_limit" in code or "rate_limit" in err_type:
        return (
            "**Limite temporario da API atingido.**\n\n"
            "Aguarde alguns instantes e tente novamente, ou reduza a frequencia das chamadas Turbo."
        )
    if status_code == 404 or "model" in code:
        return (
            "**Modelo nao encontrado ou sem permissao.**\n\n"
            "Se estiver usando OpenAI, tente `gpt-4o-mini`. "
            "Se estiver usando Grok/xAI, tente um modelo disponivel no seu console, como `grok-4` ou `grok-3-mini`."
        )
    if status_code >= 500:
        return (
            "**API temporariamente indisponivel.**\n\n"
            "Tente novamente em alguns minutos. A analise local do AutoAnalista continua funcionando."
        )
    return f"**Falha ao acionar o Modo Turbo.**\n\nStatus da API: `{status_code}`.\n\nDetalhe: {clean_report_text(message, 420)}"


def call_turbo_ai(
    turbo_config: Optional[dict],
    *,
    section: str,
    context: dict,
    instruction: str,
    max_tokens: int = 750,
) -> Tuple[bool, str]:
    if not turbo_ai_available(turbo_config):
        return False, "Modo Turbo nao esta ativo. Informe uma chave API e ative o modo turbo na sidebar."

    payload = {
        "model": turbo_config.get("model", "gpt-4o-mini"),
        "temperature": 0.2,
        "max_tokens": max_tokens,
        "messages": [
            {
                "role": "system",
                "content": (
                    "Voce e um analista de dados senior em 2026, especialista em qualidade, gestao, "
                    "dashboard executivo e leitura de planilhas. Responda em portugues do Brasil, "
                    "com objetividade, sem inventar numeros que nao estejam no contexto."
                ),
            },
            {
                "role": "user",
                "content": (
                    f"Secao: {section}\n"
                    f"Instrucao: {instruction}\n\n"
                    "Contexto agregado da analise (sem dados brutos completos):\n"
                    f"{json.dumps(context, ensure_ascii=False, default=str)[:14000]}"
                ),
            },
        ],
    }
    url = f"{str(turbo_config.get('base_url', 'https://api.openai.com/v1')).rstrip('/')}/chat/completions"
    request = urllib.request.Request(
        url,
        data=json.dumps(payload).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {turbo_config.get('api_key')}",
            "Content-Type": "application/json",
        },
        method="POST",
    )
    try:
        with urllib.request.urlopen(request, timeout=45) as response:
            result = json.loads(response.read().decode("utf-8"))
        text = result.get("choices", [{}])[0].get("message", {}).get("content", "")
        return True, text.strip() or "A API respondeu sem texto utilizavel."
    except urllib.error.HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="ignore")[:2000]
        return False, friendly_turbo_api_error(exc.code, detail)
    except Exception as exc:
        return False, f"Nao foi possivel acionar o Modo Turbo: {str(exc)[:500]}"


def compact_records(df: pd.DataFrame, limit: int = 8) -> List[dict]:
    if df is None or len(df) == 0:
        return []
    return df.head(limit).astype(object).where(pd.notna(df.head(limit)), None).to_dict("records")


def build_ai_dataset_context(
    df: pd.DataFrame,
    groups: dict[str, List[str]],
    quality_report: Optional[dict] = None,
    profile: Optional[dict] = None,
    kpis: Optional[List[dict]] = None,
    issue_catalog: Optional[pd.DataFrame] = None,
    active_filters: Optional[List[str]] = None,
) -> dict:
    metric_numeric, id_like_cols = select_numeric_metric_columns(df, groups)
    volume_col = pick_volume_id_column(df, id_like_cols)
    profile = profile or detect_spreadsheet_profile(df, groups, quality_report or {})
    top_categories = []
    for col in groups.get("categorical", [])[:6]:
        vc = df[col].fillna("NULO").astype(str).value_counts().head(8)
        top_categories.append(
            {
                "coluna": col,
                "valores": [
                    {"categoria": str(idx), "quantidade": int(val), "participacao_pct": round(float(val) / max(len(df), 1) * 100, 2)}
                    for idx, val in vc.items()
                ],
            }
        )
    numeric_summary = []
    for col in metric_numeric[:8]:
        s = pd.to_numeric(df[col], errors="coerce").dropna()
        if len(s) == 0:
            continue
        numeric_summary.append(
            {
                "coluna": col,
                "media": float(s.mean()),
                "mediana": float(s.median()),
                "min": float(s.min()),
                "max": float(s.max()),
                "n_validos": int(len(s)),
            }
        )
    issues = compact_records(issue_catalog.sort_values(["priority_rank", "affected_rows"], ascending=[True, False]), 8) if issue_catalog is not None and len(issue_catalog) > 0 else []
    return {
        "linhas": int(len(df)),
        "colunas": int(df.shape[1]),
        "grupos_colunas": {k: v[:20] for k, v in groups.items()},
        "metricas_numericas_validas": metric_numeric[:15],
        "colunas_id_codigo": id_like_cols[:15],
        "coluna_volume": volume_col,
        "perfil_detectado": {
            "tipo": profile.get("tipo"),
            "confianca": profile.get("confianca"),
            "melhor_analise": profile.get("melhor_analise"),
            "graficos_recomendados": profile.get("graficos_recomendados", [])[:8],
        },
        "qualidade": {
            "score": quality_report.get("score") if quality_report else None,
            "nivel": quality_report.get("level") if quality_report else None,
            "ausencia_pct": round(float(quality_report.get("missing_ratio", 0.0)) * 100, 2) if quality_report else None,
            "duplicadas": quality_report.get("duplicates") if quality_report else None,
        },
        "kpis": [{k: v for k, v in item.items() if k != "comparison"} for item in (kpis or [])[:6]],
        "top_categorias": top_categories,
        "resumo_numerico": numeric_summary,
        "principais_riscos": issues,
        "filtros_ativos": active_filters or [],
    }


def render_turbo_ai_box(
    *,
    turbo_config: Optional[dict],
    section: str,
    button_label: str,
    context: dict,
    instruction: str,
    cache_key: str,
) -> None:
    st.markdown("**Modo Turbo IA**")
    if not turbo_ai_available(turbo_config):
        st.info("Opcional: informe uma chave API na sidebar e ative o Modo Turbo para receber uma leitura analitica por IA.")
        return
    if st.button(button_label, use_container_width=True, key=f"turbo_btn_{cache_key}"):
        with st.spinner("Modo Turbo analisando os dados agregados..."):
            ok, text = call_turbo_ai(turbo_config, section=section, context=context, instruction=instruction)
        st.session_state[f"turbo_result_{cache_key}"] = {"ok": ok, "text": text}
    result = st.session_state.get(f"turbo_result_{cache_key}")
    if result:
        if result.get("ok"):
            st.success("Analise Turbo gerada.")
            st.markdown(result.get("text", ""))
        else:
            st.warning(result.get("text", "Falha ao gerar analise Turbo."))
            st.caption("O AutoAnalista continua operando normalmente com as analises locais, dashboards, insights e ML sem API.")


def render_database_controls(role: str) -> None:
    st.sidebar.markdown("---")
    st.sidebar.subheader("Banco local SQLite")
    st.sidebar.caption(f"Status: {database_status_label()}")
    st.sidebar.caption(f"Arquivo local: `{LOCAL_DB_FILE.name}` dentro da pasta `app_data`.")

    if st.sidebar.button("Instalar/Inicializar banco local", use_container_width=True, key="install_local_sqlite"):
        migrated = init_local_database(migrate_json=True)
        st.sidebar.success(f"Banco SQLite pronto. Historico migrado: {migrated} registro(s).")

    if role == "admin":
        confirm_clear = st.sidebar.checkbox("Confirmar limpeza do historico", key="confirm_clear_history")
        if st.sidebar.button(
            "Limpar historico de analises",
            use_container_width=True,
            key="clear_analysis_history",
            disabled=not confirm_clear,
        ):
            clear_analysis_history()
            st.session_state.analysis_ready = False
            st.session_state.analysis_bundle = None
            st.session_state.analysis_key = None
            st.session_state.export_cache = {}
            st.sidebar.success("Historico de analises limpo.")
    else:
        st.sidebar.caption("Limpeza disponivel apenas para perfil admin.")


def render_offline_install_controls() -> None:
    st.sidebar.markdown("---")
    st.sidebar.subheader("Instalacao offline")
    st.sidebar.caption("Baixe um pacote ZIP para rodar o AutoAnalista no computador local, usando SQLite proprio.")
    try:
        offline_bytes = generate_offline_package()
        st.sidebar.download_button(
            "Baixar app offline",
            data=offline_bytes,
            file_name="autoanalista_2026_offline.zip",
            mime="application/zip",
            use_container_width=True,
            key="download_offline_sidebar",
        )
        st.sidebar.caption("Depois de extrair, execute `INICIAR_AUTOANALISTA.bat`.")
    except Exception as exc:
        st.sidebar.warning(f"Pacote offline indisponivel agora: {str(exc)[:120]}")


def login_panel() -> Optional[dict]:
    st.sidebar.subheader("Acesso")
    users, demo_mode = load_auth_users()
    now = datetime.now()

    if "auth" not in st.session_state:
        st.session_state.auth = None
    if "auth_last_activity" not in st.session_state:
        st.session_state.auth_last_activity = None
    if "auth_failed_attempts" not in st.session_state:
        st.session_state.auth_failed_attempts = 0
    if "auth_locked_until" not in st.session_state:
        st.session_state.auth_locked_until = None

    # Session timeout
    if st.session_state.auth is not None and st.session_state.auth_last_activity is not None:
        elapsed = now - st.session_state.auth_last_activity
        if elapsed > timedelta(minutes=AUTH_SESSION_TIMEOUT_MINUTES):
            st.session_state.auth = None
            st.session_state.auth_last_activity = None
            st.sidebar.warning("Sessao expirada por inatividade. Faca login novamente.")
            st.rerun()

    if st.session_state.auth is not None:
        auth = st.session_state.auth
        st.session_state.auth_last_activity = now
        st.sidebar.success(f"Logado: {auth['username']} ({auth['role']})")
        st.sidebar.caption(f"Timeout de sessao: {AUTH_SESSION_TIMEOUT_MINUTES} min")
        if st.sidebar.button("Sair", use_container_width=True):
            st.session_state.auth = None
            st.session_state.auth_last_activity = None
            st.rerun()
        return auth

    locked_until = st.session_state.auth_locked_until
    is_locked = locked_until is not None and locked_until > now
    if is_locked:
        remaining = int((locked_until - now).total_seconds() // 60) + 1
        st.sidebar.error(f"Login bloqueado temporariamente. Tente novamente em {remaining} min.")

    if demo_mode:
        st.sidebar.caption("Modo demo ativo. Configure usuarios seguros via secrets/env antes de uso em producao.")

    username = st.sidebar.text_input("Usuario", value="")
    password = st.sidebar.text_input("Senha", type="password", value="")

    if st.sidebar.button("Entrar", use_container_width=True, disabled=is_locked):
        rec = users.get(username)
        valid = False
        if rec and "password_hash" in rec and "role" in rec:
            incoming_hash = hash_text(password)
            valid = hmac.compare_digest(str(incoming_hash), str(rec["password_hash"]))

        if valid:
            st.session_state.auth = {"username": username, "role": rec["role"]}
            st.session_state.auth_last_activity = now
            st.session_state.auth_failed_attempts = 0
            st.session_state.auth_locked_until = None
            st.rerun()
        else:
            st.session_state.auth_failed_attempts += 1
            attempts = st.session_state.auth_failed_attempts
            if attempts >= AUTH_MAX_ATTEMPTS:
                st.session_state.auth_locked_until = now + timedelta(minutes=AUTH_LOCK_MINUTES)
                st.session_state.auth_failed_attempts = 0
                st.sidebar.error(f"Credenciais invalidas. Login bloqueado por {AUTH_LOCK_MINUTES} min.")
            else:
                remain = AUTH_MAX_ATTEMPTS - attempts
                st.sidebar.error(f"Credenciais invalidas. Tentativas restantes: {remain}.")

    return None


def sanitize_columns(df: pd.DataFrame) -> pd.DataFrame:
    clean = df.copy()
    sanitized = []
    seen: Dict[str, int] = {}
    for i, col in enumerate(clean.columns):
        name = clean_report_text(col, 80)
        if name == "-":
            name = f"coluna_{i}"
        base = name
        count = seen.get(base, 0)
        if count:
            name = f"{base}_{count + 1}"
        seen[base] = count + 1
        sanitized.append(name)
    clean.columns = sanitized
    return clean


def read_csv_like(file_bytes: bytes, extension: str) -> pd.DataFrame:
    errors = []
    for encoding in ["utf-8", "utf-8-sig", "latin-1", "cp1252"]:
        try:
            if extension == ".tsv":
                return pd.read_csv(io.BytesIO(file_bytes), sep="\t", encoding=encoding)
            return pd.read_csv(io.BytesIO(file_bytes), sep=None, engine="python", encoding=encoding)
        except Exception as exc:
            errors.append(f"{encoding}: {exc}")
    raise ValueError("Falha ao ler arquivo texto. " + " | ".join(errors[:2]))


@st.cache_data(show_spinner=False)
def read_workbook(file_bytes: bytes, file_name: str) -> Dict[str, pd.DataFrame]:
    extension = Path(file_name).suffix.lower()

    if extension in [".csv", ".tsv", ".txt"]:
        return {"Dados": sanitize_columns(read_csv_like(file_bytes, extension))}

    if extension == ".xls":
        engines: List[Optional[str]] = [None, "xlrd", "calamine"]
    elif extension in [".xlsx", ".xlsm", ".xltx", ".xltm"]:
        engines = [None, "openpyxl", "calamine"]
    elif extension in [".xlsb", ".ods"]:
        engines = [None, "calamine"]
    else:
        engines = [None, "calamine", "openpyxl", "xlrd"]

    errors: List[str] = []
    for engine in engines:
        try:
            data = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None, engine=engine)
            return {name: sanitize_columns(df) for name, df in data.items()}
        except Exception as exc:
            engine_name = "auto" if engine is None else engine
            errors.append(f"{engine_name}: {exc}")

    if extension == ".xls":
        try:
            tables = pd.read_html(io.BytesIO(file_bytes))
            if tables:
                return {f"Tabela_{i + 1}": sanitize_columns(tbl) for i, tbl in enumerate(tables)}
        except Exception as exc:
            errors.append(f"html: {exc}")

    raise ValueError(
        "Falha ao ler planilha com os leitores disponiveis. "
        "Abra no Excel e salve como .xlsx. Detalhes: "
        + " | ".join(errors[:3])
    )


def best_numeric_parse(series: pd.Series) -> Tuple[pd.Series, float]:
    txt = series.astype(str).str.strip().replace({"": np.nan, "None": np.nan, "nan": np.nan, "NaN": np.nan, "-": np.nan})
    standard = pd.to_numeric(txt, errors="coerce")
    standard_ratio = float(standard.notna().mean())

    local = txt.str.replace(r"[R$\s%]", "", regex=True).str.replace(".", "", regex=False).str.replace(",", ".", regex=False)
    local_num = pd.to_numeric(local, errors="coerce")
    local_ratio = float(local_num.notna().mean())

    if local_ratio > standard_ratio:
        return local_num, local_ratio
    return standard, standard_ratio


@st.cache_data(show_spinner=False)
def coerce_data_types(df: pd.DataFrame) -> Tuple[pd.DataFrame, List[str]]:
    converted = df.copy()
    notes: List[str] = []

    for col in converted.columns:
        if converted[col].dtype != "object":
            continue

        non_null = converted[col].dropna()
        if len(non_null) < 5:
            continue

        dt_ratio = float(pd.to_datetime(non_null, errors="coerce", dayfirst=True).notna().mean())
        _, num_ratio = best_numeric_parse(non_null)

        if dt_ratio >= 0.9:
            converted[col] = pd.to_datetime(converted[col], errors="coerce", dayfirst=True)
            notes.append(f"Coluna `{col}` convertida para data/hora.")
            continue

        if num_ratio >= 0.9:
            num_col, _ = best_numeric_parse(converted[col])
            converted[col] = num_col
            notes.append(f"Coluna `{col}` convertida para numerica.")

    return converted, notes


def get_column_groups(df: pd.DataFrame) -> dict[str, List[str]]:
    numeric = df.select_dtypes(include=np.number).columns.tolist()
    datetime_cols = df.select_dtypes(include=["datetime64[ns]", "datetime64[ns, UTC]"]).columns.tolist()
    categorical = [c for c in df.columns if c not in numeric + datetime_cols]
    return {"numeric": numeric, "datetime": datetime_cols, "categorical": categorical}


def detect_mixed_type_columns(df: pd.DataFrame) -> List[str]:
    mixed = []
    object_cols = df.select_dtypes(include=["object"]).columns.tolist()
    for col in object_cols:
        values = df[col].dropna().astype(str).str.strip()
        if len(values) < 20:
            continue
        num_ratio = best_numeric_parse(values)[1]
        dt_ratio = float(pd.to_datetime(values, errors="coerce", dayfirst=True).notna().mean())
        m = max(num_ratio, dt_ratio)
        if 0.2 < m < 0.85:
            mixed.append(col)
    return mixed


def detect_id_like_columns(df: pd.DataFrame) -> List[str]:
    if len(df) == 0:
        return []
    out = []
    for col in df.columns:
        non_null = df[col].dropna()
        if len(non_null) < 20:
            continue
        col_key = normalize_token_text(col)
        name_tokens = set(col_key.split("_"))
        name_looks_like_id = (
            bool(name_tokens.intersection(ID_NAME_TOKENS))
            or any(col_key.startswith(prefix) for prefix in ID_NAME_PREFIXES)
            or col_key.endswith("_id")
        )
        unique_ratio = non_null.nunique(dropna=True) / max(len(non_null), 1)
        numeric_values = pd.to_numeric(non_null, errors="coerce")
        numeric_ratio = float(numeric_values.notna().mean()) if len(non_null) else 0.0
        integer_like = False
        large_code_like = False
        if numeric_ratio >= 0.9:
            numeric_clean = numeric_values.dropna()
            if len(numeric_clean) > 0:
                integer_like = bool(np.isclose(numeric_clean, np.round(numeric_clean)).mean() >= 0.98)
                large_code_like = bool(numeric_clean.abs().median() >= 1000 and integer_like)

        if name_looks_like_id:
            out.append(col)
            continue
        if unique_ratio >= 0.95 and (integer_like or df[col].dtype == "object"):
            out.append(col)
            continue
        if unique_ratio >= 0.75 and large_code_like:
            out.append(col)
    return out


def outlier_ratio_iqr(df: pd.DataFrame, numeric_cols: List[str]) -> float:
    total_points = 0
    total_outliers = 0
    for col in numeric_cols:
        s = df[col].dropna()
        if len(s) < 8:
            continue
        q1 = s.quantile(0.25)
        q3 = s.quantile(0.75)
        iqr = q3 - q1
        if iqr <= 0:
            continue
        lower = q1 - 1.5 * iqr
        upper = q3 + 1.5 * iqr
        flags = (s < lower) | (s > upper)
        total_outliers += int(flags.sum())
        total_points += int(len(s))
    return (total_outliers / total_points) if total_points else 0.0


def is_valid_email(value: str) -> bool:
    return bool(re.match(r"^[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}$", str(value).strip()))


def is_valid_cpf(value: str) -> bool:
    digits = re.sub(r"\D", "", str(value))
    if len(digits) != 11 or digits == digits[0] * 11:
        return False
    total = sum(int(digits[i]) * (10 - i) for i in range(9))
    d1 = (total * 10) % 11
    d1 = 0 if d1 == 10 else d1
    if d1 != int(digits[9]):
        return False
    total = sum(int(digits[i]) * (11 - i) for i in range(10))
    d2 = (total * 10) % 11
    d2 = 0 if d2 == 10 else d2
    return d2 == int(digits[10])


def append_violation_records(
    records: List[dict],
    df: pd.DataFrame,
    column: str,
    mask: pd.Series,
    *,
    rule_name: str,
    priority: str,
    message: str,
    suggestion: str,
    source: str,
    max_records: int = 2000,
) -> None:
    idxs = df.index[mask.fillna(False)].tolist()
    for idx in idxs[:max_records]:
        val = df.at[idx, column] if column in df.columns else None
        records.append(
            {
                "row_index": int(idx),
                "excel_row": int(idx) + 2,
                "column": column,
                "rule": rule_name,
                "priority": priority,
                "value": str(val),
                "message": message,
                "suggestion": suggestion,
                "source": source,
            }
        )


def rule_mask(df: pd.DataFrame, column: str, rule: str, params: dict) -> pd.Series:
    if column not in df.columns:
        return pd.Series(False, index=df.index)

    s = df[column]

    if rule == "required":
        if s.dtype == "object":
            return s.isna() | (s.astype(str).str.strip() == "")
        return s.isna()

    if rule == "email":
        txt = s.astype(str).str.strip()
        return (~s.isna()) & (~txt.apply(is_valid_email))

    if rule == "cpf":
        txt = s.astype(str).str.strip()
        return (~s.isna()) & (~txt.apply(is_valid_cpf))

    if rule == "non_negative":
        num, _ = best_numeric_parse(s)
        return (~num.isna()) & (num < 0)

    if rule == "no_future_date":
        dt = pd.to_datetime(s, errors="coerce", dayfirst=True)
        return (~dt.isna()) & (dt > pd.Timestamp.now())

    if rule == "min_value":
        num, _ = best_numeric_parse(s)
        v = params.get("min_value")
        return (~num.isna()) & (num < float(v)) if v is not None else pd.Series(False, index=df.index)

    if rule == "max_value":
        num, _ = best_numeric_parse(s)
        v = params.get("max_value")
        return (~num.isna()) & (num > float(v)) if v is not None else pd.Series(False, index=df.index)

    return pd.Series(False, index=df.index)


def render_custom_rules_editor(columns: List[str], can_edit: bool, username: str, area: str) -> List[dict]:
    st.sidebar.markdown("---")
    st.sidebar.subheader("Regras por Coluna")

    if "custom_rules" not in st.session_state:
        st.session_state.custom_rules = []
    if "custom_rules_context" not in st.session_state:
        st.session_state.custom_rules_context = None

    columns_signature = hash_text("|".join(map(str, columns)))[:12]
    context_key = f"{username}::{area}::{columns_signature}"
    store = load_rules_store()
    user_area_key = auth_key_for_rules(username, area)

    # Load persisted rules when context changes
    if st.session_state.custom_rules_context != context_key:
        persisted = store.get(user_area_key, [])
        st.session_state.custom_rules = normalize_rules_for_columns(persisted, columns)
        st.session_state.custom_rules_context = context_key

    if not columns:
        st.sidebar.caption("Sem colunas para regras.")
        return st.session_state.custom_rules

    def persist_rules() -> None:
        latest = load_rules_store()
        latest[user_area_key] = normalize_rules_for_columns(st.session_state.custom_rules, columns)
        save_rules_store(latest)

    if can_edit:
        top1, top2 = st.sidebar.columns(2)
        if top1.button("Template Area", use_container_width=True):
            st.session_state.custom_rules = normalize_rules_for_columns(get_domain_template_rules(area, columns), columns)
            persist_rules()
            st.sidebar.success("Template aplicado e salvo.")
        if top2.button("Regras Salvas", use_container_width=True):
            persisted = store.get(user_area_key, [])
            st.session_state.custom_rules = normalize_rules_for_columns(persisted, columns)
            st.sidebar.success("Regras salvas carregadas.")

        with st.sidebar.form("rule_form", clear_on_submit=True):
            col = st.selectbox("Coluna", options=columns, key=f"rule_form_column_{columns_signature}")
            rule = st.selectbox(
                "Regra",
                options=["required", "email", "cpf", "non_negative", "no_future_date", "min_value", "max_value"],
                key=f"rule_form_rule_{columns_signature}",
            )
            min_value = None
            max_value = None
            if rule == "min_value":
                min_value = st.number_input("Valor minimo", value=0.0, key=f"rule_form_min_value_{columns_signature}")
            if rule == "max_value":
                max_value = st.number_input("Valor maximo", value=0.0, key=f"rule_form_max_value_{columns_signature}")
            priority = st.selectbox("Prioridade", options=["critical", "high", "medium", "low"], index=2, key=f"rule_form_priority_{columns_signature}")
            suggestion = st.text_input(
                "Acao sugerida",
                value="Revisar e corrigir os registros fora do padrao.",
                key=f"rule_form_suggestion_{columns_signature}",
            )
            submitted = st.form_submit_button("Adicionar regra")
            if submitted:
                params = {}
                if min_value is not None:
                    params["min_value"] = float(min_value)
                if max_value is not None:
                    params["max_value"] = float(max_value)
                st.session_state.custom_rules.append(
                    {
                        "column": col,
                        "rule": rule,
                        "priority": priority,
                        "suggestion": suggestion,
                        "params": params,
                    }
                )
                persist_rules()
                st.sidebar.success("Regra adicionada.")

        if st.session_state.custom_rules:
            rules_df = pd.DataFrame(
                [{"#": i, "column": r["column"], "rule": r["rule"], "priority": r["priority"]} for i, r in enumerate(st.session_state.custom_rules)]
            )
            st.sidebar.dataframe(rules_df, use_container_width=True, height=180)
            idx = st.sidebar.number_input("Indice para remover", min_value=0, max_value=max(0, len(st.session_state.custom_rules) - 1), value=0)
            c1, c2 = st.sidebar.columns(2)
            if c1.button("Remover", use_container_width=True) and st.session_state.custom_rules:
                st.session_state.custom_rules.pop(int(idx))
                persist_rules()
                st.sidebar.success("Regra removida.")
            if c2.button("Limpar", use_container_width=True):
                st.session_state.custom_rules = []
                persist_rules()
                st.sidebar.success("Regras limpas.")

        if st.sidebar.button("Salvar Regras", use_container_width=True):
            persist_rules()
            st.sidebar.success("Regras salvas para este usuario/area.")
    else:
        st.sidebar.caption("Perfil viewer sem permissao para editar regras.")
        if st.session_state.custom_rules:
            preview = pd.DataFrame(
                [{"column": r["column"], "rule": r["rule"], "priority": r["priority"]} for r in st.session_state.custom_rules]
            )
            st.sidebar.dataframe(preview, use_container_width=True, height=170)

    return st.session_state.custom_rules


def run_auto_validations(df: pd.DataFrame, groups: dict[str, List[str]]) -> pd.DataFrame:
    records: List[dict] = []

    for col in [c for c in df.columns if "email" in c.lower()]:
        mask = (~df[col].isna()) & (~df[col].astype(str).str.strip().apply(is_valid_email))
        append_violation_records(
            records,
            df,
            col,
            mask,
            rule_name="email",
            priority="high",
            message="Formato de email invalido.",
            suggestion="Corrigir para usuario@dominio.",
            source="auto",
            max_records=1000,
        )

    for col in [c for c in df.columns if "cpf" in c.lower()]:
        mask = (~df[col].isna()) & (~df[col].astype(str).str.strip().apply(is_valid_cpf))
        append_violation_records(
            records,
            df,
            col,
            mask,
            rule_name="cpf",
            priority="high",
            message="CPF invalido.",
            suggestion="Corrigir digitos verificadores.",
            source="auto",
            max_records=1000,
        )

    for col in groups["datetime"]:
        dt = pd.to_datetime(df[col], errors="coerce", dayfirst=True)
        mask = (~dt.isna()) & (dt > pd.Timestamp.now())
        append_violation_records(
            records,
            df,
            col,
            mask,
            rule_name="no_future_date",
            priority="medium",
            message="Data futura detectada.",
            suggestion="Revisar data acima da data atual.",
            source="auto",
            max_records=1000,
        )

    tokens = ["valor", "preco", "quantidade", "qtd", "saldo", "idade", "total"]
    for col in [c for c in groups["numeric"] if any(tok in c.lower() for tok in tokens)]:
        mask = (~df[col].isna()) & (df[col] < 0)
        append_violation_records(
            records,
            df,
            col,
            mask,
            rule_name="non_negative",
            priority="medium",
            message="Valor negativo inesperado.",
            suggestion="Verificar regra de negocio para negativos.",
            source="auto",
            max_records=1000,
        )

    return pd.DataFrame(records)


def run_custom_validations(df: pd.DataFrame, rules: List[dict]) -> pd.DataFrame:
    records: List[dict] = []
    for rule_item in rules:
        col = rule_item.get("column")
        rule = rule_item.get("rule")
        if col not in df.columns or not rule:
            continue
        mask = rule_mask(df, col, rule, rule_item.get("params", {}))
        append_violation_records(
            records,
            df,
            col,
            mask,
            rule_name=rule,
            priority=rule_item.get("priority", "medium"),
            message=f"Violacao da regra `{rule}`.",
            suggestion=rule_item.get("suggestion", "Revisar inconsistencias."),
            source="custom",
            max_records=3000,
        )
    return pd.DataFrame(records)


def compute_quality_report(df: pd.DataFrame, groups: dict[str, List[str]], violations: pd.DataFrame) -> dict:
    rows, cols = df.shape
    total_cells = rows * cols
    missing = int(df.isna().sum().sum())
    missing_ratio = (missing / total_cells) if total_cells else 0.0
    duplicates = int(df.duplicated().sum())
    duplicate_ratio = (duplicates / rows) if rows else 0.0
    constant_cols = [c for c in df.columns if df[c].nunique(dropna=False) <= 1]
    constant_ratio = (len(constant_cols) / cols) if cols else 0.0
    mixed_cols = detect_mixed_type_columns(df)
    mixed_ratio = (len(mixed_cols) / cols) if cols else 0.0
    outlier_ratio = outlier_ratio_iqr(df, groups["numeric"])
    id_like_cols = detect_id_like_columns(df)
    violations_ratio = min(1.0, len(violations) / max(len(df), 1)) if len(violations) > 0 else 0.0

    completeness = max(0.0, 100.0 * (1.0 - missing_ratio))
    uniqueness = max(0.0, 100.0 * (1.0 - duplicate_ratio))
    consistency = max(0.0, 100.0 - (mixed_ratio * 60.0 + constant_ratio * 40.0))
    validity = max(0.0, 100.0 - (violations_ratio * 70.0 + outlier_ratio * 30.0))

    if groups["datetime"]:
        now = pd.Timestamp.now()
        scores = []
        for col in groups["datetime"]:
            dt = pd.to_datetime(df[col], errors="coerce", dayfirst=True).dropna()
            if len(dt) == 0:
                continue
            future_ratio = float((dt > now).mean())
            if (now - dt.max()).days <= 365:
                freshness = 1.0
            elif (now - dt.max()).days <= 730:
                freshness = 0.85
            else:
                freshness = 0.65
            scores.append(max(0.0, 100.0 * (1.0 - future_ratio)) * freshness)
        timeliness = float(np.mean(scores)) if scores else 80.0
    else:
        timeliness = 80.0

    score = completeness * 0.25 + validity * 0.25 + consistency * 0.20 + uniqueness * 0.15 + timeliness * 0.15
    score = float(max(0.0, min(100.0, score)))

    if score >= 85:
        level = "Excelente"
    elif score >= 70:
        level = "Boa"
    elif score >= 50:
        level = "Atencao"
    else:
        level = "Critica"

    actions = []
    if missing_ratio > 0.15:
        actions.append("Definir estrategia de imputacao nas colunas com maior ausencia.")
    if duplicate_ratio > 0.03:
        actions.append("Remover/consolidar duplicidades antes de analises preditivas.")
    if len(mixed_cols) > 0:
        actions.append("Padronizar tipos de dados em colunas mistas.")
    if outlier_ratio > 0.05:
        actions.append("Revisar outliers para separar erro operacional de evento real.")
    if len(id_like_cols) > 0:
        actions.append("Evitar usar identificadores puros como preditores.")
    if not actions:
        actions.append("Qualidade boa: manter monitoramento continuo por lote.")

    return {
        "score": score,
        "level": level,
        "missing_ratio": float(missing_ratio),
        "duplicates": int(duplicates),
        "duplicate_ratio": float(duplicate_ratio),
        "constant_cols": constant_cols,
        "mixed_cols": mixed_cols,
        "outlier_ratio": float(outlier_ratio),
        "id_like_cols": id_like_cols,
        "actions": actions,
        "pillars": {
            "Completude": float(completeness),
            "Consistencia": float(consistency),
            "Unicidade": float(uniqueness),
            "Validade": float(validity),
            "Atualidade": float(timeliness),
        },
    }


def strongest_correlations(df: pd.DataFrame, top_n: int = 10) -> pd.DataFrame:
    id_like_cols = set(detect_id_like_columns(df))
    numeric_cols = [c for c in df.select_dtypes(include=np.number).columns.tolist() if c not in id_like_cols]
    if len(numeric_cols) < 2:
        return pd.DataFrame()
    corr = df[numeric_cols].corr(numeric_only=True)
    pairs = []
    for i, c1 in enumerate(corr.columns):
        for c2 in corr.columns[i + 1 :]:
            value = corr.loc[c1, c2]
            if pd.notna(value):
                pairs.append({"variavel_1": c1, "variavel_2": c2, "correlacao": float(value), "abs": abs(float(value))})
    if not pairs:
        return pd.DataFrame()
    return pd.DataFrame(pairs).sort_values("abs", ascending=False).head(top_n).drop(columns="abs")


def select_time_frequency(span_days: int) -> Tuple[str, str]:
    if span_days > 365:
        return "ME", "mes"
    if span_days > 90:
        return "W-SUN", "semana"
    return "D", "dia"


def safe_resample_series(series: pd.Series, freq: str, agg_mode: str) -> pd.Series:
    candidates = [freq]
    legacy_alias = {"ME": "M", "QE": "Q", "YE": "Y", "BME": "BM", "BQE": "BQ", "BYE": "BA", "W-SUN": "W"}
    old = legacy_alias.get(freq)
    if old:
        candidates.append(old)

    last_exc: Optional[Exception] = None
    for cand in candidates:
        try:
            rs = series.resample(cand)
            if agg_mode == "mean":
                return rs.mean()
            if agg_mode == "median":
                return rs.median()
            if agg_mode == "count":
                return rs.count()
            if agg_mode == "nunique":
                return rs.nunique()
            return rs.sum()
        except Exception as exc:  # pragma: no cover - defensive for pandas alias changes
            last_exc = exc

    if last_exc:
        raise last_exc
    raise ValueError("Falha ao agregar serie temporal.")


def select_numeric_metric_columns(df: pd.DataFrame, groups: dict[str, List[str]]) -> Tuple[List[str], List[str]]:
    id_like_cols = set(detect_id_like_columns(df))
    metric_cols = [c for c in groups["numeric"] if c not in id_like_cols and df[c].notna().sum() > 0]
    return metric_cols, sorted(list(id_like_cols))


def pick_volume_id_column(df: pd.DataFrame, id_like_cols: List[str]) -> Optional[str]:
    if not id_like_cols:
        return None
    tokens = ["chamado", "ticket", "protocolo", "os", "ordem", "id"]
    ranked = sorted(
        id_like_cols,
        key=lambda c: (
            any(tok in c.lower() for tok in tokens),
            df[c].notna().sum(),
            df[c].nunique(dropna=True),
        ),
        reverse=True,
    )
    return ranked[0] if ranked else None


def period_comparison(df: pd.DataFrame, date_col: str, value_col: str, agg_mode: str = "sum") -> Optional[dict]:
    base = df[[date_col, value_col]].dropna().sort_values(date_col)
    if len(base) < 10:
        return None
    span_days = int((base[date_col].max() - base[date_col].min()).days)
    freq, label = select_time_frequency(span_days)
    series = base.set_index(date_col)[value_col]
    agg = safe_resample_series(series, freq, agg_mode).dropna()
    if len(agg) < 2:
        return None
    current, previous = float(agg.iloc[-1]), float(agg.iloc[-2])
    if abs(previous) < 1e-9:
        delta_pct = None
        delta_txt = f"variacao absoluta de {current - previous:.2f}"
    else:
        delta_pct = ((current - previous) / abs(previous)) * 100.0
        direction = "alta" if delta_pct >= 0 else "queda"
        delta_txt = f"{direction} de {abs(delta_pct):.1f}%"
    return {"label": label, "delta_pct": delta_pct, "delta_txt": delta_txt, "series": agg.reset_index()}


def detect_kpis(df: pd.DataFrame, groups: dict[str, List[str]]) -> List[dict]:
    kpis = []
    metric_numeric, id_like_cols = select_numeric_metric_columns(df, groups)
    date_col = max(groups["datetime"], key=lambda c: df[c].notna().sum()) if groups["datetime"] else None

    volume_col = pick_volume_id_column(df, id_like_cols)
    if volume_col:
        comp = period_comparison(df, date_col, volume_col, agg_mode="nunique") if date_col else None
        item = {
            "name": f"IDs unicos ({volume_col})",
            "metric": float(df[volume_col].nunique(dropna=True)),
            "value_type": "count",
            "delta_text": "",
            "alert": "",
            "comparison": comp,
        }
        if comp:
            item["delta_text"] = comp["delta_txt"]
            if comp["delta_pct"] is not None and abs(comp["delta_pct"]) >= 15:
                item["alert"] = "Variacao relevante detectada"
        kpis.append(item)
    elif len(df) > 0:
        item = {
            "name": "Volume de linhas",
            "metric": float(len(df)),
            "value_type": "count",
            "delta_text": "",
            "alert": "",
            "comparison": None,
        }
        kpis.append(item)

    candidates = sorted(
        metric_numeric,
        key=lambda c: (df[c].notna().sum(), df[c].var(skipna=True) if pd.notna(df[c].var(skipna=True)) else -1),
        reverse=True,
    )[:3]

    for col in candidates:
        item = {
            "name": f"Media de {col}",
            "metric": float(df[col].mean(skipna=True)) if df[col].notna().sum() else 0.0,
            "value_type": "avg",
            "delta_text": "",
            "alert": "",
            "comparison": None,
        }
        if date_col:
            comp = period_comparison(df, date_col, col)
            if comp:
                item["comparison"] = comp
                item["delta_text"] = comp["delta_txt"]
                if comp["delta_pct"] is not None and abs(comp["delta_pct"]) >= 15:
                    item["alert"] = "Variacao relevante detectada"
        kpis.append(item)

    return kpis[:4]


def generate_professional_insights(
    df: pd.DataFrame,
    groups: dict[str, List[str]],
    quality_report: dict,
    top_corr: pd.DataFrame,
    kpis: List[dict],
    treatment_report: Optional[dict] = None,
) -> List[str]:
    insights = [
        f"Base analisada com {len(df)} linhas e {df.shape[1]} colunas.",
        f"Indice de qualidade: {quality_report['score']:.1f}/100 ({quality_report['level']}).",
    ]

    missing_by_col = (df.isna().mean() * 100).sort_values(ascending=False)
    high_missing = missing_by_col[missing_by_col >= 30]
    if not high_missing.empty:
        detail = ", ".join([f"{i} ({v:.1f}%)" for i, v in high_missing.head(3).items()])
        insights.append(f"Colunas com ausencia elevada: {detail}.")

    if quality_report["duplicates"] > 0:
        insights.append(f"Foram detectadas {quality_report['duplicates']} linhas duplicadas.")

    if treatment_report and treatment_report.get("total_removed", 0) > 0:
        insights.append(
            "Tratativa automatica removeu "
            f"{treatment_report['total_removed']} linhas "
            f"(vazias={treatment_report['empty_rows_removed']}, "
            f"duplicadas={treatment_report['duplicates_removed']}, "
            f"outliers={treatment_report['outliers_removed']})."
        )

    if quality_report["mixed_cols"]:
        insights.append("Colunas com tipos mistos: " + ", ".join(quality_report["mixed_cols"][:3]) + ".")

    if not top_corr.empty:
        row = top_corr.iloc[0]
        direction = "positiva" if row["correlacao"] >= 0 else "negativa"
        insights.append(f"Maior correlacao {direction}: `{row['variavel_1']}` x `{row['variavel_2']}` = {row['correlacao']:.3f}.")

    for item in kpis:
        if item["delta_text"]:
            insights.append(f"KPI `{item['name']}` com {item['delta_text']} no periodo mais recente.")
            break

    for col in groups["categorical"][:8]:
        vc = df[col].fillna("NULO").astype(str).value_counts(normalize=True)
        if len(vc) > 0 and float(vc.iloc[0]) >= 0.6:
            insights.append(f"Concentracao alta em `{col}`: categoria principal representa {vc.iloc[0] * 100:.1f}%.")
            break

    return insights[:10]


def build_action_plan(area: str, quality_report: dict, issue_catalog: pd.DataFrame) -> List[str]:
    action_df = build_structured_action_plan(area, quality_report, issue_catalog)
    return [
        f"{row['prioridade']} | Problema: {row.get('problema', row.get('foco', '-'))} | "
        f"Impacto: {row['impacto']} | Acao: {row['acao']} | "
        f"Responsavel: {row.get('responsavel', 'Dono do dado')} | Prazo: {row['prazo']}"
        for _, row in action_df.head(8).iterrows()
    ]


def suggested_action_owner(category: str, area: str) -> str:
    category_key = normalize_token_text(category)
    if category_key in ["completude", "unicidade", "consistencia"]:
        return "Dono do dado + Dados"
    if category_key in ["validade", "regras"]:
        return "Area de negocio + Qualidade"
    if area and area != "Geral":
        return f"Gestor de {area}"
    return "Governanca de dados"


def build_structured_action_plan(area: str, quality_report: dict, issue_catalog: pd.DataFrame) -> pd.DataFrame:
    rows = []
    priority_rank = {"critical": "P0", "high": "P1", "medium": "P2", "low": "P3"}

    if len(issue_catalog) > 0:
        sorted_issues = issue_catalog.sort_values(["priority_rank", "affected_rows"], ascending=[True, False]).head(6)
        for _, row in sorted_issues.iterrows():
            pri = priority_rank.get(str(row.get("priority", "medium")).lower(), "P2")
            issue = clean_report_text(row.get("issue", "-"), 80)
            column = clean_report_text(row.get("column", "-"), 80)
            affected = int(row.get("affected_rows", 0) or 0)
            pct = row.get("affected_pct", None)
            pct_txt = f" ({float(pct):.1f}%)" if pct is not None and pd.notna(pct) else ""
            rows.append(
                {
                    "prioridade": pri,
                    "problema": f"{issue} - {column}",
                    "foco": f"{issue} - {column}",
                    "impacto": f"{affected:,} registros afetados{pct_txt}".replace(",", "."),
                    "acao": clean_report_text(row.get("suggestion", "Corrigir a causa raiz na origem dos dados."), 140),
                    "responsavel": suggested_action_owner(str(row.get("category", "")), area),
                    "prazo": "0-7 dias" if pri in ["P0", "P1"] else "8-30 dias",
                }
            )

    area_actions = {
        "Geral": ("Governanca de dados", "Reduz retrabalho e aumenta confianca gerencial.", "Criar checklist de qualidade antes da publicacao de dados."),
        "Financeiro": ("Controle financeiro", "Reduz risco de valores incorretos e auditoria manual.", "Validar limites, negativos e datas futuras em metricas sensiveis."),
        "Vendas": ("Cadastro e receita", "Melhora leitura de funil, carteira e ticket.", "Padronizar cadastro de clientes e acompanhar ticket medio por periodo."),
        "Operacoes": ("SLA e atendimento", "Ajuda a priorizar gargalos e reincidencias.", "Padronizar status, falhas e datas de atendimento."),
        "RH": ("Cadastro de pessoas", "Evita inconsistencias em headcount e movimentacoes.", "Padronizar campos obrigatorios e datas de admissao/desligamento."),
    }
    foco, impacto, acao = area_actions.get(area, area_actions["Geral"])
    rows.append(
        {
            "prioridade": "P2",
            "problema": foco,
            "foco": foco,
            "impacto": impacto,
            "acao": acao,
            "responsavel": suggested_action_owner("governanca", area),
            "prazo": "8-30 dias",
        }
    )

    for action in quality_report.get("actions", [])[:3]:
        rows.append(
            {
                "prioridade": "P2",
                "problema": "Qualidade de dados",
                "foco": "Qualidade de dados",
                "impacto": "Mantem os 5 pilares de qualidade sob controle.",
                "acao": clean_report_text(action, 140),
                "responsavel": "Analista de dados + Dono do dado",
                "prazo": "Recorrente",
            }
        )

    dedup = []
    seen = set()
    for row in rows:
        key = (row["prioridade"], row["foco"], row["acao"])
        if key not in seen:
            dedup.append(row)
            seen.add(key)
    return pd.DataFrame(dedup[:8])


def compute_column_quality(df: pd.DataFrame, violations: pd.DataFrame, mixed_cols: List[str]) -> pd.DataFrame:
    rows = max(len(df), 1)
    vio_map = violations.groupby("column").size().to_dict() if len(violations) > 0 else {}
    result = []
    for col in df.columns:
        missing_pct = float(df[col].isna().mean() * 100)
        unique_ratio = float(df[col].nunique(dropna=True) / max(df[col].notna().sum(), 1))
        constant_penalty = 25.0 if df[col].nunique(dropna=False) <= 1 else 0.0
        mixed_penalty = 20.0 if col in mixed_cols else 0.0
        vio = int(vio_map.get(col, 0))
        vio_pct = (vio / rows) * 100
        score = 100.0 - min(50.0, missing_pct * 0.5) - constant_penalty - mixed_penalty - min(30.0, vio_pct * 0.7)
        score -= 5.0 if unique_ratio >= 0.99 else 0.0
        score = max(0.0, min(100.0, score))
        status = "Verde" if score >= 85 else "Amarelo" if score >= 70 else "Vermelho"
        result.append({"column": col, "score": round(score, 2), "status": status, "missing_pct": round(missing_pct, 2), "violations": vio})
    return pd.DataFrame(result).sort_values("score")


def build_issue_catalog(df: pd.DataFrame, quality_report: dict, violations: pd.DataFrame, missing_by_col: pd.Series) -> pd.DataFrame:
    issues = []
    rank = {"critical": 0, "high": 1, "medium": 2, "low": 3}

    for col, pct in missing_by_col.items():
        if pct >= 30:
            issues.append(
                {
                    "category": "Completude",
                    "issue": "Alta ausencia",
                    "column": col,
                    "priority": "high",
                    "priority_rank": rank["high"],
                    "affected_rows": int(df[col].isna().sum()),
                    "affected_pct": round(float(pct), 2),
                    "suggestion": "Definir imputacao ou tornar campo obrigatorio na origem.",
                }
            )

    if quality_report["duplicates"] > 0:
        pri = "high" if quality_report["duplicate_ratio"] > 0.1 else "medium"
        issues.append(
            {
                "category": "Unicidade",
                "issue": "Duplicidade de linhas",
                "column": "(linha)",
                "priority": pri,
                "priority_rank": rank[pri],
                "affected_rows": int(quality_report["duplicates"]),
                "affected_pct": round(quality_report["duplicate_ratio"] * 100, 2),
                "suggestion": "Aplicar chave unica e deduplicacao.",
            }
        )

    for col in quality_report["mixed_cols"][:8]:
        issues.append(
            {
                "category": "Consistencia",
                "issue": "Tipo misto",
                "column": col,
                "priority": "medium",
                "priority_rank": rank["medium"],
                "affected_rows": int(df[col].notna().sum()),
                "affected_pct": None,
                "suggestion": "Padronizar formato em ETL e validacao de entrada.",
            }
        )

    if len(violations) > 0:
        grouped = (
            violations.groupby(["column", "rule", "priority"], dropna=False)
            .size()
            .reset_index(name="count")
            .sort_values("count", ascending=False)
        )
        for _, row in grouped.head(50).iterrows():
            pri = str(row["priority"])
            issues.append(
                {
                    "category": "Validade",
                    "issue": f"Regra {row['rule']} violada",
                    "column": row["column"],
                    "priority": pri,
                    "priority_rank": rank.get(pri, 2),
                    "affected_rows": int(row["count"]),
                    "affected_pct": round((int(row["count"]) / max(len(df), 1)) * 100, 2),
                    "suggestion": "Corrigir registros invalidados pela regra configurada.",
                }
            )

    if not issues:
        issues.append(
            {
                "category": "Geral",
                "issue": "Nenhuma inconsistencia critica detectada",
                "column": "-",
                "priority": "low",
                "priority_rank": 3,
                "affected_rows": 0,
                "affected_pct": 0.0,
                "suggestion": "Manter monitoramento continuo.",
            }
        )

    return pd.DataFrame(issues)


def quality_status(score: float) -> str:
    if score >= 85:
        return "Verde"
    if score >= 70:
        return "Amarelo"
    return "Vermelho"


def render_quality_overview(quality_report: dict, groups: dict[str, List[str]], key_prefix: str = "quality") -> None:
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Score Qualidade", f"{quality_report['score']:.1f}/100")
    c2.metric("Semaforo", quality_status(quality_report["score"]))
    c3.metric("Duplicados", f"{quality_report['duplicates']:,}".replace(",", "."))
    c4.metric("Outlier Ratio", f"{quality_report['outlier_ratio'] * 100:.2f}%")

    gauge = go.Figure(
        go.Indicator(
            mode="gauge+number",
            value=quality_report["score"],
            title={"text": "Indice de Qualidade"},
            gauge={
                "axis": {"range": [0, 100]},
                "bar": {"color": "#16A34A" if quality_report["score"] >= 85 else "#D97706" if quality_report["score"] >= 70 else "#DC2626"},
                "steps": [
                    {"range": [0, 50], "color": "#FEE2E2"},
                    {"range": [50, 70], "color": "#FEF3C7"},
                    {"range": [70, 85], "color": "#DBEAFE"},
                    {"range": [85, 100], "color": "#DCFCE7"},
                ],
            },
        )
    )
    gauge.update_layout(height=280, margin=dict(l=20, r=20, t=40, b=20))

    pillars_df = pd.DataFrame({"Pilar": list(quality_report["pillars"].keys()), "Score": list(quality_report["pillars"].values())})
    fig_pillars = px.bar(
        pillars_df,
        x="Pilar",
        y="Score",
        color="Score",
        title="Score por Pilar",
        range_color=[0, 100],
        color_continuous_scale=["#DC2626", "#D97706", "#16A34A"],
    )

    left, right = st.columns(2)
    with left:
        st.plotly_chart(gauge, use_container_width=True, key=f"{key_prefix}_gauge")
    with right:
        st.plotly_chart(fig_pillars, use_container_width=True, key=f"{key_prefix}_pillars")

    type_df = pd.DataFrame(
        {
            "Tipo": ["Numericas", "Categoricas", "Datas"],
            "Quantidade": [len(groups["numeric"]), len(groups["categorical"]), len(groups["datetime"])],
        }
    )
    fig_types = px.pie(type_df, names="Tipo", values="Quantidade", title="Composicao de Tipos")
    st.plotly_chart(fig_types, use_container_width=True, key=f"{key_prefix}_types")


def render_global_filters(df: pd.DataFrame, groups: dict[str, List[str]]) -> Tuple[pd.DataFrame, List[str]]:
    st.sidebar.markdown("---")
    st.sidebar.subheader("Filtros Globais")
    filter_context = hash_text("|".join(map(str, df.columns)))[:12]
    cols = st.sidebar.multiselect(
        "Selecionar colunas para filtrar",
        options=df.columns.tolist(),
        default=[],
        key=f"global_filter_columns_{filter_context}",
    )
    filtered = df.copy()
    active = []

    for col in cols:
        if col in groups["numeric"]:
            non_null = filtered[col].dropna()
            if len(non_null) == 0:
                continue
            min_v = float(non_null.min())
            max_v = float(non_null.max())
            if min_v == max_v:
                continue
            chosen = st.sidebar.slider(
                f"{col} (intervalo)",
                min_value=min_v,
                max_value=max_v,
                value=(min_v, max_v),
                key=f"flt_num_{filter_context}_{col}",
            )
            filtered = filtered[(filtered[col].isna()) | ((filtered[col] >= chosen[0]) & (filtered[col] <= chosen[1]))]
            active.append(f"{col}: {chosen[0]:.2f} ate {chosen[1]:.2f}")
        elif col in groups["datetime"]:
            valid = pd.to_datetime(filtered[col], errors="coerce", dayfirst=True).dropna()
            if len(valid) == 0:
                continue
            dmin, dmax = valid.min().date(), valid.max().date()
            chosen = st.sidebar.date_input(f"{col} (periodo)", value=(dmin, dmax), key=f"flt_dt_{filter_context}_{col}")
            if isinstance(chosen, tuple) and len(chosen) == 2:
                start, end = chosen
                dt = pd.to_datetime(filtered[col], errors="coerce", dayfirst=True)
                filtered = filtered[(dt.isna()) | ((dt.dt.date >= start) & (dt.dt.date <= end))]
                active.append(f"{col}: {start} ate {end}")
        else:
            values = filtered[col].fillna("NULO").astype(str)
            opts = values.value_counts().head(80).index.tolist()
            selected = st.sidebar.multiselect(
                f"{col} (categorias)",
                options=opts,
                default=opts,
                key=f"flt_cat_{filter_context}_{col}",
            )
            if selected:
                filtered = filtered[values.isin(selected)]
                active.append(f"{col}: {len(selected)} categorias")
            else:
                filtered = filtered.iloc[0:0].copy()
                active.append(f"{col}: nenhuma categoria selecionada")

    return filtered, active


def render_data_treatment_controls() -> dict:
    st.sidebar.markdown("---")
    st.sidebar.subheader("Tratativa de Dados")
    remove_empty_rows = st.sidebar.checkbox("Excluir linhas totalmente vazias", value=True, key="treat_empty_rows")
    remove_duplicates = st.sidebar.checkbox("Excluir linhas duplicadas exatas", value=True, key="treat_duplicates")
    remove_outliers = st.sidebar.checkbox("Excluir outliers (IQR)", value=True, key="treat_outliers")

    outlier_iqr_multiplier = DEFAULT_OUTLIER_IQR_MULTIPLIER
    outlier_min_hits = 1
    max_outlier_drop_pct = DEFAULT_MAX_OUTLIER_DROP_PCT
    min_rows_after_treatment = DEFAULT_MIN_ROWS_AFTER_TREATMENT

    if remove_outliers:
        outlier_iqr_multiplier = st.sidebar.slider(
            "Sensibilidade de outlier (IQR)",
            min_value=1.0,
            max_value=3.0,
            value=DEFAULT_OUTLIER_IQR_MULTIPLIER,
            step=0.1,
            key="treat_outlier_iqr_multiplier",
        )
        outlier_min_hits = st.sidebar.slider(
            "Minimo de colunas anomalas por linha",
            min_value=1,
            max_value=3,
            value=1,
            step=1,
            key="treat_outlier_min_hits",
        )
        max_outlier_drop_pct = (
            st.sidebar.slider(
                "Limite maximo para remocao de outliers (%)",
                min_value=5,
                max_value=40,
                value=int(DEFAULT_MAX_OUTLIER_DROP_PCT * 100),
                step=1,
                key="treat_max_outlier_drop_pct",
            )
            / 100.0
        )
        min_rows_after_treatment = st.sidebar.number_input(
            "Minimo de linhas apos tratativa",
            min_value=20,
            max_value=5000,
            value=DEFAULT_MIN_ROWS_AFTER_TREATMENT,
            step=10,
            key="treat_min_rows_after",
        )

    st.sidebar.caption("A tratativa e aplicada antes da analise, do ML e dos dashboards.")
    return {
        "remove_empty_rows": bool(remove_empty_rows),
        "remove_duplicates": bool(remove_duplicates),
        "remove_outliers": bool(remove_outliers),
        "outlier_iqr_multiplier": float(outlier_iqr_multiplier),
        "outlier_min_hits": int(outlier_min_hits),
        "max_outlier_drop_pct": float(max_outlier_drop_pct),
        "min_rows_after_treatment": int(min_rows_after_treatment),
    }


def normalize_missing_markers(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    markers = {"", "-", "--", "n/a", "na", "null", "none", "nan", "sem informacao"}
    for col in out.select_dtypes(include=["object"]).columns:
        txt = out[col].astype(str).str.strip()
        mask = txt.str.lower().isin(markers)
        if mask.any():
            out.loc[mask, col] = np.nan
    return out


def detect_outlier_rows_iqr(
    df: pd.DataFrame,
    numeric_cols: List[str],
    *,
    iqr_multiplier: float,
    min_hits: int,
) -> Tuple[pd.Series, List[dict]]:
    if not numeric_cols or len(df) == 0:
        return pd.Series(False, index=df.index), []

    hits = pd.Series(0, index=df.index, dtype="int32")
    details: List[dict] = []

    for col in numeric_cols:
        s = df[col]
        non_null = s.dropna()
        if len(non_null) < 12 or non_null.nunique(dropna=True) < 4:
            continue
        q1 = non_null.quantile(0.25)
        q3 = non_null.quantile(0.75)
        iqr = q3 - q1
        if pd.isna(iqr) or iqr <= 0:
            continue
        lower = q1 - iqr_multiplier * iqr
        upper = q3 + iqr_multiplier * iqr
        mask = s.notna() & ((s < lower) | (s > upper))
        count = int(mask.sum())
        if count > 0:
            hits += mask.astype("int32")
            details.append(
                {
                    "coluna": col,
                    "outliers_detectados": count,
                    "limite_inferior": float(lower),
                    "limite_superior": float(upper),
                }
            )

    row_mask = hits >= max(1, int(min_hits))
    details = sorted(details, key=lambda x: x["outliers_detectados"], reverse=True)
    return row_mask, details


def apply_data_treatment(df: pd.DataFrame, options: dict) -> Tuple[pd.DataFrame, dict]:
    base = normalize_missing_markers(df)
    treated = base.copy()
    warnings: List[str] = []

    report = {
        "rows_before": int(len(base)),
        "rows_after": int(len(base)),
        "empty_rows_removed": 0,
        "duplicates_removed": 0,
        "outliers_removed": 0,
        "outlier_candidates": 0,
        "outlier_columns": [],
        "total_removed": 0,
        "missing_ratio_before": float(base.isna().sum().sum() / max(base.size, 1)),
        "missing_ratio_after": float(base.isna().sum().sum() / max(base.size, 1)),
        "warnings": warnings,
    }

    if options.get("remove_empty_rows", True):
        empty_mask = treated.isna().all(axis=1)
        report["empty_rows_removed"] = int(empty_mask.sum())
        treated = treated.loc[~empty_mask].copy()

    if options.get("remove_duplicates", True) and len(treated) > 0:
        dup_count = int(treated.duplicated(keep="first").sum())
        report["duplicates_removed"] = dup_count
        treated = treated.drop_duplicates(keep="first").copy()

    if options.get("remove_outliers", True) and len(treated) > 0:
        id_like_cols = set(detect_id_like_columns(treated))
        numeric_cols = [c for c in treated.select_dtypes(include=np.number).columns.tolist() if c not in id_like_cols]
        row_mask, details = detect_outlier_rows_iqr(
            treated,
            numeric_cols,
            iqr_multiplier=float(options.get("outlier_iqr_multiplier", DEFAULT_OUTLIER_IQR_MULTIPLIER)),
            min_hits=int(options.get("outlier_min_hits", 1)),
        )
        outlier_candidates = int(row_mask.sum())
        report["outlier_candidates"] = outlier_candidates
        report["outlier_columns"] = details[:12]

        max_drop = int(len(treated) * float(options.get("max_outlier_drop_pct", DEFAULT_MAX_OUTLIER_DROP_PCT)))
        min_rows_after = int(options.get("min_rows_after_treatment", DEFAULT_MIN_ROWS_AFTER_TREATMENT))
        remaining = int(len(treated) - outlier_candidates)

        if outlier_candidates > 0 and outlier_candidates <= max_drop and remaining >= min_rows_after:
            report["outliers_removed"] = outlier_candidates
            treated = treated.loc[~row_mask].copy()
        elif outlier_candidates > 0:
            warnings.append(
                "Outliers detectados, mas a remocao foi ignorada para preservar volume minimo e evitar perda excessiva."
            )

    report["rows_after"] = int(len(treated))
    report["total_removed"] = int(report["rows_before"] - report["rows_after"])
    report["missing_ratio_after"] = float(treated.isna().sum().sum() / max(treated.size, 1))
    treated = treated.reset_index(drop=True)
    return treated, report


def render_treatment_summary(report: dict, key_prefix: str = "treatment") -> None:
    st.markdown("**Impacto da tratativa de dados**")
    c1, c2, c3, c4, c5 = st.columns(5)
    c1.metric("Linhas originais", f"{report['rows_before']:,}".replace(",", "."))
    c2.metric("Linhas tratadas", f"{report['rows_after']:,}".replace(",", "."))
    c3.metric("Vazias removidas", f"{report['empty_rows_removed']:,}".replace(",", "."))
    c4.metric("Duplicadas removidas", f"{report['duplicates_removed']:,}".replace(",", "."))
    c5.metric("Outliers removidos", f"{report['outliers_removed']:,}".replace(",", "."))

    impact_df = pd.DataFrame(
        [
            {"Etapa": "Linhas vazias", "Removidas": int(report["empty_rows_removed"])},
            {"Etapa": "Duplicadas", "Removidas": int(report["duplicates_removed"])},
            {"Etapa": "Outliers", "Removidas": int(report["outliers_removed"])},
        ]
    )
    st.plotly_chart(
        px.bar(impact_df, x="Etapa", y="Removidas", text="Removidas", title="Remocoes por etapa"),
        use_container_width=True,
        key=f"{key_prefix}_impact",
    )

    if report.get("outlier_columns"):
        st.caption("Colunas com mais outliers detectados (IQR):")
        st.dataframe(pd.DataFrame(report["outlier_columns"]), use_container_width=True, height=180)
    if report.get("warnings"):
        for item in report["warnings"]:
            st.warning(item)


def format_kpi_metric(item: dict) -> str:
    value = float(item.get("metric", 0.0))
    if item.get("value_type") == "count":
        return f"{int(round(value)):,}".replace(",", ".")
    return f"{value:.2f}"


def format_br_number(value: object, decimals: int = 2) -> str:
    try:
        num = float(value)
    except Exception:
        return clean_report_text(value, 80)
    formatted = f"{num:,.{decimals}f}"
    return formatted.replace(",", "X").replace(".", ",").replace("X", ".")


def apply_inline_filters(
    df: pd.DataFrame,
    groups: dict[str, List[str]],
    *,
    key_prefix: str,
    title: str,
    expanded: bool = False,
    max_categories: int = 80,
) -> Tuple[pd.DataFrame, List[str]]:
    filtered = df.copy()
    active_filters: List[str] = []
    metric_numeric, _ = select_numeric_metric_columns(df, groups)

    with st.expander(title, expanded=expanded):
        c1, c2, c3 = st.columns(3)
        date_col = c1.selectbox(
            "Filtro por data",
            options=["(sem filtro)"] + groups["datetime"],
            key=f"{key_prefix}_date_col",
        )
        cat_col = c2.selectbox(
            "Filtro por categoria",
            options=["(sem filtro)"] + groups["categorical"],
            key=f"{key_prefix}_cat_col",
        )
        num_col = c3.selectbox(
            "Filtro por metrica numerica",
            options=["(sem filtro)"] + metric_numeric,
            key=f"{key_prefix}_num_col",
        )

        if date_col != "(sem filtro)":
            dt = pd.to_datetime(filtered[date_col], errors="coerce", dayfirst=True)
            valid = dt.dropna()
            if len(valid) > 0:
                dmin, dmax = valid.min().date(), valid.max().date()
                chosen = st.date_input(
                    "Periodo",
                    value=(dmin, dmax),
                    key=f"{key_prefix}_date_range_{normalize_token_text(date_col)}",
                )
                if isinstance(chosen, tuple) and len(chosen) == 2:
                    start, end = chosen
                    filtered = filtered[(dt.isna()) | ((dt.dt.date >= start) & (dt.dt.date <= end))]
                    active_filters.append(f"{date_col}: {start} ate {end}")

        if cat_col != "(sem filtro)":
            values = filtered[cat_col].fillna("NULO").astype(str)
            opts = values.value_counts().head(max_categories).index.tolist()
            selected = st.multiselect(
                "Categorias selecionadas",
                options=opts,
                default=opts,
                key=f"{key_prefix}_cat_values_{normalize_token_text(cat_col)}",
            )
            if selected:
                filtered = filtered[values.isin(selected)]
                active_filters.append(f"{cat_col}: {len(selected)} categorias")
            else:
                filtered = filtered.iloc[0:0].copy()
                active_filters.append(f"{cat_col}: nenhuma categoria selecionada")

        if num_col != "(sem filtro)" and not filtered.empty:
            non_null = pd.to_numeric(filtered[num_col], errors="coerce").dropna()
            if len(non_null) > 0:
                min_v = float(non_null.min())
                max_v = float(non_null.max())
                if min_v != max_v:
                    chosen = st.slider(
                        "Intervalo numerico",
                        min_value=min_v,
                        max_value=max_v,
                        value=(min_v, max_v),
                        key=f"{key_prefix}_num_range_{normalize_token_text(num_col)}",
                    )
                    numeric_series = pd.to_numeric(filtered[num_col], errors="coerce")
                    filtered = filtered[(numeric_series.isna()) | ((numeric_series >= chosen[0]) & (numeric_series <= chosen[1]))]
                    active_filters.append(f"{num_col}: {chosen[0]:.2f} ate {chosen[1]:.2f}")

        st.caption(f"Registros no recorte: {len(filtered):,}".replace(",", "."))

    return filtered, active_filters


def apply_dashboard_filters(
    df: pd.DataFrame,
    groups: dict[str, List[str]],
    *,
    key_prefix: str = "dash",
) -> Tuple[pd.DataFrame, List[str]]:
    return apply_inline_filters(
        df,
        groups,
        key_prefix=f"{key_prefix}_flt",
        title="Filtros do Dashboard",
        expanded=False,
        max_categories=40,
    )


def build_dashboard_narrative(
    df: pd.DataFrame,
    groups: dict[str, List[str]],
    quality_report: dict,
    issue_catalog: pd.DataFrame,
    kpis: List[dict],
    active_filters: List[str],
    volume_col: Optional[str],
    spreadsheet_profile: Optional[dict] = None,
) -> List[str]:
    notes = [
        f"A base atual contem {len(df):,} registros e {df.shape[1]} colunas.".replace(",", "."),
        f"O score de qualidade esta em {quality_report['score']:.1f}/100 ({quality_report['level']}).",
    ]
    if spreadsheet_profile:
        notes.insert(
            0,
            f"Tipo detectado: `{spreadsheet_profile.get('tipo', '-')}` "
            f"({float(spreadsheet_profile.get('confianca', 0.0)):.1f}% de confianca). "
            f"Melhor abordagem: {spreadsheet_profile.get('melhor_analise', '-')}",
        )

    if active_filters:
        notes.append("Esta leitura considera filtros ativos no dashboard: " + " | ".join(active_filters[:3]) + ".")

    if volume_col:
        notes.append(f"Foram identificados {df[volume_col].nunique(dropna=True):,} IDs unicos em `{volume_col}`.".replace(",", "."))

    if groups["datetime"]:
        date_col = max(groups["datetime"], key=lambda c: df[c].notna().sum())
        dt = pd.to_datetime(df[date_col], errors="coerce", dayfirst=True).dropna()
        if len(dt) > 0:
            notes.append(
                f"Cobertura temporal de `{date_col}`: {dt.min().strftime('%d/%m/%Y')} ate {dt.max().strftime('%d/%m/%Y')}."
            )

    if groups["categorical"]:
        cat_col = max(groups["categorical"], key=lambda c: df[c].notna().sum())
        vc = df[cat_col].fillna("NULO").astype(str).value_counts(normalize=True)
        if len(vc) > 0:
            top_cat = vc.index[0]
            top_share = float(vc.iloc[0]) * 100.0
            notes.append(f"A categoria dominante em `{cat_col}` e `{top_cat}`, representando {top_share:.1f}% dos registros.")

    strongest = strongest_correlations(df, top_n=1)
    if len(strongest) > 0:
        row = strongest.iloc[0]
        notes.append(
            f"A relacao numerica mais forte ocorre entre `{row['variavel_1']}` e `{row['variavel_2']}` "
            f"(correlacao {row['correlacao']:.3f})."
        )

    if kpis:
        notes.append(f"KPI em destaque: `{kpis[0]['name']}` = {format_kpi_metric(kpis[0])}.")

    if len(issue_catalog) > 0:
        top_risk = issue_catalog.sort_values(["priority_rank", "affected_rows"], ascending=[True, False]).iloc[0]
        notes.append(
            f"Maior risco atual: `{top_risk['issue']}` na coluna `{top_risk['column']}` ({top_risk['affected_rows']} registros afetados)."
        )

    return notes[:8]


ANALYTIC_DOMAIN_PROFILES = [
    {
        "tipo": "Atendimento / Chamados / Suporte",
        "tokens": ["chamado", "ticket", "protocolo", "atendimento", "solicitacao", "status", "sla", "motivo", "causa", "problema", "prioridade", "encerramento", "abertura"],
        "foco": "volume, status, fila, prioridade, tempo de atendimento, reincidencia e causa raiz.",
        "analises": [
            "Contar chamados/IDs unicos, nunca somar codigos.",
            "Analisar status, motivos, cidade/equipe e evolucao temporal.",
            "Priorizar gargalos por Pareto e verificar SLA ou tempo de resolucao quando existir.",
        ],
        "graficos": ["KPI de volume", "Pizza/status", "Pareto por motivo/cidade/equipe", "Serie temporal de chamados", "Boxplot de tempo/SLA"],
        "dim_tokens": ["status", "motivo", "causa", "cidade", "equipe", "fila", "canal", "prioridade", "tipo"],
        "metric_tokens": ["tempo", "sla", "prazo", "duracao", "dias", "horas", "nota"],
    },
    {
        "tipo": "Financeiro / Faturamento / Cobranca",
        "tokens": ["valor", "receita", "fatura", "boleto", "pagamento", "vencimento", "cobranca", "inadimplencia", "imposto", "custo", "despesa", "saldo", "liquido", "bruto"],
        "foco": "valores, receitas/custos, vencimentos, inadimplencia, concentracao financeira e variacao no tempo.",
        "analises": [
            "Somar valores financeiros reais e separar IDs/codigos de metricas.",
            "Comparar valores por periodo, categoria, cliente, cidade e status de pagamento.",
            "Detectar outliers, negativos, vencidos e concentracao de receita/custo.",
        ],
        "graficos": ["Cards de valor total/media", "Ranking por categoria/cliente", "Serie temporal financeira", "Distribuicao de valores", "Correlacao entre metricas"],
        "dim_tokens": ["status", "cliente", "cidade", "categoria", "tipo", "conta", "centro", "vendedor"],
        "metric_tokens": ["valor", "receita", "custo", "despesa", "saldo", "liquido", "bruto", "imposto", "multa"],
    },
    {
        "tipo": "Vendas / Comercial / Clientes",
        "tokens": ["venda", "pedido", "cliente", "produto", "servico", "contrato", "vendedor", "funil", "lead", "proposta", "ticket", "receita", "cidade", "segmento"],
        "foco": "volume de vendas, carteira, produtos, funil, ticket medio, canais e desempenho comercial.",
        "analises": [
            "Avaliar volume por vendedor/produto/cidade/canal.",
            "Medir ticket medio, receita, conversao e concentracao por cliente ou segmento.",
            "Comparar evolucao temporal e identificar oportunidades de crescimento.",
        ],
        "graficos": ["Ranking de vendas", "Pareto de produto/cliente", "Serie temporal", "Distribuicao de ticket", "Mapa de concentracao por segmento"],
        "dim_tokens": ["cliente", "produto", "servico", "vendedor", "cidade", "canal", "segmento", "status"],
        "metric_tokens": ["valor", "receita", "ticket", "quantidade", "qtd", "desconto"],
    },
    {
        "tipo": "RH / Pessoas / Jornada",
        "tokens": ["colaborador", "funcionario", "matricula", "cargo", "setor", "departamento", "admissao", "demissao", "salario", "ferias", "absenteismo", "treinamento"],
        "foco": "headcount, movimentacoes, cargos, setores, custos, presenca e desenvolvimento.",
        "analises": [
            "Contar pessoas/matriculas unicas e segmentar por setor/cargo/local.",
            "Analisar admissoes/desligamentos por periodo quando houver datas.",
            "Monitorar ausencias, custos e concentracao por area.",
        ],
        "graficos": ["Headcount por setor", "Distribuicao por cargo", "Serie temporal de movimentacoes", "Boxplot salarial", "Qualidade cadastral"],
        "dim_tokens": ["setor", "departamento", "cargo", "cidade", "unidade", "status"],
        "metric_tokens": ["salario", "idade", "horas", "dias", "custo", "nota"],
    },
    {
        "tipo": "Estoque / Inventario / Compras",
        "tokens": ["estoque", "produto", "sku", "item", "quantidade", "qtd", "entrada", "saida", "compra", "fornecedor", "almoxarifado", "lote", "validade"],
        "foco": "saldo, giro, ruptura, fornecedores, itens criticos e valor parado.",
        "analises": [
            "Avaliar quantidade/saldo por produto, fornecedor e local.",
            "Identificar itens sem giro, baixo estoque, outliers e vencimentos.",
            "Separar codigos/SKU de metricas reais de quantidade e valor.",
        ],
        "graficos": ["Ranking de itens", "Pareto por fornecedor", "Distribuicao de saldo", "Serie de entradas/saidas", "Alertas de baixo estoque"],
        "dim_tokens": ["produto", "sku", "item", "fornecedor", "local", "categoria", "lote"],
        "metric_tokens": ["quantidade", "qtd", "saldo", "valor", "custo", "entrada", "saida"],
    },
    {
        "tipo": "Pesquisa / Satisfacao / Qualidade percebida",
        "tokens": ["nota", "nps", "csat", "satisfacao", "avaliacao", "pesquisa", "resposta", "feedback", "comentario", "canal", "promotor", "detrator"],
        "foco": "notas, experiencia, sentimento operacional, canais, motivos e evolucao da satisfacao.",
        "analises": [
            "Medir notas medias/medianas e distribuicao por faixa.",
            "Segmentar por canal, cidade, equipe, produto ou motivo.",
            "Identificar detratores, concentracoes e mudancas temporais.",
        ],
        "graficos": ["Distribuicao de notas", "Ranking por canal/cidade", "Serie temporal da nota", "Pareto de motivos", "Boxplot por segmento"],
        "dim_tokens": ["canal", "cidade", "equipe", "produto", "motivo", "tipo", "status"],
        "metric_tokens": ["nota", "nps", "csat", "satisfacao", "score"],
    },
    {
        "tipo": "Qualidade / Auditoria / Conformidade",
        "tokens": ["qualidade", "auditoria", "conformidade", "nao_conformidade", "nc", "erro", "falha", "risco", "criticidade", "plano", "acao", "evidencia"],
        "foco": "nao conformidades, riscos, criticidade, recorrencia, responsaveis e plano de acao.",
        "analises": [
            "Contar ocorrencias e priorizar por criticidade/impacto.",
            "Analisar recorrencia por area, responsavel, processo e causa.",
            "Monitorar prazos, status do plano de acao e evolucao temporal.",
        ],
        "graficos": ["Riscos por criticidade", "Pareto de causas", "Status do plano", "Serie temporal de NCs", "Ranking por area"],
        "dim_tokens": ["criticidade", "risco", "status", "area", "responsavel", "causa", "processo"],
        "metric_tokens": ["prazo", "dias", "impacto", "score", "nota"],
    },
]


def score_column_by_tokens(col: str, preferred_tokens: List[str]) -> int:
    col_key = normalize_token_text(col)
    tokens = set(col_key.split("_"))
    score = 0
    for i, token in enumerate(preferred_tokens):
        weight = max(1, len(preferred_tokens) - i)
        token_key = normalize_token_text(token)
        if token_key in tokens or token_key in col_key:
            score += weight
    return score


def pick_recommended_dimension(df: pd.DataFrame, groups: dict[str, List[str]], preferred_tokens: List[str]) -> Optional[str]:
    candidates = []
    id_like = set(detect_id_like_columns(df))
    for col in groups.get("categorical", []):
        if col in id_like:
            continue
        non_null = df[col].dropna()
        if len(non_null) == 0:
            continue
        unique_count = int(non_null.nunique(dropna=True))
        if unique_count < 2 or unique_count > max(120, int(len(non_null) * 0.6)):
            continue
        score = score_column_by_tokens(col, preferred_tokens)
        coverage = float(len(non_null) / max(len(df), 1))
        balance_bonus = 5 if 3 <= unique_count <= 30 else 2
        candidates.append((score, coverage, balance_bonus, -unique_count, col))
    if not candidates:
        return groups.get("categorical", [None])[0] if groups.get("categorical") else None
    return sorted(candidates, reverse=True)[0][-1]


def pick_recommended_metric(df: pd.DataFrame, metric_numeric: List[str], preferred_tokens: List[str]) -> Optional[str]:
    candidates = []
    for col in metric_numeric:
        non_null = pd.to_numeric(df[col], errors="coerce").dropna()
        if len(non_null) == 0:
            continue
        variance = float(non_null.var()) if pd.notna(non_null.var()) else 0.0
        candidates.append((score_column_by_tokens(col, preferred_tokens), len(non_null), variance, col))
    if not candidates:
        return None
    return sorted(candidates, reverse=True)[0][-1]


def detect_spreadsheet_profile(
    df: pd.DataFrame,
    groups: dict[str, List[str]],
    quality_report: Optional[dict] = None,
) -> dict:
    metric_numeric, id_like_cols = select_numeric_metric_columns(df, groups)
    volume_col = pick_volume_id_column(df, id_like_cols)
    normalized_cols = {col: normalize_token_text(col) for col in df.columns}
    profile_rows = []

    for profile in ANALYTIC_DOMAIN_PROFILES:
        score = 0
        matched_cols: List[str] = []
        for col, col_key in normalized_cols.items():
            col_tokens = set(col_key.split("_"))
            for token in profile["tokens"]:
                token_key = normalize_token_text(token)
                if token_key in col_tokens or token_key in col_key:
                    score += 3 if token_key in col_tokens else 2
                    matched_cols.append(col)
                    break
        if groups.get("datetime"):
            score += 2
        if metric_numeric:
            score += 2
        if volume_col and any(tok in normalize_token_text(volume_col) for tok in profile["tokens"]):
            score += 3
        profile_rows.append((score, profile, sorted(set(matched_cols))))

    best_score, best_profile, matched_cols = sorted(profile_rows, key=lambda item: item[0], reverse=True)[0]
    max_possible = max(10, min(36, len(df.columns) * 3))
    confidence = max(0.15, min(0.98, best_score / max_possible))
    if best_score < 5:
        best_profile = {
            "tipo": "Planilha geral / base mista",
            "tokens": [],
            "foco": "entender estrutura, qualidade, principais categorias, metricas numericas e possibilidade de tendencia temporal.",
            "analises": [
                "Comecar por qualidade, volumetria, categorias dominantes e metricas numericas reais.",
                "Usar contagem para IDs/codigos e soma/media apenas para metricas de negocio.",
                "Aplicar filtros para descobrir recortes relevantes antes de conclusoes gerenciais.",
            ],
            "graficos": ["Resumo de qualidade", "Top categorias", "Distribuicao numerica", "Serie temporal quando houver data", "Correlacao quando houver metricas"],
            "dim_tokens": ["status", "categoria", "tipo", "cidade", "area", "produto", "cliente"],
            "metric_tokens": ["valor", "quantidade", "qtd", "nota", "tempo", "dias"],
        }
        matched_cols = []
        confidence = 0.35

    recommended_dimension = pick_recommended_dimension(df, groups, best_profile.get("dim_tokens", []))
    recommended_metric = pick_recommended_metric(df, metric_numeric, best_profile.get("metric_tokens", []))
    recommended_date = max(groups["datetime"], key=lambda c: df[c].notna().sum()) if groups.get("datetime") else None

    contains = [
        f"{len(df):,} registros".replace(",", "."),
        f"{df.shape[1]} colunas",
        f"{len(groups.get('numeric', []))} numericas",
        f"{len(groups.get('categorical', []))} categoricas",
        f"{len(groups.get('datetime', []))} temporais",
    ]
    if volume_col:
        contains.append(f"ID/volume principal: {volume_col}")
    if recommended_dimension:
        contains.append(f"dimensao recomendada: {recommended_dimension}")
    if recommended_metric:
        contains.append(f"metrica recomendada: {recommended_metric}")
    if quality_report:
        contains.append(f"qualidade: {float(quality_report.get('score', 0.0)):.1f}/100")

    if recommended_metric:
        best_analysis = f"Analisar `{recommended_metric}` por `{recommended_dimension}` e por periodo, usando ranking, distribuicao e tendencia."
    elif volume_col:
        best_analysis = f"Analisar contagem/IDs unicos de `{volume_col}` por `{recommended_dimension or 'categoria principal'}` e por periodo."
    else:
        best_analysis = f"Analisar volume de linhas por `{recommended_dimension or 'categoria principal'}` com foco em concentracao, qualidade e filtros."

    return {
        "tipo": best_profile["tipo"],
        "confianca": round(float(confidence) * 100.0, 1),
        "foco": best_profile["foco"],
        "contem": contains,
        "melhor_analise": best_analysis,
        "analises_recomendadas": best_profile["analises"],
        "graficos_recomendados": best_profile["graficos"],
        "colunas_detectadas": matched_cols[:12],
        "recommended_dimension": recommended_dimension,
        "recommended_metric": recommended_metric,
        "recommended_date": recommended_date,
        "volume_col": volume_col,
    }


def render_spreadsheet_profile(profile: dict, key_prefix: str) -> None:
    st.markdown("**Diagnostico inteligente da planilha**")
    c1, c2, c3 = st.columns([1.1, 0.7, 2.0])
    c1.metric("Tipo detectado", clean_report_text(profile.get("tipo", "-"), 42))
    c2.metric("Confianca", f"{float(profile.get('confianca', 0.0)):.1f}%")
    c3.info(clean_report_text(profile.get("melhor_analise", "-"), 320))

    with st.expander("O que o AutoAnalista identificou e como vai analisar", expanded=True):
        a, b, c = st.columns(3)
        with a:
            st.markdown("**O que contem**")
            for item in profile.get("contem", [])[:8]:
                st.caption(f"- {item}")
        with b:
            st.markdown("**Analises indicadas**")
            for item in profile.get("analises_recomendadas", [])[:5]:
                st.caption(f"- {item}")
        with c:
            st.markdown("**Graficos recomendados**")
            for item in profile.get("graficos_recomendados", [])[:6]:
                st.caption(f"- {item}")

        matched = profile.get("colunas_detectadas", [])
        if matched:
            st.caption("Colunas que ajudaram na deteccao: " + ", ".join([clean_report_text(c, 40) for c in matched[:10]]))


def build_numeric_distribution_table(df: pd.DataFrame, metric_numeric: List[str], preferred_metric: Optional[str] = None) -> pd.DataFrame:
    if not metric_numeric or df.empty:
        return pd.DataFrame()
    if preferred_metric in metric_numeric:
        dist_col = preferred_metric
    else:
        ranked = sorted(
            metric_numeric,
            key=lambda c: (
                df[c].notna().sum(),
                df[c].var(skipna=True) if pd.notna(df[c].var(skipna=True)) else -1,
            ),
            reverse=True,
        )
        dist_col = ranked[0] if ranked else None
    if not dist_col:
        return pd.DataFrame()
    series = pd.to_numeric(df[dist_col], errors="coerce").dropna()
    if len(series) < 3:
        return pd.DataFrame()

    if series.nunique(dropna=True) <= 12:
        counts = series.value_counts().sort_index()
        rows = [
            {
                "coluna": dist_col,
                "faixa": clean_report_text(idx, 40),
                "quantidade": int(value),
                "participacao_pct": round(float(value) / max(len(series), 1) * 100.0, 2),
            }
            for idx, value in counts.items()
        ]
        return pd.DataFrame(rows).head(12)

    bin_count = min(12, max(4, int(np.sqrt(len(series)))))
    counts, edges = np.histogram(series, bins=bin_count)
    rows = []
    for i, count in enumerate(counts):
        rows.append(
            {
                "coluna": dist_col,
                "faixa": f"{edges[i]:.2f} ate {edges[i + 1]:.2f}",
                "quantidade": int(count),
                "participacao_pct": round(float(count) / max(len(series), 1) * 100.0, 2),
            }
        )
    return pd.DataFrame(rows)


def apply_analytics_chart_layout(fig, *, title: Optional[str] = None, height: int = 420):
    if title:
        fig.update_layout(title=title)
    fig.update_layout(
        height=height,
        margin=dict(l=24, r=24, t=70, b=42),
        legend_title_text="",
        hovermode="x unified",
        font=dict(size=12),
    )
    return fig


def add_chart_reading(title: str, bullets: List[str]) -> None:
    clean_bullets = [clean_report_text(item, 260) for item in bullets if item]
    if not clean_bullets:
        return
    with st.container():
        st.markdown(f"**Leitura do grafico - {title}**")
        for item in clean_bullets[:4]:
            st.caption(f"- {item}")


def categorical_distribution(df: pd.DataFrame, cat_col: str, top_n: int = 10) -> pd.DataFrame:
    dist = df[cat_col].fillna("NULO").astype(str).value_counts().head(top_n).reset_index()
    dist.columns = [cat_col, "quantidade"]
    total = max(int(len(df)), 1)
    dist["participacao_pct"] = dist["quantidade"] / total * 100.0
    dist["acumulado_pct"] = dist["participacao_pct"].cumsum()
    return dist


def category_reading(dist: pd.DataFrame, cat_col: str) -> List[str]:
    if dist.empty:
        return []
    top = dist.iloc[0]
    top_name = str(top[cat_col])
    top_share = float(top.get("participacao_pct", 0.0))
    pareto_count = int((dist["acumulado_pct"] <= 80).sum())
    pareto_count = max(1, pareto_count)
    bullets = [
        f"`{top_name}` lidera `{cat_col}` com {top_share:.1f}% dos registros.",
        f"As {pareto_count} primeiras categorias concentram aproximadamente {float(dist.iloc[pareto_count - 1]['acumulado_pct']):.1f}% do volume exibido.",
    ]
    if top_share >= 60:
        bullets.append("Ha alta concentracao: vale investigar dependencia, gargalo ou categoria dominante no processo.")
    elif top_share >= 30:
        bullets.append("Ha concentracao moderada: priorize as maiores categorias antes de atuar em casos pulverizados.")
    else:
        bullets.append("A distribuicao esta pulverizada: a leitura deve considerar varios grupos, nao apenas o primeiro.")
    return bullets


def default_aggregation_for_profile(profile: Optional[dict], metric_col: Optional[str] = None) -> str:
    tipo = normalize_token_text((profile or {}).get("tipo", ""))
    metric_key = normalize_token_text(metric_col or "")
    sum_domains = ["financeiro", "faturamento", "cobranca", "vendas", "comercial", "estoque", "inventario"]
    sum_metrics = ["valor", "receita", "custo", "despesa", "saldo", "liquido", "bruto", "imposto", "quantidade", "qtd"]
    if any(token in tipo for token in sum_domains) and any(token in metric_key for token in sum_metrics):
        return "Soma"
    return "Contagem"


def choose_secondary_category(df: pd.DataFrame, groups: dict[str, List[str]], primary_col: Optional[str]) -> Optional[str]:
    candidates = []
    id_like = set(detect_id_like_columns(df))
    for col in groups.get("categorical", []):
        if col == primary_col or col in id_like:
            continue
        non_null = df[col].dropna()
        if len(non_null) == 0:
            continue
        unique_count = int(non_null.nunique(dropna=True))
        if 2 <= unique_count <= 20:
            candidates.append((score_column_by_tokens(col, ["status", "motivo", "cidade", "canal", "tipo", "categoria", "prioridade", "area"]), len(non_null), -unique_count, col))
    if not candidates:
        return None
    return sorted(candidates, reverse=True)[0][-1]


def render_category_cross_analysis(
    df: pd.DataFrame,
    groups: dict[str, List[str]],
    primary_col: Optional[str],
    *,
    key_prefix: str,
    top_n: int = 10,
) -> None:
    secondary_col = choose_secondary_category(df, groups, primary_col)
    if not primary_col or not secondary_col:
        return
    base = df[[primary_col, secondary_col]].dropna().copy()
    if len(base) < 5:
        return
    base[primary_col] = base[primary_col].astype(str)
    base[secondary_col] = base[secondary_col].astype(str)
    top_primary = base[primary_col].value_counts().head(top_n).index.tolist()
    top_secondary = base[secondary_col].value_counts().head(8).index.tolist()
    base = base[base[primary_col].isin(top_primary) & base[secondary_col].isin(top_secondary)]
    if base.empty:
        return

    cross = pd.crosstab(base[primary_col], base[secondary_col]).reindex(top_primary).fillna(0)
    long = cross.reset_index().melt(id_vars=primary_col, var_name=secondary_col, value_name="quantidade")
    safe_primary = normalize_token_text(primary_col)
    safe_secondary = normalize_token_text(secondary_col)
    left, right = st.columns(2)

    stacked = px.bar(
        long,
        x=primary_col,
        y="quantidade",
        color=secondary_col,
        title=f"Cruzamento: {primary_col} x {secondary_col}",
        text="quantidade",
    )
    stacked = apply_analytics_chart_layout(stacked, height=430)
    left.plotly_chart(stacked, use_container_width=True, key=f"{key_prefix}_cross_stack_{safe_primary}_{safe_secondary}")
    with left:
        add_chart_reading(
            f"Cruzamento {primary_col} x {secondary_col}",
            [
                "Este grafico cruza duas dimensoes para revelar onde o volume se concentra.",
                "Use para encontrar combinacoes criticas, como status por cidade, motivo por equipe ou prioridade por canal.",
            ],
        )

    heat = px.imshow(
        cross,
        text_auto=True,
        aspect="auto",
        title=f"Mapa de calor: quantidade por {primary_col} e {secondary_col}",
        color_continuous_scale=["#EFF6FF", "#1D4ED8"],
    )
    heat = apply_analytics_chart_layout(heat, height=430)
    right.plotly_chart(heat, use_container_width=True, key=f"{key_prefix}_cross_heat_{safe_primary}_{safe_secondary}")


def render_count_based_insight_charts(
    df: pd.DataFrame,
    groups: dict[str, List[str]],
    profile: dict,
    *,
    key_prefix: str,
) -> None:
    metric_numeric, id_like_cols = select_numeric_metric_columns(df, groups)
    volume_col = pick_volume_id_column(df, id_like_cols)
    preferred_cat = profile.get("recommended_dimension")
    cat_cols = []
    if preferred_cat in groups.get("categorical", []):
        cat_cols.append(preferred_cat)
    for col in groups.get("categorical", []):
        if col not in cat_cols and 2 <= df[col].nunique(dropna=True) <= 50:
            cat_cols.append(col)
        if len(cat_cols) >= 3:
            break

    if cat_cols:
        st.markdown("**Rankings por contagem (sem somar IDs/codigos)**")
        cols = st.columns(len(cat_cols))
        for i, col in enumerate(cat_cols):
            dist = categorical_distribution(df, col, top_n=10)
            fig = px.bar(
                dist,
                x=col,
                y="quantidade",
                text="quantidade",
                title=f"Quantidade por {col}",
                color="participacao_pct",
                color_continuous_scale=["#DBEAFE", "#0369A1"],
            )
            fig.add_scatter(
                x=dist[col],
                y=dist["acumulado_pct"],
                mode="lines+markers",
                name="% acumulado",
                yaxis="y2",
                line=dict(color="#F97316", width=3),
            )
            fig.update_layout(yaxis2=dict(title="% acumulado", overlaying="y", side="right", range=[0, 105]))
            fig = apply_analytics_chart_layout(fig, height=390)
            cols[i].plotly_chart(fig, use_container_width=True, key=f"{key_prefix}_count_rank_{i}_{normalize_token_text(col)}")
            with cols[i]:
                add_chart_reading(f"Quantidade por {col}", category_reading(dist, col))

    if groups.get("datetime"):
        date_col = profile.get("recommended_date") if profile.get("recommended_date") in groups["datetime"] else max(groups["datetime"], key=lambda c: df[c].notna().sum())
        time_df = df.copy()
        time_df[date_col] = pd.to_datetime(time_df[date_col], errors="coerce", dayfirst=True)
        time_df = time_df.dropna(subset=[date_col]).sort_values(date_col)
        if len(time_df) >= 8:
            span_days = int((time_df[date_col].max() - time_df[date_col].min()).days)
            freq, freq_label = select_time_frequency(span_days)
            base = time_df.set_index(date_col)
            if volume_col:
                trend = safe_resample_series(base[volume_col], freq, "nunique").reset_index(name="quantidade")
                title = f"IDs unicos por {freq_label} ({volume_col})"
            else:
                trend = safe_resample_series(pd.Series(1, index=base.index), freq, "count").reset_index(name="quantidade")
                title = f"Registros por {freq_label}"
            if date_col not in trend.columns:
                trend = trend.rename(columns={trend.columns[0]: date_col})
            trend = enrich_trend_dataframe(trend, date_col, "quantidade")
            fig = px.area(trend, x=date_col, y="quantidade", title=title)
            fig.add_scatter(x=trend[date_col], y=trend["media_movel"], mode="lines", name="media movel", line=dict(color="#F97316", width=3))
            fig = apply_analytics_chart_layout(fig, height=390)
            st.plotly_chart(fig, use_container_width=True, key=f"{key_prefix}_count_trend_{normalize_token_text(date_col)}")
            add_chart_reading(title, trend_reading(trend, "quantidade"))

    render_category_cross_analysis(df, groups, cat_cols[0] if cat_cols else profile.get("recommended_dimension"), key_prefix=key_prefix, top_n=10)

    if metric_numeric:
        metric_col = profile.get("recommended_metric") if profile.get("recommended_metric") in metric_numeric else metric_numeric[0]
        st.markdown("**Distribuicao da metrica recomendada**")
        left, right = st.columns(2)
        hist = px.histogram(df, x=metric_col, nbins=35, title=f"Distribuicao de {metric_col}", marginal="box")
        add_mean_median_lines(hist, df[metric_col], axis="x")
        hist = apply_analytics_chart_layout(hist, height=390)
        left.plotly_chart(hist, use_container_width=True, key=f"{key_prefix}_metric_distribution_{normalize_token_text(metric_col)}")
        with left:
            add_chart_reading(f"Distribuicao de {metric_col}", numeric_distribution_reading(df, metric_col))
        if cat_cols:
            bx = df[[cat_cols[0], metric_col]].dropna().copy()
            bx[cat_cols[0]] = bx[cat_cols[0]].astype(str)
            top_values = bx[cat_cols[0]].value_counts().head(10).index.tolist()
            bx = bx[bx[cat_cols[0]].isin(top_values)]
            fig_box = px.box(bx, x=cat_cols[0], y=metric_col, title=f"{metric_col} por {cat_cols[0]}", points="outliers")
            fig_box = apply_analytics_chart_layout(fig_box, height=390)
            right.plotly_chart(fig_box, use_container_width=True, key=f"{key_prefix}_metric_box_{normalize_token_text(metric_col)}_{normalize_token_text(cat_cols[0])}")


def build_metric_by_category(
    df: pd.DataFrame,
    cat_col: str,
    metric_col: Optional[str],
    agg_mode: str,
    top_n: int,
) -> Tuple[pd.DataFrame, str, str]:
    base = df[[cat_col] + ([metric_col] if metric_col else [])].dropna(subset=[cat_col]).copy()
    base[cat_col] = base[cat_col].astype(str)
    if metric_col and metric_col in base.columns and agg_mode != "Contagem":
        if agg_mode == "Media":
            out = base.groupby(cat_col, as_index=False)[metric_col].mean()
        elif agg_mode == "Mediana":
            out = base.groupby(cat_col, as_index=False)[metric_col].median()
        else:
            out = base.groupby(cat_col, as_index=False)[metric_col].sum()
        value_col = metric_col
        title = f"{agg_mode} de {metric_col} por {cat_col}"
    else:
        out = base.groupby(cat_col, as_index=False).size().rename(columns={"size": "contagem"})
        value_col = "contagem"
        title = f"Contagem de registros por {cat_col}"
    out = out.sort_values(value_col, ascending=False).head(top_n)
    return out, value_col, title


def metric_category_reading(out: pd.DataFrame, cat_col: str, value_col: str, title: str) -> List[str]:
    if out.empty:
        return []
    total = float(out[value_col].sum()) if value_col in out.columns else 0.0
    top = out.iloc[0]
    top_value = float(top[value_col])
    share = (top_value / total * 100.0) if total else 0.0
    bullets = [
        f"O maior grupo e `{top[cat_col]}` com {top_value:,.2f} em `{value_col}`.".replace(",", "."),
        f"Esse grupo representa {share:.1f}% do total exibido no ranking.",
    ]
    if len(out) >= 3:
        third_value = float(out.iloc[min(2, len(out) - 1)][value_col])
        if third_value > 0 and top_value / third_value >= 2:
            bullets.append("O primeiro colocado esta pelo menos 2x acima do terceiro, indicando diferenca relevante de desempenho/volume.")
    bullets.append(f"Use este ranking para priorizar acao sobre os maiores impactos em `{title}`.")
    return bullets


def add_mean_median_lines(fig, series: pd.Series, axis: str = "x") -> None:
    clean = pd.to_numeric(series, errors="coerce").dropna()
    if len(clean) == 0:
        return
    mean_v = float(clean.mean())
    median_v = float(clean.median())
    if axis == "x":
        fig.add_vline(x=mean_v, line_dash="dash", line_color="#F97316", annotation_text="media")
        fig.add_vline(x=median_v, line_dash="dot", line_color="#22C55E", annotation_text="mediana")
    else:
        fig.add_hline(y=mean_v, line_dash="dash", line_color="#F97316", annotation_text="media")
        fig.add_hline(y=median_v, line_dash="dot", line_color="#22C55E", annotation_text="mediana")


def numeric_distribution_reading(df: pd.DataFrame, col: str) -> List[str]:
    s = pd.to_numeric(df[col], errors="coerce").dropna()
    if len(s) == 0:
        return []
    q1, q3 = float(s.quantile(0.25)), float(s.quantile(0.75))
    mean_v, median_v = float(s.mean()), float(s.median())
    skew_hint = "acima" if mean_v > median_v else "abaixo"
    bullets = [
        f"Media {mean_v:,.2f} e mediana {median_v:,.2f}; media {skew_hint} da mediana sugere assimetria na distribuicao.".replace(",", "."),
        f"Faixa interquartil: {q1:,.2f} a {q3:,.2f}, onde fica a metade central dos registros.".replace(",", "."),
    ]
    outlier_ratio = float(((s < q1 - 1.5 * (q3 - q1)) | (s > q3 + 1.5 * (q3 - q1))).mean()) if q3 > q1 else 0.0
    if outlier_ratio >= 0.05:
        bullets.append(f"Outliers potenciais representam {outlier_ratio * 100:.1f}% dos valores; revisar antes de decisoes automaticas.")
    return bullets


def enrich_trend_dataframe(trend: pd.DataFrame, date_col: str, value_col: str = "valor") -> pd.DataFrame:
    out = trend.copy()
    out = out.sort_values(date_col)
    window = min(4, max(2, len(out) // 4))
    out["media_movel"] = out[value_col].rolling(window=window, min_periods=1).mean()
    out["variacao_pct"] = out[value_col].pct_change() * 100.0
    return out


def trend_reading(trend: pd.DataFrame, value_col: str = "valor") -> List[str]:
    if len(trend) < 2 or value_col not in trend.columns:
        return []
    first = float(trend[value_col].iloc[0])
    last = float(trend[value_col].iloc[-1])
    delta_abs = last - first
    delta_pct = (delta_abs / abs(first) * 100.0) if abs(first) > 1e-9 else None
    recent = trend[value_col].tail(min(4, len(trend))).astype(float)
    direction = "alta" if delta_abs > 0 else "queda" if delta_abs < 0 else "estabilidade"
    bullets = [f"A serie mostra {direction} entre o primeiro e o ultimo periodo."]
    if delta_pct is not None:
        bullets.append(f"Variacao acumulada aproximada: {delta_pct:.1f}%.")
    if len(recent) >= 3:
        recent_slope = recent.iloc[-1] - recent.iloc[0]
        if recent_slope > 0:
            bullets.append("Nos periodos recentes, a tendencia esta acelerando para cima.")
        elif recent_slope < 0:
            bullets.append("Nos periodos recentes, a tendencia esta desacelerando ou caindo.")
    return bullets


def scatter_reading(df: pd.DataFrame, x_col: str, y_col: str) -> List[str]:
    base = df[[x_col, y_col]].dropna()
    if len(base) < 10:
        return []
    corr = base[x_col].corr(base[y_col])
    if pd.isna(corr):
        return []
    strength = "forte" if abs(float(corr)) >= 0.7 else "moderada" if abs(float(corr)) >= 0.4 else "fraca"
    direction = "positiva" if corr >= 0 else "negativa"
    return [
        f"Correlacao {direction} {strength}: {float(corr):.3f}.",
        "Use como sinal de associacao; correlacao nao confirma causalidade sem validacao de negocio.",
    ]


def build_analytic_reading(
    df: pd.DataFrame,
    groups: dict[str, List[str]],
    quality_report: dict,
    issue_catalog: pd.DataFrame,
    top_categories_df: pd.DataFrame,
    numeric_distribution_df: pd.DataFrame,
    active_filters: Optional[List[str]] = None,
) -> List[str]:
    readings: List[str] = []
    if active_filters:
        readings.append(
            "A leitura analitica abaixo considera os filtros ativos; portanto, os numeros representam o recorte selecionado no dashboard."
        )

    if len(top_categories_df) > 0:
        top = top_categories_df.iloc[0]
        col = clean_report_text(top.get("coluna_categoria", "categoria"), 80)
        cat = clean_report_text(top.get("categoria", "-"), 80)
        share = float(top.get("participacao_pct", 0.0) or 0.0)
        if share >= 60:
            readings.append(
                f"A concentracao em `{cat}` na coluna `{col}` ({share:.1f}%) sugere dependencia operacional forte e deve ser monitorada como risco de gargalo."
            )
        elif share >= 30:
            readings.append(
                f"A categoria `{cat}` lidera a coluna `{col}` com {share:.1f}% dos registros, indicando prioridade natural para investigacao gerencial."
            )
        else:
            readings.append(
                f"A distribuicao de `{col}` esta mais pulverizada; a gestao deve comparar as principais categorias antes de escolher uma unica causa raiz."
            )

    missing_by_col = (df.isna().mean() * 100).sort_values(ascending=False) if len(df) > 0 else pd.Series(dtype=float)
    causal_cols = [
        col
        for col, pct in missing_by_col.items()
        if pct >= 80 and any(tok in normalize_token_text(col).split("_") for tok in ["motivo", "causa", "justificativa", "reason"])
    ]
    if causal_cols:
        readings.append(
            f"A ausencia elevada da coluna `{causal_cols[0]}` reduz a capacidade de analise causal e limita a explicacao do porque dos eventos."
        )
    elif len(missing_by_col) > 0 and float(missing_by_col.iloc[0]) >= 30:
        readings.append(
            f"A coluna `{missing_by_col.index[0]}` tem {float(missing_by_col.iloc[0]):.1f}% de ausencia, o que reduz confiabilidade em filtros, segmentacoes e modelos."
        )

    if float(quality_report.get("outlier_ratio", 0.0)) >= 0.05:
        readings.append(
            "A presenca relevante de outliers indica que parte dos registros pode representar erro de entrada ou eventos excepcionais; a decisao de excluir deve ser validada com o negocio."
        )

    if len(issue_catalog) > 0:
        top_risk = issue_catalog.sort_values(["priority_rank", "affected_rows"], ascending=[True, False]).iloc[0]
        readings.append(
            f"O risco mais urgente e `{top_risk['issue']}` em `{top_risk['column']}`, afetando {int(top_risk['affected_rows'])} registro(s)."
        )

    strongest = strongest_correlations(df, top_n=1)
    if len(strongest) > 0:
        row = strongest.iloc[0]
        readings.append(
            f"A correlacao entre `{row['variavel_1']}` e `{row['variavel_2']}` deve ser tratada como sinal analitico, nao como causalidade automatica."
        )

    if len(numeric_distribution_df) > 0:
        dist_col = str(numeric_distribution_df["coluna"].iloc[0])
        peak = numeric_distribution_df.sort_values("quantidade", ascending=False).iloc[0]
        readings.append(
            f"A distribuicao numerica de `{dist_col}` concentra mais registros na faixa `{peak['faixa']}`, util para calibrar metas, faixas de alerta e excecoes."
        )

    if not groups.get("datetime"):
        readings.append(
            "A ausencia de uma coluna temporal limita analises de tendencia, sazonalidade e comparacao periodo contra periodo."
        )

    unique_readings = []
    seen = set()
    for item in readings:
        key = clean_report_text(item, 220)
        if key not in seen:
            unique_readings.append(item)
            seen.add(key)
    return unique_readings[:8]


def build_dashboard_export_bundle(
    df: pd.DataFrame,
    groups: dict[str, List[str]],
    quality_report: dict,
    issue_catalog: pd.DataFrame,
    kpis: List[dict],
    active_filters: Optional[List[str]] = None,
    dashboard_config: Optional[dict] = None,
) -> dict:
    metric_numeric, id_like_cols = select_numeric_metric_columns(df, groups)
    volume_col = pick_volume_id_column(df, id_like_cols)
    spreadsheet_profile = detect_spreadsheet_profile(df, groups, quality_report)
    dashboard_config = dashboard_config or {}
    configured_date = dashboard_config.get("date_col")
    preferred_date = spreadsheet_profile.get("recommended_date")
    date_col = (
        configured_date
        if configured_date in groups["datetime"]
        else preferred_date
        if preferred_date in groups["datetime"]
        else (max(groups["datetime"], key=lambda c: df[c].notna().sum()) if groups["datetime"] else None)
    )
    configured_cat = dashboard_config.get("cat_col")
    configured_metric = dashboard_config.get("dist_col")
    preferred_cat = spreadsheet_profile.get("recommended_dimension")
    preferred_metric = spreadsheet_profile.get("recommended_metric")
    top_n = int(dashboard_config.get("top_n", 15) or 15)
    top_n = min(25, max(5, top_n))
    active_filters = active_filters or []

    summary_rows = [
        {"metrica": "linhas_analisadas", "valor": int(len(df))},
        {"metrica": "colunas_analisadas", "valor": int(df.shape[1])},
        {"metrica": "score_qualidade", "valor": float(quality_report.get("score", 0.0))},
        {"metrica": "nivel_qualidade", "valor": str(quality_report.get("level", "-"))},
        {"metrica": "riscos_criticos_altos", "valor": int(issue_catalog["priority"].astype(str).str.lower().isin(["critical", "high"]).sum()) if len(issue_catalog) > 0 and "priority" in issue_catalog.columns else 0},
    ]
    if volume_col:
        summary_rows.append({"metrica": f"ids_unicos_{volume_col}", "valor": int(df[volume_col].nunique(dropna=True))})
    summary_df = pd.DataFrame(summary_rows)

    narrative = build_dashboard_narrative(
        df,
        groups,
        quality_report,
        issue_catalog,
        kpis,
        active_filters=active_filters,
        volume_col=volume_col,
        spreadsheet_profile=spreadsheet_profile,
    )
    narrative_df = pd.DataFrame({"analise": narrative})
    filters_df = pd.DataFrame({"filtro": active_filters}) if active_filters else pd.DataFrame({"filtro": ["Nenhum filtro ativo no dashboard."]})
    profile_df = pd.DataFrame(
        [
            {"campo": "tipo_detectado", "valor": spreadsheet_profile.get("tipo", "-")},
            {"campo": "confianca_pct", "valor": spreadsheet_profile.get("confianca", 0.0)},
            {"campo": "foco", "valor": spreadsheet_profile.get("foco", "-")},
            {"campo": "melhor_analise", "valor": spreadsheet_profile.get("melhor_analise", "-")},
            {"campo": "dimensao_recomendada", "valor": spreadsheet_profile.get("recommended_dimension", "-")},
            {"campo": "metrica_recomendada", "valor": spreadsheet_profile.get("recommended_metric", "-")},
            {"campo": "data_recomendada", "valor": spreadsheet_profile.get("recommended_date", "-")},
            {"campo": "graficos_recomendados", "valor": "; ".join(spreadsheet_profile.get("graficos_recomendados", []))},
            {"campo": "analises_recomendadas", "valor": "; ".join(spreadsheet_profile.get("analises_recomendadas", []))},
            {"campo": "colunas_detectadas", "valor": ", ".join(spreadsheet_profile.get("colunas_detectadas", []))},
        ]
    )

    kpi_df = pd.DataFrame(
        [
            {
                "kpi": item["name"],
                "valor": format_kpi_metric(item),
                "variacao_recente": item["delta_text"] if item["delta_text"] else "-",
                "alerta": item["alert"] if item["alert"] else "-",
            }
            for item in kpis
        ]
    )

    top_categories_df = pd.DataFrame()
    metric_by_category_df = pd.DataFrame()
    if groups["categorical"]:
        cat_col = (
            configured_cat
            if configured_cat in groups["categorical"]
            else preferred_cat
            if preferred_cat in groups["categorical"]
            else max(groups["categorical"], key=lambda c: df[c].notna().sum())
        )
        dist = df[cat_col].fillna("NULO").astype(str).value_counts().head(top_n).reset_index()
        dist.columns = ["categoria", "quantidade"]
        dist["coluna_categoria"] = cat_col
        dist["participacao_pct"] = (dist["quantidade"] / max(len(df), 1) * 100.0).round(2)
        top_categories_df = dist[["coluna_categoria", "categoria", "quantidade", "participacao_pct"]]

        if metric_numeric:
            metric_col = (
                configured_metric
                if configured_metric in metric_numeric
                else preferred_metric
                if preferred_metric in metric_numeric
                else max(metric_numeric, key=lambda c: df[c].notna().sum())
            )
            if default_aggregation_for_profile(spreadsheet_profile, metric_col) == "Contagem":
                grp = (
                    df[[cat_col]]
                    .dropna(subset=[cat_col])
                    .assign(**{cat_col: lambda x: x[cat_col].astype(str)})
                    .groupby(cat_col, as_index=False)
                    .size()
                    .rename(columns={"size": "valor_metrica"})
                    .sort_values("valor_metrica", ascending=False)
                    .head(top_n)
                )
                grp["coluna_metrica"] = "contagem_registros"
                metric_by_category_df = grp[["coluna_metrica", cat_col, "valor_metrica"]].copy()
                metric_by_category_df.insert(0, "coluna_categoria", cat_col)
                metric_by_category_df.columns = ["coluna_categoria", "coluna_metrica", "categoria", "valor_metrica"]
                metric_by_category_df = metric_by_category_df[["coluna_categoria", "categoria", "coluna_metrica", "valor_metrica"]]
            else:
                grp = (
                    df[[cat_col, metric_col]]
                    .dropna(subset=[cat_col, metric_col])
                    .assign(**{cat_col: lambda x: x[cat_col].astype(str)})
                    .groupby(cat_col, as_index=False)[metric_col]
                    .sum()
                    .sort_values(metric_col, ascending=False)
                    .head(top_n)
                )
                grp["coluna_categoria"] = cat_col
                grp["coluna_metrica"] = metric_col
                metric_by_category_df = grp[["coluna_categoria", cat_col, "coluna_metrica", metric_col]]
                metric_by_category_df.columns = ["coluna_categoria", "categoria", "coluna_metrica", "valor_metrica"]
        elif volume_col:
            grp = (
                df[[cat_col, volume_col]]
                .dropna(subset=[cat_col])
                .assign(**{cat_col: lambda x: x[cat_col].astype(str)})
                .groupby(cat_col, as_index=False)[volume_col]
                .nunique()
                .sort_values(volume_col, ascending=False)
                .head(top_n)
            )
            grp["coluna_categoria"] = cat_col
            grp["coluna_metrica"] = volume_col
            metric_by_category_df = grp[["coluna_categoria", cat_col, "coluna_metrica", volume_col]]
            metric_by_category_df.columns = ["coluna_categoria", "categoria", "coluna_metrica", "valor_metrica"]

    trend_df = pd.DataFrame()
    trend_label = ""
    if date_col:
        time_df = df.copy()
        time_df[date_col] = pd.to_datetime(time_df[date_col], errors="coerce", dayfirst=True)
        time_df = time_df.dropna(subset=[date_col]).sort_values(date_col)
        if len(time_df) >= 8:
            span_days = int((time_df[date_col].max() - time_df[date_col].min()).days)
            freq, freq_label = select_time_frequency(span_days)
            base = time_df.set_index(date_col)

            if volume_col:
                trend_series = safe_resample_series(base[volume_col], freq, "nunique")
                trend = trend_series.reset_index(name="valor")
                trend_label = f"IDs unicos por {freq_label} ({volume_col})"
            elif metric_numeric:
                metric_col = max(metric_numeric, key=lambda c: time_df[c].notna().sum())
                if default_aggregation_for_profile(spreadsheet_profile, metric_col) == "Contagem":
                    trend_series = safe_resample_series(pd.Series(1, index=base.index), freq, "count")
                    trend = trend_series.reset_index(name="valor")
                    trend_label = f"Volume de registros por {freq_label}"
                else:
                    trend_series = safe_resample_series(base[metric_col], freq, "sum")
                    trend = trend_series.reset_index(name="valor")
                    trend_label = f"Soma de {metric_col} por {freq_label}"
            else:
                trend_series = safe_resample_series(pd.Series(1, index=base.index), freq, "count")
                trend = trend_series.reset_index(name="valor")
                trend_label = f"Volume de linhas por {freq_label}"

            if date_col not in trend.columns:
                trend = trend.rename(columns={trend.columns[0]: date_col})
            trend["indicador"] = trend_label
            trend_df = trend[[date_col, "indicador", "valor"]]

    distribution_preference = configured_metric if configured_metric in metric_numeric else preferred_metric
    distribution_metrics = ([distribution_preference] if distribution_preference in metric_numeric else []) + [c for c in metric_numeric if c != distribution_preference]
    numeric_distribution_df = build_numeric_distribution_table(df, distribution_metrics, preferred_metric=distribution_preference)
    analytic_reading = build_analytic_reading(
        df,
        groups,
        quality_report,
        issue_catalog,
        top_categories_df,
        numeric_distribution_df,
        active_filters=active_filters,
    )
    analytic_reading_df = pd.DataFrame({"leitura_analitica": analytic_reading})

    return {
        "summary_df": summary_df,
        "narrative_df": narrative_df,
        "profile_df": profile_df,
        "analytic_reading_df": analytic_reading_df,
        "filters_df": filters_df,
        "kpi_df": kpi_df,
        "top_categories_df": top_categories_df,
        "metric_by_category_df": metric_by_category_df,
        "numeric_distribution_df": numeric_distribution_df,
        "trend_df": trend_df,
        "trend_label": trend_label,
        "volume_col": volume_col,
    }


def render_auto_charts(df: pd.DataFrame, groups: dict[str, List[str]], top_corr: pd.DataFrame, key_prefix: str = "auto_charts") -> None:
    st.markdown("**Graficos automaticos por perfil de dados**")
    metric_numeric, id_like_cols = select_numeric_metric_columns(df, groups)
    volume_col = pick_volume_id_column(df, id_like_cols)

    if metric_numeric:
        vars_ = df[metric_numeric].var(numeric_only=True).sort_values(ascending=False)
        top_nums = vars_.index[: min(3, len(vars_))].tolist()
        if top_nums:
            c = st.columns(len(top_nums))
            for i, col in enumerate(top_nums):
                fig = px.histogram(df, x=col, nbins=35, title=f"Distribuicao - {col}", marginal="box")
                add_mean_median_lines(fig, df[col], axis="x")
                fig = apply_analytics_chart_layout(fig, height=390)
                c[i].plotly_chart(fig, use_container_width=True, key=f"{key_prefix}_hist_{i}_{col}")
                with c[i]:
                    add_chart_reading(f"Distribuicao - {col}", numeric_distribution_reading(df, col))
            fig_box = px.box(df[top_nums], y=top_nums, title="Comparativo de dispersao")
            fig_box = apply_analytics_chart_layout(fig_box)
            st.plotly_chart(fig_box, use_container_width=True, key=f"{key_prefix}_box_compare")

    if groups["categorical"]:
        chosen = []
        for col in groups["categorical"]:
            n = df[col].nunique(dropna=True)
            if 2 <= n <= 30:
                chosen.append((col, n))
        chosen = [x[0] for x in sorted(chosen, key=lambda t: t[1], reverse=True)[:3]]
        if chosen:
            c = st.columns(len(chosen))
            for i, col in enumerate(chosen):
                counts = categorical_distribution(df, col, top_n=12)
                fig = px.bar(
                    counts,
                    x=col,
                    y="quantidade",
                    title=f"Top categorias - {col}",
                    text="quantidade",
                    color="participacao_pct",
                    color_continuous_scale=["#BAE6FD", "#0369A1"],
                )
                fig.add_scatter(
                    x=counts[col],
                    y=counts["acumulado_pct"],
                    mode="lines+markers",
                    name="% acumulado",
                    yaxis="y2",
                    line=dict(color="#F97316", width=3),
                )
                fig.update_layout(yaxis2=dict(title="% acumulado", overlaying="y", side="right", range=[0, 105]))
                fig = apply_analytics_chart_layout(fig, height=420)
                c[i].plotly_chart(fig, use_container_width=True, key=f"{key_prefix}_cat_{i}_{col}")
                with c[i]:
                    add_chart_reading(f"Top categorias - {col}", category_reading(counts, col))

    if groups["datetime"] and metric_numeric:
        date_col = max(groups["datetime"], key=lambda x: df[x].notna().sum())
        metric_col = max(metric_numeric, key=lambda x: df[x].var(skipna=True) if pd.notna(df[x].var(skipna=True)) else -1)
        comp = period_comparison(df, date_col, metric_col)
        if comp:
            series = comp["series"].copy()
            if metric_col not in series.columns and "valor" in series.columns:
                series = series.rename(columns={"valor": metric_col})
            trend = enrich_trend_dataframe(series, date_col, metric_col)
            fig = px.line(trend, x=date_col, y=metric_col, markers=True, title=f"Serie temporal de {metric_col}")
            fig.add_scatter(x=trend[date_col], y=trend["media_movel"], mode="lines", name="media movel", line=dict(color="#F97316", width=3))
            fig = apply_analytics_chart_layout(fig)
            st.plotly_chart(fig, use_container_width=True, key=f"{key_prefix}_time_metric_{date_col}_{metric_col}")
            add_chart_reading(f"Serie temporal de {metric_col}", trend_reading(trend, metric_col))
    elif groups["datetime"] and volume_col:
        date_col = max(groups["datetime"], key=lambda x: df[x].notna().sum())
        comp = period_comparison(df, date_col, volume_col, agg_mode="nunique")
        if comp:
            series = comp["series"].copy()
            value_name = "quantidade_registros"
            series.columns = [date_col, value_name]
            trend = enrich_trend_dataframe(series, date_col, value_name)
            fig = px.line(trend, x=date_col, y=value_name, markers=True, title=f"Serie temporal de quantidade ({volume_col})")
            fig.add_scatter(x=trend[date_col], y=trend["media_movel"], mode="lines", name="media movel", line=dict(color="#F97316", width=3))
            fig = apply_analytics_chart_layout(fig)
            st.plotly_chart(fig, use_container_width=True, key=f"{key_prefix}_time_volume_{date_col}_{volume_col}")
            add_chart_reading(f"Serie temporal de quantidade ({volume_col})", trend_reading(trend, value_name))

    if not top_corr.empty:
        c1 = top_corr.iloc[0]["variavel_1"]
        c2 = top_corr.iloc[0]["variavel_2"]
        base = df[[c1, c2]].dropna().head(5000)
        if len(base) >= 10:
            fig = px.scatter(base, x=c1, y=c2, title=f"Relacao entre {c1} e {c2}", opacity=0.7)
            if base[c1].nunique(dropna=True) > 1:
                coef = np.polyfit(base[c1].astype(float), base[c2].astype(float), 1)
                x_line = np.linspace(float(base[c1].min()), float(base[c1].max()), 50)
                y_line = coef[0] * x_line + coef[1]
                fig.add_scatter(x=x_line, y=y_line, mode="lines", name="tendencia", line=dict(color="#F97316", width=3))
            fig = apply_analytics_chart_layout(fig)
            st.plotly_chart(fig, use_container_width=True, key=f"{key_prefix}_scatter_corr_{c1}_{c2}")
            add_chart_reading(f"Relacao entre {c1} e {c2}", scatter_reading(base, c1, c2))


def render_management_dashboard(
    df: pd.DataFrame,
    groups: dict[str, List[str]],
    quality_report: dict,
    issue_catalog: pd.DataFrame,
    kpis: List[dict],
    treatment_report: dict,
    violations: pd.DataFrame,
    *,
    key_prefix: str = "dash",
    turbo_config: Optional[dict] = None,
) -> None:
    st.subheader("Dashboard Gerencial")

    dashboard_df, dashboard_filters = apply_dashboard_filters(df, groups, key_prefix=key_prefix)
    st.session_state[f"{key_prefix}_active_filters"] = dashboard_filters
    st.session_state[f"{key_prefix}_filtered_df"] = dashboard_df.copy()
    if dashboard_df.empty:
        st.warning("Os filtros do dashboard removeram todos os registros. Ajuste os filtros desta aba.")
        return
    if dashboard_filters:
        st.caption("Filtros ativos no dashboard: " + " | ".join(dashboard_filters))

    dash_groups = get_column_groups(dashboard_df)
    metric_numeric, id_like_cols = select_numeric_metric_columns(dashboard_df, dash_groups)
    volume_col = pick_volume_id_column(dashboard_df, id_like_cols)
    dash_kpis = detect_kpis(dashboard_df, dash_groups) or kpis
    spreadsheet_profile = detect_spreadsheet_profile(dashboard_df, dash_groups, quality_report)
    st.session_state[f"{key_prefix}_profile"] = spreadsheet_profile

    critical_or_high = (
        issue_catalog["priority"].astype(str).str.lower().isin(["critical", "high"]).sum()
        if len(issue_catalog) > 0 and "priority" in issue_catalog.columns
        else 0
    )
    violation_ratio = (len(violations) / max(len(dashboard_df), 1)) * 100 if len(dashboard_df) > 0 else 0.0

    c1, c2, c3, c4, c5, c6 = st.columns(6)
    c1.metric("Volume analisado", f"{len(dashboard_df):,}".replace(",", "."))
    c2.metric("Score qualidade", f"{quality_report['score']:.1f}")
    c3.metric("Semaforo", quality_status(quality_report["score"]))
    c4.metric("Riscos crit./altos", int(critical_or_high))
    c5.metric("Violacoes/linha", f"{violation_ratio:.1f}%")
    c6.metric("Linhas saneadas", int(treatment_report.get("total_removed", 0)))

    render_spreadsheet_profile(spreadsheet_profile, key_prefix)
    render_treatment_summary(treatment_report, key_prefix=f"{key_prefix}_treatment")

    st.markdown("**Analise automatica da planilha**")
    dashboard_notes = build_dashboard_narrative(
        dashboard_df,
        dash_groups,
        quality_report,
        issue_catalog,
        dash_kpis,
        dashboard_filters,
        volume_col,
        spreadsheet_profile=spreadsheet_profile,
    )
    for note in dashboard_notes:
        st.markdown(f"- {note}")

    render_turbo_ai_box(
        turbo_config=turbo_config,
        section="Dashboard",
        button_label="Gerar leitura Turbo do Dashboard",
        context=build_ai_dataset_context(
            dashboard_df,
            dash_groups,
            quality_report=quality_report,
            profile=spreadsheet_profile,
            kpis=dash_kpis,
            issue_catalog=issue_catalog,
            active_filters=dashboard_filters,
        ),
        instruction=(
            "Explique para a gerencia o que esta acontecendo no dashboard. "
            "Destaque riscos, oportunidades, gargalos, leituras dos graficos recomendados e proximas acoes. "
            "Quando houver IDs/codigos, reforce que devem ser contados, nao somados."
        ),
        cache_key=f"{key_prefix}_dashboard_{hash_text('|'.join(dashboard_filters))[:8]}",
    )

    if dash_kpis:
        st.markdown("**KPI board executivo**")
        kpi_df = pd.DataFrame(
            [
                {
                    "KPI": item["name"],
                    "Valor": format_kpi_metric(item),
                    "Variacao recente": item["delta_text"] if item["delta_text"] else "-",
                    "Alerta": item["alert"] if item["alert"] else "-",
                }
                for item in dash_kpis
            ]
        )
        st.dataframe(kpi_df, use_container_width=True)

    st.markdown("**Explorador visual interativo**")
    controls = st.columns([1.2, 1, 1.2, 1.2, 1])
    preferred_cat = spreadsheet_profile.get("recommended_dimension")
    preferred_date = spreadsheet_profile.get("recommended_date")
    preferred_metric = spreadsheet_profile.get("recommended_metric")
    cat_index = dash_groups["categorical"].index(preferred_cat) if preferred_cat in dash_groups["categorical"] else 0
    date_index = dash_groups["datetime"].index(preferred_date) if preferred_date in dash_groups["datetime"] else 0
    metric_index = metric_numeric.index(preferred_metric) if preferred_metric in metric_numeric else 0
    cat_col = (
        controls[0].selectbox(
            "Dimensao categorica",
            options=dash_groups["categorical"],
            index=cat_index,
            key=f"{key_prefix}_cat_dim",
        )
        if dash_groups["categorical"]
        else None
    )
    top_n = controls[1].slider("Top categorias", min_value=5, max_value=25, value=10, step=1, key=f"{key_prefix}_topn")
    date_col = (
        controls[2].selectbox(
            "Coluna temporal",
            options=dash_groups["datetime"],
            index=date_index,
            key=f"{key_prefix}_date_dim",
        )
        if dash_groups["datetime"]
        else None
    )
    dist_col = (
        controls[3].selectbox(
            "Metrica numerica principal",
            options=metric_numeric,
            index=metric_index,
            key=f"{key_prefix}_primary_metric",
        )
        if metric_numeric
        else None
    )
    agg_options = ["Contagem", "Soma", "Media", "Mediana"]
    default_agg = default_aggregation_for_profile(spreadsheet_profile, dist_col)
    agg_key = f"{key_prefix}_primary_metric_agg"
    if agg_key not in st.session_state or st.session_state.get(agg_key) not in agg_options:
        st.session_state[agg_key] = default_agg
    agg_mode = (
        controls[4].selectbox(
            "Agregacao",
            options=agg_options,
            key=agg_key,
        )
        if metric_numeric
        else "Contagem"
    )
    if agg_mode == "Soma":
        st.caption("Atencao: soma deve ser usada apenas para metricas reais de negocio, como valor, receita, custo ou quantidade. IDs/codigos devem ser analisados por contagem.")
    st.session_state[f"{key_prefix}_config"] = {
        "cat_col": cat_col,
        "top_n": top_n,
        "date_col": date_col,
        "dist_col": dist_col,
        "agg_mode": agg_mode,
    }

    if cat_col:
        dist = categorical_distribution(dashboard_df, cat_col, top_n=top_n)
        left, right = st.columns(2)
        st.caption("A pizza mostra participacao por quantidade. A barra mostra ranking/Pareto com percentual acumulado.")
        pie_fig = px.pie(dist, names=cat_col, values="quantidade", title=f"Participacao por {cat_col}", hole=0.42)
        pie_fig = apply_analytics_chart_layout(pie_fig, height=420)
        left.plotly_chart(
            pie_fig,
            use_container_width=True,
            key=f"{key_prefix}_pie_{cat_col}",
        )
        with left:
            add_chart_reading(f"Participacao por {cat_col}", category_reading(dist, cat_col))

        if metric_numeric and dist_col:
            metric_col = dist_col
            out, value_col, title = build_metric_by_category(dashboard_df, cat_col, metric_col, agg_mode, top_n)
            out["participacao_pct"] = out[value_col] / max(float(out[value_col].sum()), 1e-9) * 100.0
            out["acumulado_pct"] = out["participacao_pct"].cumsum()
            bar_fig = px.bar(
                out,
                x=cat_col,
                y=value_col,
                text=value_col,
                title=title,
                color="participacao_pct",
                color_continuous_scale=["#DCFCE7", "#15803D"],
            )
            bar_fig.add_scatter(
                x=out[cat_col],
                y=out["acumulado_pct"],
                mode="lines+markers",
                name="% acumulado",
                yaxis="y2",
                line=dict(color="#F97316", width=3),
            )
            bar_fig.update_layout(yaxis2=dict(title="% acumulado", overlaying="y", side="right", range=[0, 105]))
            bar_fig = apply_analytics_chart_layout(bar_fig, height=420)
            right.plotly_chart(
                bar_fig,
                use_container_width=True,
                key=f"{key_prefix}_metric_top_{cat_col}_{value_col}_{agg_mode}",
            )
            with right:
                add_chart_reading(title, metric_category_reading(out, cat_col, value_col, title))
        else:
            bar_fig = px.bar(
                dist,
                x=cat_col,
                y="quantidade",
                text="quantidade",
                title=f"Top {top_n} categorias - {cat_col}",
                color="participacao_pct",
                color_continuous_scale=["#DBEAFE", "#1D4ED8"],
            )
            bar_fig.add_scatter(
                x=dist[cat_col],
                y=dist["acumulado_pct"],
                mode="lines+markers",
                name="% acumulado",
                yaxis="y2",
                line=dict(color="#F97316", width=3),
            )
            bar_fig.update_layout(yaxis2=dict(title="% acumulado", overlaying="y", side="right", range=[0, 105]))
            bar_fig = apply_analytics_chart_layout(bar_fig, height=420)
            right.plotly_chart(
                bar_fig,
                use_container_width=True,
                key=f"{key_prefix}_bar_top_{cat_col}",
            )
            with right:
                add_chart_reading(f"Top {top_n} categorias - {cat_col}", category_reading(dist, cat_col))

        treemap_fig = px.treemap(
            dist,
            path=[cat_col],
            values="quantidade",
            color="participacao_pct",
            color_continuous_scale=["#E0F2FE", "#075985"],
            title=f"Mapa de participacao por {cat_col}",
        )
        treemap_fig = apply_analytics_chart_layout(treemap_fig, height=380)
        st.plotly_chart(treemap_fig, use_container_width=True, key=f"{key_prefix}_treemap_{normalize_token_text(cat_col)}")
        render_category_cross_analysis(dashboard_df, dash_groups, cat_col, key_prefix=key_prefix, top_n=top_n)

    if date_col:
        trend_choices = ["Linhas por periodo"]
        if volume_col:
            trend_choices.append(f"IDs unicos ({volume_col})")
        trend_choices.extend([f"Metrica numerica ({c})" for c in metric_numeric])
        t1, t2 = st.columns(2)
        trend_choice_key = f"{key_prefix}_trend_choice"
        if trend_choice_key not in st.session_state or st.session_state.get(trend_choice_key) not in trend_choices:
            st.session_state[trend_choice_key] = f"IDs unicos ({volume_col})" if volume_col else "Linhas por periodo"
        trend_choice = t1.selectbox("Indicador da serie temporal", options=trend_choices, key=trend_choice_key)
        trend_agg_options = ["Contagem", "Soma", "Media", "Mediana"]
        trend_agg_key = f"{key_prefix}_trend_agg"
        if trend_agg_key not in st.session_state or st.session_state.get(trend_agg_key) not in trend_agg_options:
            st.session_state[trend_agg_key] = default_aggregation_for_profile(spreadsheet_profile, dist_col)
        metric_agg = t2.selectbox("Agregacao numerica temporal", options=trend_agg_options, key=trend_agg_key)
        st.session_state[f"{key_prefix}_config"].update({"trend_choice": trend_choice, "trend_agg": metric_agg})

        dt_series = pd.to_datetime(dashboard_df[date_col], errors="coerce", dayfirst=True)
        time_df = dashboard_df.copy()
        time_df[date_col] = dt_series
        time_df = time_df.dropna(subset=[date_col]).sort_values(date_col)
        if len(time_df) >= 8:
            span_days = int((time_df[date_col].max() - time_df[date_col].min()).days)
            freq, freq_label = select_time_frequency(span_days)
            base = time_df.set_index(date_col)

            if trend_choice == "Linhas por periodo":
                trend = safe_resample_series(pd.Series(1, index=base.index), freq, "count").reset_index(name="valor")
                y_name = "valor"
            elif trend_choice.startswith("IDs unicos"):
                trend = safe_resample_series(base[volume_col], freq, "nunique").reset_index(name="valor")
                y_name = "valor"
            else:
                metric_col = trend_choice.replace("Metrica numerica (", "").replace(")", "")
                if metric_agg == "Media":
                    trend = safe_resample_series(base[metric_col], freq, "mean").reset_index(name="valor")
                elif metric_agg == "Mediana":
                    trend = safe_resample_series(base[metric_col], freq, "median").reset_index(name="valor")
                elif metric_agg == "Contagem":
                    trend = safe_resample_series(base[metric_col], freq, "count").reset_index(name="valor")
                else:
                    trend = safe_resample_series(base[metric_col], freq, "sum").reset_index(name="valor")
                y_name = "valor"

            if date_col not in trend.columns:
                trend = trend.rename(columns={trend.columns[0]: date_col})
            trend = enrich_trend_dataframe(trend, date_col, y_name)
            chart_left, chart_right = st.columns(2)
            line_fig = px.line(trend, x=date_col, y=y_name, markers=True, title=f"Tendencia por {freq_label}")
            line_fig.add_scatter(x=trend[date_col], y=trend["media_movel"], mode="lines", name="media movel", line=dict(color="#F97316", width=3))
            line_fig = apply_analytics_chart_layout(line_fig)
            chart_left.plotly_chart(
                line_fig,
                use_container_width=True,
                key=f"{key_prefix}_trend_line_{date_col}_{trend_choice}_{metric_agg}",
            )
            with chart_left:
                add_chart_reading(f"Tendencia por {freq_label}", trend_reading(trend, y_name))
            bar_trend_fig = px.bar(trend, x=date_col, y=y_name, title=f"Volume por {freq_label}", color="variacao_pct", color_continuous_scale=["#DC2626", "#F8FAFC", "#16A34A"])
            bar_trend_fig = apply_analytics_chart_layout(bar_trend_fig)
            chart_right.plotly_chart(
                bar_trend_fig,
                use_container_width=True,
                key=f"{key_prefix}_trend_bar_{date_col}_{trend_choice}_{metric_agg}",
            )

    if dist_col:
        dist_left, dist_right = st.columns(2)
        hist_fig = px.histogram(dashboard_df, x=dist_col, nbins=35, title=f"Distribuicao de {dist_col}", marginal="box")
        add_mean_median_lines(hist_fig, dashboard_df[dist_col], axis="x")
        hist_fig = apply_analytics_chart_layout(hist_fig)
        dist_left.plotly_chart(
            hist_fig,
            use_container_width=True,
            key=f"{key_prefix}_dist_hist_{dist_col}",
        )
        with dist_left:
            add_chart_reading(f"Distribuicao de {dist_col}", numeric_distribution_reading(dashboard_df, dist_col))
        if cat_col:
            bx = dashboard_df[[cat_col, dist_col]].dropna(subset=[cat_col, dist_col]).copy()
            bx[cat_col] = bx[cat_col].astype(str)
            top_cat_values = bx[cat_col].value_counts().head(top_n).index.tolist()
            bx = bx[bx[cat_col].isin(top_cat_values)]
            box_fig = px.box(bx, x=cat_col, y=dist_col, title=f"Dispersao de {dist_col} por {cat_col}", points="outliers")
            box_fig = apply_analytics_chart_layout(box_fig)
            dist_right.plotly_chart(
                box_fig,
                use_container_width=True,
                key=f"{key_prefix}_dist_box_{cat_col}_{dist_col}",
            )
            with dist_right:
                add_chart_reading(
                    f"Dispersao de {dist_col} por {cat_col}",
                    [
                        "O boxplot compara mediana, dispersao e outliers entre grupos.",
                        "Grupos com caixas mais altas ou muitos pontos extremos merecem investigacao operacional.",
                    ],
                )

    if len(metric_numeric) >= 2:
        scatter_left, scatter_right = st.columns(2)
        x_col = scatter_left.selectbox("Eixo X (numerico)", options=metric_numeric, key=f"{key_prefix}_scatter_x")
        y_choices = [c for c in metric_numeric if c != x_col]
        y_col = scatter_right.selectbox(
            "Eixo Y (numerico)",
            options=y_choices if y_choices else metric_numeric,
            key=f"{key_prefix}_scatter_y",
        )
        base = dashboard_df[[x_col, y_col]].dropna().head(4000)
        if len(base) > 0:
            scatter_fig = px.scatter(base, x=x_col, y=y_col, opacity=0.7, title=f"Relacao entre {x_col} e {y_col}")
            if base[x_col].nunique(dropna=True) > 1:
                coef = np.polyfit(base[x_col].astype(float), base[y_col].astype(float), 1)
                x_line = np.linspace(float(base[x_col].min()), float(base[x_col].max()), 50)
                y_line = coef[0] * x_line + coef[1]
                scatter_fig.add_scatter(x=x_line, y=y_line, mode="lines", name="tendencia", line=dict(color="#F97316", width=3))
            scatter_fig = apply_analytics_chart_layout(scatter_fig)
            st.plotly_chart(
                scatter_fig,
                use_container_width=True,
                key=f"{key_prefix}_scatter_{x_col}_{y_col}",
            )
            add_chart_reading(f"Relacao entre {x_col} e {y_col}", scatter_reading(base, x_col, y_col))

    if len(metric_numeric) >= 2:
        top_numeric = metric_numeric[: min(12, len(metric_numeric))]
        corr = dashboard_df[top_numeric].corr(numeric_only=True)
        corr_fig = px.imshow(corr, text_auto=True, title="Heatmap de correlacao (top numericas)", color_continuous_scale="RdBu_r", zmin=-1, zmax=1)
        corr_fig = apply_analytics_chart_layout(corr_fig, height=520)
        st.plotly_chart(
            corr_fig,
            use_container_width=True,
            key=f"{key_prefix}_heatmap_corr",
        )
        corr_pairs = strongest_correlations(dashboard_df[top_numeric], top_n=3)
        if len(corr_pairs) > 0:
            add_chart_reading(
                "Heatmap de correlacao",
                [
                    f"{row['variavel_1']} x {row['variavel_2']}: correlacao {float(row['correlacao']):.3f}."
                    for _, row in corr_pairs.iterrows()
                ],
            )

    if len(issue_catalog) > 0:
        st.markdown("**Top riscos para a gerencia**")
        st.dataframe(
            issue_catalog.sort_values(["priority_rank", "affected_rows"], ascending=[True, False]).head(15),
            use_container_width=True,
        )


def suggest_best_k(x_scaled: np.ndarray) -> Optional[Tuple[int, float]]:
    max_k = min(8, len(x_scaled) - 1)
    if max_k < 2:
        return None
    best_k = 2
    best_score = -1.0
    for k in range(2, max_k + 1):
        model = KMeans(n_clusters=k, random_state=SEED, n_init="auto")
        labels = model.fit_predict(x_scaled)
        if len(np.unique(labels)) < 2:
            continue
        score = silhouette_score(x_scaled, labels)
        if score > best_score:
            best_k, best_score = k, score
    return (best_k, float(best_score)) if best_score >= 0 else None


def run_unsupervised_ml(
    df: pd.DataFrame,
    numeric_cols: List[str],
    contamination: float = 0.05,
    key_prefix: str = "ml_unsup",
) -> Tuple[Optional[dict], pd.DataFrame]:
    st.subheader("Machine Learning Sem Alvo")
    if len(numeric_cols) < 2:
        st.info("ML sem alvo requer pelo menos 2 colunas numericas.")
        return None, pd.DataFrame()
    ml_base = df[numeric_cols].copy().dropna(how="all")
    if len(ml_base) < 20:
        st.info("Use ao menos 20 linhas para clusterizacao/anomalias.")
        return None, pd.DataFrame()
    if len(ml_base) > MAX_ML_ROWS:
        ml_base = ml_base.sample(MAX_ML_ROWS, random_state=SEED)
        st.caption(f"ML sem alvo executado em amostra de {MAX_ML_ROWS} linhas.")

    x = StandardScaler().fit_transform(SimpleImputer(strategy="median").fit_transform(ml_base))
    k_res = suggest_best_k(x)
    if k_res is None:
        st.warning("Nao foi possivel estimar clusters estaveis.")
        return None, pd.DataFrame()
    best_k, best_score = k_res
    clusters = KMeans(n_clusters=best_k, random_state=SEED, n_init="auto").fit_predict(x)
    st.markdown(f"**Clusterizacao automatica:** melhor k = `{best_k}` (silhouette = `{best_score:.3f}`).")

    counts = pd.Series(clusters, name="cluster").value_counts().sort_index().reset_index()
    counts.columns = ["Cluster", "Quantidade"]
    st.plotly_chart(
        px.bar(counts, x="Cluster", y="Quantidade", text="Quantidade", title="Distribuicao por Cluster"),
        use_container_width=True,
        key=f"{key_prefix}_cluster_counts",
    )

    pca = PCA(n_components=2, random_state=SEED).fit_transform(x)
    proj_df = pd.DataFrame({"PCA 1": pca[:, 0], "PCA 2": pca[:, 1], "Cluster": clusters.astype(str)})
    st.plotly_chart(
        px.scatter(proj_df, x="PCA 1", y="PCA 2", color="Cluster", title="Mapa de Clusters (PCA 2D)"),
        use_container_width=True,
        key=f"{key_prefix}_cluster_map",
    )

    anom_model = IsolationForest(contamination=contamination, random_state=SEED)
    flags = anom_model.fit_predict(x)
    scores = anom_model.decision_function(x)
    out = ml_base.copy()
    out["cluster"] = clusters
    out["anomalia"] = np.where(flags == -1, "Sim", "Nao")
    out["score_anomalia"] = scores
    anomalias_df = out[out["anomalia"] == "Sim"].sort_values("score_anomalia")
    st.metric("Registros potencialmente anomalos", int(len(anomalias_df)))
    if len(anomalias_df):
        st.dataframe(anomalias_df.head(20), use_container_width=True)
    else:
        st.success("Nenhuma anomalia relevante detectada.")
    return {
        "best_k": int(best_k),
        "silhouette": float(best_score),
        "anomalias": int(len(anomalias_df)),
        "contamination": float(contamination),
    }, anomalias_df


def infer_problem_type(target: pd.Series) -> str:
    if pd.api.types.is_numeric_dtype(target) and target.nunique(dropna=True) > 10:
        return "regression"
    return "classification"


def detect_leakage(x: pd.DataFrame, y: pd.Series, target_col: str) -> List[dict]:
    flags = []
    y_non_null = y.dropna()
    t = target_col.lower()
    for col in x.columns:
        cl = col.lower()
        if cl == t or t in cl:
            flags.append({"column": col, "reason": "Nome da feature sugere leakage.", "severity": "high"})
        joined = pd.concat([x[col], y], axis=1).dropna()
        if len(joined) >= 30:
            same_ratio = float((joined.iloc[:, 0].astype(str) == joined.iloc[:, 1].astype(str)).mean())
            if same_ratio >= 0.98:
                flags.append({"column": col, "reason": "Feature quase identica ao alvo.", "severity": "critical"})
        if pd.api.types.is_numeric_dtype(x[col]) and pd.api.types.is_numeric_dtype(y_non_null):
            joined_num = pd.concat([x[col], y], axis=1).dropna()
            if len(joined_num) >= 30:
                corr = joined_num.iloc[:, 0].corr(joined_num.iloc[:, 1])
                if pd.notna(corr) and abs(float(corr)) >= 0.98:
                    flags.append({"column": col, "reason": f"Correlacao extrema com alvo ({corr:.3f}).", "severity": "high"})
    dedup = []
    seen = set()
    for item in flags:
        key = (item["column"], item["reason"])
        if key not in seen:
            dedup.append(item)
            seen.add(key)
    return dedup


def serialize_feature_importance(imp_df: pd.DataFrame, top_n: int = 12) -> List[dict]:
    if imp_df is None or len(imp_df) == 0:
        return []
    out = []
    for _, row in imp_df.head(top_n).iterrows():
        out.append({"feature": clean_report_text(row.get("feature", "-"), 90), "importance": float(row.get("importance", 0.0) or 0.0)})
    return out


def ml_reliability_alert(problem_type: str, metrics: dict) -> str:
    if problem_type == "classification":
        score = float(metrics.get("holdout_f1", metrics.get("cv_f1", 0.0)) or 0.0)
        if score >= 0.75:
            return "Confiabilidade boa para apoio gerencial, mantendo validacao humana antes de decisoes automaticas."
        if score >= 0.60:
            return "Confiabilidade moderada: use como sinal de apoio e melhore volume/qualidade dos dados antes de automatizar."
        return "Confiabilidade baixa: resultado exploratorio, recomendado revisar alvo, features e balanceamento das classes."
    score = float(metrics.get("holdout_r2", metrics.get("cv_r2", 0.0)) or 0.0)
    if score >= 0.60:
        return "Confiabilidade boa para estimativas direcionais, mantendo monitoramento de erro por periodo."
    if score >= 0.30:
        return "Confiabilidade moderada: modelo explica parte da variacao, mas ainda requer mais dados ou variaveis explicativas."
    return "Confiabilidade baixa: modelo tem baixo poder explicativo para previsao operacional neste conjunto de dados."


def suggest_ml_targets(df: pd.DataFrame, groups: dict[str, List[str]], max_items: int = 6) -> List[dict]:
    id_like_cols = set(detect_id_like_columns(df))
    suggestions = []
    preferred_tokens = [
        "status",
        "situacao",
        "resultado",
        "risco",
        "atraso",
        "sla",
        "cancelado",
        "inadimplente",
        "valor",
        "custo",
        "tempo",
        "prazo",
        "score",
        "nota",
    ]
    for col in df.columns:
        if col in id_like_cols:
            continue
        non_null = df[col].dropna()
        if len(non_null) < 30:
            continue
        missing_pct = 100.0 * (1.0 - len(non_null) / max(len(df), 1))
        nunique = int(non_null.nunique(dropna=True))
        unique_ratio = nunique / max(len(non_null), 1)
        key = normalize_token_text(col)
        score = 0.0
        reasons = []
        if any(tok in key for tok in preferred_tokens):
            score += 35
            reasons.append("nome indica alvo de negocio")
        if pd.api.types.is_numeric_dtype(df[col]) and nunique > 10:
            score += 20
            reasons.append("numerico continuo para regressao")
        elif 2 <= nunique <= 20:
            score += 25
            reasons.append("classes manejaveis para classificacao")
        if missing_pct <= 10:
            score += 15
            reasons.append("baixo percentual de nulos")
        if unique_ratio < 0.5:
            score += 10
        if len(non_null) >= 100:
            score += 10
        if score <= 0:
            continue
        suggestions.append(
            {
                "coluna": col,
                "tipo": infer_problem_type(df[col]),
                "score": round(score, 1),
                "nulos_pct": round(missing_pct, 2),
                "valores_unicos": nunique,
                "motivo": "; ".join(reasons[:3]) or "potencial estatistico razoavel",
            }
        )
    return sorted(suggestions, key=lambda item: item["score"], reverse=True)[:max_items]


def diagnose_ml_target(df: pd.DataFrame, target_col: str) -> Tuple[List[str], List[str]]:
    warnings: List[str] = []
    strengths: List[str] = []
    if target_col == "(nao usar)" or target_col not in df.columns:
        return warnings, strengths
    y = df[target_col]
    non_null = y.dropna()
    missing_pct = 100.0 * (1.0 - len(non_null) / max(len(df), 1))
    nunique = int(non_null.nunique(dropna=True))
    if len(non_null) < 100:
        warnings.append("Poucos registros com alvo preenchido; o modelo pode ficar instavel.")
    else:
        strengths.append("Volume minimo de registros preenchidos adequado para treino inicial.")
    if missing_pct > 20:
        warnings.append(f"Alvo com {missing_pct:.1f}% de nulos; isso reduz a base util de treinamento.")
    if infer_problem_type(y) == "classification":
        counts = non_null.astype(str).value_counts()
        if len(counts) < 2:
            warnings.append("Alvo possui menos de 2 classes; nao e possivel classificar.")
        elif len(counts) > 30:
            warnings.append("Alvo possui muitas classes; recomenda agrupar categorias antes de treinar.")
        else:
            strengths.append(f"Classificacao com {len(counts)} classes detectadas.")
        if len(counts) > 1:
            imbalance = float(counts.iloc[0] / max(counts.sum(), 1))
            if imbalance >= 0.85:
                warnings.append(f"Classe dominante representa {imbalance * 100:.1f}% dos casos; avaliar desbalanceamento.")
    else:
        numeric = pd.to_numeric(non_null, errors="coerce").dropna()
        if len(numeric) < 30:
            warnings.append("Alvo numerico tem poucos valores validos para regressao.")
        if len(numeric) > 0 and numeric.nunique(dropna=True) < 10:
            warnings.append("Alvo numerico tem baixa variacao; classificacao ou faixas podem ser melhores.")
        if len(numeric) >= 30:
            strengths.append("Alvo numerico possui variacao suficiente para uma regressao inicial.")
    id_like = detect_id_like_columns(df[[target_col]]) if target_col in df.columns else []
    if id_like:
        warnings.append("A coluna parece identificador/codigo; normalmente nao e um bom alvo preditivo.")
    return warnings, strengths


def ml_business_explanation(result: dict) -> List[str]:
    if not result:
        return []
    target = result.get("target_col", "alvo")
    model = result.get("model_choice", "modelo")
    kind = result.get("problem_type", "-")
    bullets = [f"O modelo `{model}` foi treinado para prever `{target}` usando as demais colunas validas da planilha."]
    if result.get("requested_model_choice") == "Auto (comparar modelos)":
        bullets.append("A opcao automatica comparou modelos por validacao cruzada e manteve o melhor desempenho para o treino final.")
    if kind == "classification":
        bullets.append(
            f"Como o alvo e categorico, a leitura principal e F1/acuracia. F1 holdout: {float(result.get('holdout_f1', 0.0)):.3f}."
        )
    else:
        bullets.append(
            f"Como o alvo e numerico, a leitura principal e R2/RMSE. R2 holdout: {float(result.get('holdout_r2', 0.0)):.3f}."
        )
    if result.get("model_reliability_alert"):
        bullets.append(result["model_reliability_alert"])
    top_features = result.get("feature_importance", [])[:3]
    if top_features:
        feats = ", ".join([item["feature"] for item in top_features])
        bullets.append(f"Variaveis com maior influencia estimada: {feats}.")
    if result.get("removed_leakage_features"):
        bullets.append("Algumas colunas foram removidas por risco de vazamento de informacao do alvo.")
    return bullets


def prepare_prediction_input(df: pd.DataFrame, training_features: List[str]) -> pd.DataFrame:
    prepared = df.copy()
    for col in training_features:
        if col not in prepared.columns:
            prepared[col] = np.nan
    return prepared[training_features]


def public_ml_summary(result: Optional[dict]) -> dict:
    if not result:
        return {}
    hidden_keys = {"pipeline"}
    return {k: v for k, v in result.items() if k not in hidden_keys}


def get_classification_candidates() -> List[Tuple[str, object]]:
    return [
        ("Random Forest", RandomForestClassifier(n_estimators=300, random_state=SEED, class_weight="balanced_subsample")),
        ("Regressao Logistica", LogisticRegression(max_iter=1000, class_weight="balanced", random_state=SEED)),
        ("Gradient Boosting", GradientBoostingClassifier(random_state=SEED)),
    ]


def get_regression_candidates() -> List[Tuple[str, object]]:
    return [
        ("Random Forest", RandomForestRegressor(n_estimators=300, random_state=SEED)),
        ("Regressao Linear", LinearRegression()),
        ("Gradient Boosting", GradientBoostingRegressor(random_state=SEED)),
    ]


def select_model_candidates(problem_type: str, model_choice: str) -> List[Tuple[str, object]]:
    if problem_type == "classification":
        candidates = get_classification_candidates()
        manual_alias = "Regressao Logistica"
    else:
        candidates = get_regression_candidates()
        manual_alias = "Regressao Linear"

    if model_choice == "Auto (comparar modelos)":
        return candidates
    if model_choice == "Random Forest":
        return [item for item in candidates if item[0] == "Random Forest"]
    if model_choice == "Regressao Linear / Logistica":
        return [item for item in candidates if item[0] == manual_alias]
    return candidates[:1]


def compare_model_candidates(
    *,
    x: pd.DataFrame,
    y: pd.Series,
    preprocess: ColumnTransformer,
    candidates: List[Tuple[str, object]],
    cv_folds: int,
    problem_type: str,
) -> Tuple[str, object, pd.DataFrame]:
    rows = []
    best_name = candidates[0][0]
    best_model = candidates[0][1]
    best_score = -np.inf

    for name, model in candidates:
        pipe = Pipeline(steps=[("preprocess", preprocess), ("model", model)])
        try:
            if problem_type == "classification":
                cv = cross_validate(pipe, x, y, cv=cv_folds, scoring={"acc": "accuracy", "f1": "f1_weighted"}, error_score="raise")
                score = float(np.mean(cv["test_f1"]))
                rows.append(
                    {
                        "modelo": name,
                        "metrica_principal": "F1 ponderado",
                        "score": score,
                        "acuracia_cv": float(np.mean(cv["test_acc"])),
                        "f1_cv": score,
                        "status": "ok",
                    }
                )
            else:
                cv = cross_validate(
                    pipe,
                    x,
                    y,
                    cv=cv_folds,
                    scoring={"r2": "r2", "rmse": "neg_root_mean_squared_error", "mae": "neg_mean_absolute_error"},
                    error_score="raise",
                )
                score = float(np.mean(cv["test_r2"]))
                rows.append(
                    {
                        "modelo": name,
                        "metrica_principal": "R2",
                        "score": score,
                        "r2_cv": score,
                        "rmse_cv": float(-np.mean(cv["test_rmse"])),
                        "mae_cv": float(-np.mean(cv["test_mae"])),
                        "status": "ok",
                    }
                )
            if score > best_score:
                best_score = score
                best_name = name
                best_model = model
        except Exception as exc:
            rows.append(
                {
                    "modelo": name,
                    "metrica_principal": "erro",
                    "score": None,
                    "status": f"falhou: {str(exc)[:120]}",
                }
            )

    comparison_df = pd.DataFrame(rows)
    if len(comparison_df) > 0 and "score" in comparison_df.columns:
        comparison_df = comparison_df.sort_values("score", ascending=False, na_position="last")
        comparison_df = comparison_df.astype(object).where(pd.notna(comparison_df), None)
    return best_name, best_model, comparison_df


def sanitize_ml_features(x: pd.DataFrame) -> Tuple[pd.DataFrame, List[str], List[str], List[dict]]:
    clean = x.copy().dropna(axis=1, how="all")
    warnings = []
    num = clean.select_dtypes(include=np.number).columns.tolist()
    cat = [c for c in clean.columns if c not in num]

    for col in cat:
        raw = clean[col]
        if pd.api.types.is_datetime64_any_dtype(raw):
            values = pd.to_datetime(raw, errors="coerce").dt.strftime("%Y-%m-%d")
        else:
            values = raw.astype("object")
        values = values.where(pd.notna(values), ML_MISSING_CATEGORY).astype(str)
        values = values.replace(
            {
                "nan": ML_MISSING_CATEGORY,
                "NaN": ML_MISSING_CATEGORY,
                "NaT": ML_MISSING_CATEGORY,
                "None": ML_MISSING_CATEGORY,
                "<NA>": ML_MISSING_CATEGORY,
            }
        )
        unique_count = int(pd.Series(values).nunique(dropna=False))
        if unique_count > MAX_ML_CATEGORICAL_LEVELS:
            top_values = pd.Series(values).value_counts(dropna=False).head(MAX_ML_CATEGORICAL_LEVELS).index
            values = pd.Series(np.where(pd.Series(values).isin(top_values), values, ML_OTHER_CATEGORY), index=clean.index)
            warnings.append(
                {
                    "coluna": col,
                    "acao": f"Categorias agrupadas em {ML_OTHER_CATEGORY}",
                    "categorias_originais": unique_count,
                    "limite": MAX_ML_CATEGORICAL_LEVELS,
                }
            )
        clean[col] = values

    return clean, num, cat, warnings


def run_supervised_ml(
    df: pd.DataFrame,
    target_col: str,
    model_choice: str = "Random Forest",
    key_prefix: str = "ml_sup",
) -> Optional[dict]:
    st.subheader("Modelo Preditivo")
    if df.shape[1] < 2:
        st.info("Necessario pelo menos duas colunas para modelo supervisionado.")
        return None
    if target_col == "(nao usar)":
        st.caption("Selecione uma coluna alvo para ativar o treinamento.")
        return None

    features = [c for c in df.columns if c != target_col]
    model_df = df[features + [target_col]].copy().dropna(subset=[target_col])
    if len(model_df) < 30:
        st.warning("Poucos registros com alvo preenchido (minimo recomendado: 30).")
        return None
    if len(model_df) > MAX_ML_ROWS:
        model_df = model_df.sample(MAX_ML_ROWS, random_state=SEED)
        st.caption(f"Treinamento supervisionado executado em amostra de {MAX_ML_ROWS} linhas.")

    x = model_df[features]
    y = model_df[target_col]
    leakage = detect_leakage(x, y, target_col)
    leak_cols = sorted(list({item["column"] for item in leakage}))
    if leakage:
        st.warning("Possivel leakage detectado. Colunas abaixo removidas automaticamente:")
        st.dataframe(pd.DataFrame(leakage), use_container_width=True)
        x = x.drop(columns=[c for c in leak_cols if c in x.columns])
    if x.shape[1] == 0:
        st.error("Todas as features foram removidas por risco de leakage.")
        return None

    x, num, cat, feature_warnings = sanitize_ml_features(x)
    training_features = x.columns.tolist()
    if feature_warnings:
        st.info("Algumas colunas categoricas tinham muitos valores unicos e foram agrupadas para estabilizar o modelo.")
        st.dataframe(pd.DataFrame(feature_warnings), use_container_width=True)
    if x.shape[1] == 0:
        st.error("Nao ha features validas para treinar o modelo apos a preparacao dos dados.")
        return None
    preprocess = ColumnTransformer(
        transformers=[
            ("num", Pipeline(steps=[("imputer", SimpleImputer(strategy="median")), ("scaler", StandardScaler())]), num),
            ("cat", Pipeline(steps=[("onehot", OneHotEncoder(handle_unknown="ignore"))]), cat),
        ],
        remainder="drop",
    )

    kind = infer_problem_type(y)
    st.caption(f"Alvo: `{target_col}` | Tipo detectado: `{kind}` | Modelo selecionado: `{model_choice}`")
    if kind == "classification":
        y = y.where(pd.notna(y), ML_MISSING_CATEGORY).astype(str)
        class_counts = y.value_counts(dropna=True)
        rare_classes = class_counts[class_counts < 2].index.tolist()
        if rare_classes:
            keep_mask = ~y.isin(rare_classes)
            x = x.loc[keep_mask]
            y = y.loc[keep_mask]
            st.info(
                f"{len(rare_classes)} classe(s) com apenas 1 registro foram removidas do treino para permitir validacao."
            )
        if y.nunique(dropna=True) < 2:
            st.warning("Alvo para classificacao precisa ter pelo menos 2 classes.")
            return None
        class_counts = y.value_counts(dropna=True)
        min_count = int(class_counts.min()) if len(class_counts) else 1
        cv_folds = int(min(5, min_count))
        if cv_folds < 2:
            st.warning("Nao ha registros suficientes por classe para validacao cruzada.")
            return None
        estimated_test_rows = max(1, int(np.ceil(len(y) * 0.2)))
        strat = y if min_count >= 2 and y.nunique(dropna=True) <= min(30, estimated_test_rows) else None
        candidates = select_model_candidates("classification", model_choice)
        selected_model_name, model, comparison_df = compare_model_candidates(
            x=x,
            y=y,
            preprocess=preprocess,
            candidates=candidates,
            cv_folds=cv_folds,
            problem_type="classification",
        )
        if comparison_df.empty or not (comparison_df.get("status") == "ok").any():
            st.error("Nenhum modelo de classificacao conseguiu treinar com esta configuracao.")
            if len(comparison_df) > 0:
                st.dataframe(comparison_df, use_container_width=True)
            return None
        if model_choice == "Auto (comparar modelos)":
            st.markdown("**Comparacao automatica de modelos**")
            st.dataframe(comparison_df, use_container_width=True)
            st.success(f"Melhor modelo selecionado: `{selected_model_name}`.")
        else:
            st.caption(f"Modelo treinado: `{selected_model_name}`.")
        pipe = Pipeline(steps=[("preprocess", preprocess), ("model", model)])
        try:
            cv = cross_validate(pipe, x, y, cv=cv_folds, scoring={"acc": "accuracy", "f1": "f1_weighted"}, error_score="raise")
            x_train, x_test, y_train, y_test = train_test_split(x, y, test_size=0.2, random_state=SEED, stratify=strat)
            pipe.fit(x_train, y_train)
            pred = pipe.predict(x_test)
        except Exception as exc:
            st.error("Nao foi possivel treinar o modelo de classificacao com esta configuracao.")
            st.caption(f"Detalhe tecnico: {str(exc)[:500]}")
            return None
        hold_acc = accuracy_score(y_test, pred)
        hold_f1 = f1_score(y_test, pred, average="weighted")
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("CV Acuracia", f"{np.mean(cv['test_acc']):.3f}")
        c2.metric("CV F1", f"{np.mean(cv['test_f1']):.3f}")
        c3.metric("Holdout Acuracia", f"{hold_acc:.3f}")
        c4.metric("Holdout F1", f"{hold_f1:.3f}")
        sample_x = x_test.sample(min(len(x_test), 1200), random_state=SEED) if len(x_test) > 1200 else x_test
        sample_y = y_test.loc[sample_x.index]
        imp = permutation_importance(pipe, sample_x, sample_y, scoring="f1_weighted", n_repeats=5, random_state=SEED)
        imp_df = pd.DataFrame({"feature": sample_x.columns, "importance": imp.importances_mean}).sort_values("importance", ascending=False)
        labels = sorted(pd.Series(y_test).astype(str).unique().tolist())
        cm = confusion_matrix(pd.Series(y_test).astype(str), pd.Series(pred).astype(str), labels=labels)
        cm_df = pd.DataFrame(cm, index=[f"real_{x}" for x in labels], columns=[f"previsto_{x}" for x in labels])
        st.markdown("**Matriz de confusao (holdout)**")
        st.dataframe(cm_df, use_container_width=True)
        st.plotly_chart(
            px.bar(imp_df.head(12), x="feature", y="importance", title="Top features (permutation importance)"),
            use_container_width=True,
            key=f"{key_prefix}_classification_importance",
        )
        st.success("Classificacao treinada com validacao cruzada e explicabilidade.")
        metrics = {
            "cv_accuracy": float(np.mean(cv["test_acc"])),
            "cv_f1": float(np.mean(cv["test_f1"])),
            "holdout_accuracy": float(hold_acc),
            "holdout_f1": float(hold_f1),
        }
        return {
            "problem_type": "classification",
            "target_col": target_col,
            "model_choice": selected_model_name,
            "requested_model_choice": model_choice,
            **metrics,
            "model_comparison": comparison_df.to_dict("records"),
            "removed_leakage_features": leak_cols,
            "feature_importance": serialize_feature_importance(imp_df),
            "model_reliability_alert": ml_reliability_alert("classification", metrics),
            "confusion_matrix": cm_df.reset_index().rename(columns={"index": "real"}).to_dict("records"),
            "training_features": training_features,
            "pipeline": pipe,
        }

    cv_folds = 5 if len(x) >= 100 else 3
    try:
        y = pd.to_numeric(y, errors="coerce")
        valid_target = y.notna()
        x = x.loc[valid_target]
        y = y.loc[valid_target]
        if len(x) < 30:
            st.warning("Poucos registros com alvo numerico valido para regressao (minimo recomendado: 30).")
            return None
        candidates = select_model_candidates("regression", model_choice)
        selected_model_name, model, comparison_df = compare_model_candidates(
            x=x,
            y=y,
            preprocess=preprocess,
            candidates=candidates,
            cv_folds=cv_folds,
            problem_type="regression",
        )
        if comparison_df.empty or not (comparison_df.get("status") == "ok").any():
            st.error("Nenhum modelo de regressao conseguiu treinar com esta configuracao.")
            if len(comparison_df) > 0:
                st.dataframe(comparison_df, use_container_width=True)
            return None
        if model_choice == "Auto (comparar modelos)":
            st.markdown("**Comparacao automatica de modelos**")
            st.dataframe(comparison_df, use_container_width=True)
            st.success(f"Melhor modelo selecionado: `{selected_model_name}`.")
        else:
            st.caption(f"Modelo treinado: `{selected_model_name}`.")
        pipe = Pipeline(steps=[("preprocess", preprocess), ("model", model)])
        cv = cross_validate(
            pipe,
            x,
            y,
            cv=cv_folds,
            scoring={"r2": "r2", "rmse": "neg_root_mean_squared_error", "mae": "neg_mean_absolute_error"},
            error_score="raise",
        )
        x_train, x_test, y_train, y_test = train_test_split(x, y, test_size=0.2, random_state=SEED)
        pipe.fit(x_train, y_train)
        pred = pipe.predict(x_test)
    except Exception as exc:
        st.error("Nao foi possivel treinar o modelo de regressao com esta configuracao.")
        st.caption(f"Detalhe tecnico: {str(exc)[:500]}")
        return None
    hold_mae = mean_absolute_error(y_test, pred)
    hold_rmse = mean_squared_error(y_test, pred) ** 0.5
    hold_r2 = r2_score(y_test, pred)
    c1, c2, c3, c4, c5, c6 = st.columns(6)
    c1.metric("CV R2", f"{np.mean(cv['test_r2']):.3f}")
    c2.metric("CV RMSE", f"{-np.mean(cv['test_rmse']):.3f}")
    c3.metric("CV MAE", f"{-np.mean(cv['test_mae']):.3f}")
    c4.metric("Holdout R2", f"{hold_r2:.3f}")
    c5.metric("Holdout RMSE", f"{hold_rmse:.3f}")
    c6.metric("Holdout MAE", f"{hold_mae:.3f}")
    sample_x = x_test.sample(min(len(x_test), 1200), random_state=SEED) if len(x_test) > 1200 else x_test
    sample_y = y_test.loc[sample_x.index]
    imp = permutation_importance(pipe, sample_x, sample_y, scoring="r2", n_repeats=5, random_state=SEED)
    imp_df = pd.DataFrame({"feature": sample_x.columns, "importance": imp.importances_mean}).sort_values("importance", ascending=False)
    pred_df = pd.DataFrame({"real": y_test, "previsto": pred})
    pred_df["erro"] = pred_df["previsto"] - pred_df["real"]
    fig_pred = px.scatter(pred_df, x="real", y="previsto", title="Real vs previsto (holdout)", opacity=0.7)
    min_line = float(min(pred_df["real"].min(), pred_df["previsto"].min()))
    max_line = float(max(pred_df["real"].max(), pred_df["previsto"].max()))
    fig_pred.add_scatter(x=[min_line, max_line], y=[min_line, max_line], mode="lines", name="ideal", line=dict(color="#F97316", dash="dash"))
    fig_pred = apply_analytics_chart_layout(fig_pred)
    st.plotly_chart(fig_pred, use_container_width=True, key=f"{key_prefix}_regression_real_vs_pred")
    st.plotly_chart(
        px.bar(imp_df.head(12), x="feature", y="importance", title="Top features (permutation importance)"),
        use_container_width=True,
        key=f"{key_prefix}_regression_importance",
    )
    st.success("Regressao treinada com validacao cruzada e explicabilidade.")
    metrics = {
        "cv_r2": float(np.mean(cv["test_r2"])),
        "cv_rmse": float(-np.mean(cv["test_rmse"])),
        "cv_mae": float(-np.mean(cv["test_mae"])),
        "holdout_r2": float(hold_r2),
        "holdout_rmse": float(hold_rmse),
        "holdout_mae": float(hold_mae),
    }
    return {
        "problem_type": "regression",
        "target_col": target_col,
        "model_choice": selected_model_name,
        "requested_model_choice": model_choice,
        **metrics,
        "model_comparison": comparison_df.to_dict("records"),
        "removed_leakage_features": leak_cols,
        "feature_importance": serialize_feature_importance(imp_df),
        "model_reliability_alert": ml_reliability_alert("regression", metrics),
        "prediction_sample": pred_df.head(200).to_dict("records"),
        "training_features": training_features,
        "pipeline": pipe,
    }


def compare_dataframes(current_df: pd.DataFrame, other_df: pd.DataFrame, current_name: str = "Atual", other_name: str = "Comparacao") -> dict:
    current_groups = get_column_groups(current_df)
    other_groups = get_column_groups(other_df)
    current_quality = compute_quality_report(current_df, current_groups, pd.DataFrame())
    other_quality = compute_quality_report(other_df, other_groups, pd.DataFrame())

    common_cols = [c for c in current_df.columns if c in other_df.columns]
    new_cols = [c for c in current_df.columns if c not in other_df.columns]
    removed_cols = [c for c in other_df.columns if c not in current_df.columns]
    current_metric_numeric, current_ids = select_numeric_metric_columns(current_df, current_groups)
    other_metric_numeric, other_ids = select_numeric_metric_columns(other_df, other_groups)
    id_like = set(current_ids + other_ids)
    common_numeric = [c for c in common_cols if c in current_metric_numeric and c in other_metric_numeric]
    common_categorical = [c for c in common_cols if c in current_groups["categorical"] and c in other_groups["categorical"]]

    summary_df = pd.DataFrame(
        [
            {"indicador": "Linhas", current_name: len(current_df), other_name: len(other_df), "delta": len(current_df) - len(other_df)},
            {"indicador": "Colunas", current_name: current_df.shape[1], other_name: other_df.shape[1], "delta": current_df.shape[1] - other_df.shape[1]},
            {
                "indicador": "Score qualidade",
                current_name: round(current_quality["score"], 2),
                other_name: round(other_quality["score"], 2),
                "delta": round(current_quality["score"] - other_quality["score"], 2),
            },
            {
                "indicador": "Ausencia %",
                current_name: round(current_quality["missing_ratio"] * 100, 2),
                other_name: round(other_quality["missing_ratio"] * 100, 2),
                "delta": round((current_quality["missing_ratio"] - other_quality["missing_ratio"]) * 100, 2),
            },
            {
                "indicador": "Duplicadas",
                current_name: current_quality["duplicates"],
                other_name: other_quality["duplicates"],
                "delta": current_quality["duplicates"] - other_quality["duplicates"],
            },
        ]
    )

    numeric_rows = []
    for col in common_numeric[:30]:
        cur = pd.to_numeric(current_df[col], errors="coerce")
        oth = pd.to_numeric(other_df[col], errors="coerce")
        cur_sum, oth_sum = float(cur.sum(skipna=True)), float(oth.sum(skipna=True))
        cur_mean, oth_mean = float(cur.mean(skipna=True)), float(oth.mean(skipna=True))
        cur_median, oth_median = float(cur.median(skipna=True)), float(oth.median(skipna=True))
        delta_sum = cur_sum - oth_sum
        delta_sum_pct = (delta_sum / abs(oth_sum) * 100.0) if abs(oth_sum) > 1e-9 else None
        numeric_rows.append(
            {
                "coluna": col,
                f"soma_{current_name}": cur_sum,
                f"soma_{other_name}": oth_sum,
                "delta_soma": delta_sum,
                "delta_soma_pct": round(delta_sum_pct, 2) if delta_sum_pct is not None else None,
                f"media_{current_name}": cur_mean,
                f"media_{other_name}": oth_mean,
                "delta_media": cur_mean - oth_mean,
                f"mediana_{current_name}": cur_median,
                f"mediana_{other_name}": oth_median,
                "delta_mediana": cur_median - oth_median,
            }
        )
    numeric_df = pd.DataFrame(numeric_rows).sort_values("delta_soma", key=lambda s: s.abs(), ascending=False) if numeric_rows else pd.DataFrame()

    category_rows = []
    category_delta_rows = []
    ranking_rows = []
    for col in common_categorical[:12]:
        cur_counts = current_df[col].fillna("NULO").astype(str).value_counts()
        oth_counts = other_df[col].fillna("NULO").astype(str).value_counts()
        cur_top = cur_counts.div(max(len(current_df), 1)).head(1)
        oth_top = oth_counts.div(max(len(other_df), 1)).head(1)
        if len(cur_top) == 0 or len(oth_top) == 0:
            continue
        category_rows.append(
            {
                "coluna": col,
                f"top_{current_name}": cur_top.index[0],
                f"participacao_{current_name}_pct": round(float(cur_top.iloc[0]) * 100, 2),
                f"top_{other_name}": oth_top.index[0],
                f"participacao_{other_name}_pct": round(float(oth_top.iloc[0]) * 100, 2),
                "mudou_lider": cur_top.index[0] != oth_top.index[0],
            }
        )
        all_cats = sorted(set(cur_counts.index).union(set(oth_counts.index)))
        for cat in all_cats:
            cur_qtd = int(cur_counts.get(cat, 0))
            oth_qtd = int(oth_counts.get(cat, 0))
            cur_share = cur_qtd / max(len(current_df), 1) * 100.0
            oth_share = oth_qtd / max(len(other_df), 1) * 100.0
            delta_qtd = cur_qtd - oth_qtd
            delta_pct = (delta_qtd / abs(oth_qtd) * 100.0) if oth_qtd else None
            category_delta_rows.append(
                {
                    "coluna": col,
                    "categoria": cat,
                    f"qtd_{current_name}": cur_qtd,
                    f"qtd_{other_name}": oth_qtd,
                    "delta_qtd": delta_qtd,
                    "delta_pct": round(delta_pct, 2) if delta_pct is not None else None,
                    f"participacao_{current_name}_pct": round(cur_share, 2),
                    f"participacao_{other_name}_pct": round(oth_share, 2),
                    "delta_participacao_pct": round(cur_share - oth_share, 2),
                    "status_mudanca": "nova" if oth_qtd == 0 and cur_qtd > 0 else "desapareceu" if cur_qtd == 0 and oth_qtd > 0 else "permanece",
                }
            )
        cur_rank = {cat: i + 1 for i, cat in enumerate(cur_counts.head(10).index)}
        oth_rank = {cat: i + 1 for i, cat in enumerate(oth_counts.head(10).index)}
        for cat in sorted(set(cur_rank).union(oth_rank)):
            ranking_rows.append(
                {
                    "coluna": col,
                    "categoria": cat,
                    f"ranking_{current_name}": cur_rank.get(cat),
                    f"ranking_{other_name}": oth_rank.get(cat),
                    "mudanca_ranking": (oth_rank.get(cat, 99) - cur_rank.get(cat, 99)),
                    "entrou_top10": cat in cur_rank and cat not in oth_rank,
                    "saiu_top10": cat in oth_rank and cat not in cur_rank,
                }
            )
    category_df = pd.DataFrame(category_rows)
    category_delta_df = (
        pd.DataFrame(category_delta_rows).sort_values("delta_qtd", key=lambda s: s.abs(), ascending=False)
        if category_delta_rows
        else pd.DataFrame()
    )
    ranking_df = (
        pd.DataFrame(ranking_rows).sort_values(["entrou_top10", "saiu_top10", "mudanca_ranking"], ascending=[False, False, False])
        if ranking_rows
        else pd.DataFrame()
    )

    current_volume_col = pick_volume_id_column(current_df, current_ids)
    other_volume_col = current_volume_col if current_volume_col in other_df.columns else pick_volume_id_column(other_df, other_ids)
    current_volume = int(current_df[current_volume_col].nunique(dropna=True)) if current_volume_col else int(len(current_df))
    other_volume = int(other_df[other_volume_col].nunique(dropna=True)) if other_volume_col else int(len(other_df))
    volume_delta = current_volume - other_volume
    volume_delta_pct = (volume_delta / abs(other_volume) * 100.0) if other_volume else None
    volume_df = pd.DataFrame(
        [
            {
                "indicador": f"Volume por {'ID ' + current_volume_col if current_volume_col else 'linhas'}",
                current_name: current_volume,
                other_name: other_volume,
                "delta": volume_delta,
                "delta_pct": round(volume_delta_pct, 2) if volume_delta_pct is not None else None,
            }
        ]
    )

    insights = []
    if volume_delta != 0:
        direction = "aumentou" if volume_delta > 0 else "reduziu"
        pct_txt = f" ({abs(volume_delta_pct):.1f}%)" if volume_delta_pct is not None else ""
        insights.append(f"O volume analisado {direction} em {abs(volume_delta):,} registro(s){pct_txt}.".replace(",", "."))
    score_delta = current_quality["score"] - other_quality["score"]
    if abs(score_delta) >= 5:
        direction = "melhorou" if score_delta > 0 else "piorou"
        insights.append(f"O score de qualidade {direction} {abs(score_delta):.1f} ponto(s).")
    if len(category_delta_df) > 0:
        top_change = category_delta_df.iloc[0]
        direction = "cresceu" if int(top_change["delta_qtd"]) > 0 else "reduziu"
        insights.append(
            f"A maior mudanca de perfil ocorreu em `{top_change['coluna']}` = `{top_change['categoria']}`: {direction} {abs(int(top_change['delta_qtd']))} registro(s)."
        )
    if len(category_df) > 0 and bool(category_df["mudou_lider"].any()):
        changed = category_df[category_df["mudou_lider"]].iloc[0]
        insights.append(
            f"A lideranca mudou em `{changed['coluna']}`: antes `{changed[f'top_{other_name}']}`, agora `{changed[f'top_{current_name}']}`."
        )
    if new_cols:
        insights.append("Novas colunas na base atual: " + ", ".join(new_cols[:5]) + ".")
    if removed_cols:
        insights.append("Colunas ausentes na base atual: " + ", ".join(removed_cols[:5]) + ".")
    if len(numeric_df) > 0:
        row = numeric_df.iloc[0]
        insights.append(f"Maior variacao numerica em `{row['coluna']}`: delta de soma {format_br_number(row['delta_soma'])}.")
    if not insights:
        insights.append("As bases comparadas apresentam variacao baixa nos principais indicadores automaticos.")

    alerts = []
    if volume_delta_pct is not None and abs(volume_delta_pct) >= 20:
        alerts.append(
            {
                "prioridade": "P1",
                "alerta": "Variacao relevante de volume",
                "evidencia": f"Delta de {volume_delta} registro(s), {volume_delta_pct:.1f}%.",
                "impacto": "Pode indicar mudanca operacional, recorte diferente ou evolucao real da demanda.",
                "acao": "Validar periodo, filtros e causa operacional da variacao.",
            }
        )
    if len(category_delta_df) > 0:
        material = category_delta_df[
            (category_delta_df["delta_participacao_pct"].abs() >= 10)
            | (category_delta_df["status_mudanca"].isin(["nova", "desapareceu"]))
        ].head(5)
        for _, row in material.iterrows():
            alerts.append(
                {
                    "prioridade": "P1" if abs(float(row["delta_participacao_pct"])) >= 20 else "P2",
                    "alerta": "Mudanca de perfil por categoria",
                    "evidencia": f"{row['coluna']}={row['categoria']} | delta participacao {row['delta_participacao_pct']:.1f}%.",
                    "impacto": "Pode alterar prioridade gerencial e distribuicao de recursos.",
                    "acao": "Investigar causa da mudanca e revisar plano de atendimento/operacao.",
                }
            )
    if abs(score_delta) >= 5:
        alerts.append(
            {
                "prioridade": "P1" if score_delta < 0 else "P2",
                "alerta": "Mudanca no score de qualidade",
                "evidencia": f"Delta de qualidade {score_delta:.1f} ponto(s).",
                "impacto": "Afeta confiabilidade da leitura e das decisoes baseadas na planilha.",
                "acao": "Revisar ausencias, duplicidades, tipos mistos e regras de qualidade.",
            }
        )
    action_df = pd.DataFrame(alerts[:8]) if alerts else pd.DataFrame(
        [
            {
                "prioridade": "P3",
                "alerta": "Sem mudanca critica detectada",
                "evidencia": "A comparacao automatica nao identificou variacoes materiais.",
                "impacto": "Manter acompanhamento recorrente.",
                "acao": "Repetir comparacao no proximo ciclo e monitorar indicadores-chave.",
            }
        ]
    )

    return {
        "summary_df": summary_df,
        "volume_df": volume_df,
        "numeric_df": numeric_df,
        "category_df": category_df,
        "category_delta_df": category_delta_df,
        "ranking_df": ranking_df,
        "action_df": action_df,
        "insights": insights,
        "new_cols": new_cols,
        "removed_cols": removed_cols,
        "current_quality": current_quality,
        "other_quality": other_quality,
        "id_like_cols": sorted(id_like),
    }


def generate_comparison_excel(comparison: dict) -> bytes:
    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        comparison["summary_df"].to_excel(writer, sheet_name="resumo", index=False)
        comparison.get("volume_df", pd.DataFrame()).to_excel(writer, sheet_name="volume", index=False)
        if len(comparison["numeric_df"]) > 0:
            comparison["numeric_df"].to_excel(writer, sheet_name="numericas", index=False)
        if len(comparison["category_df"]) > 0:
            comparison["category_df"].to_excel(writer, sheet_name="categorias", index=False)
        if len(comparison.get("category_delta_df", pd.DataFrame())) > 0:
            comparison["category_delta_df"].to_excel(writer, sheet_name="mudancas_categoria", index=False)
        if len(comparison.get("ranking_df", pd.DataFrame())) > 0:
            comparison["ranking_df"].to_excel(writer, sheet_name="ranking", index=False)
        comparison.get("action_df", pd.DataFrame()).to_excel(writer, sheet_name="plano_acao", index=False)
        pd.DataFrame({"insight": comparison["insights"]}).to_excel(writer, sheet_name="insights", index=False)
        pd.DataFrame({"colunas_novas": comparison["new_cols"]}).to_excel(writer, sheet_name="colunas_novas", index=False)
        pd.DataFrame({"colunas_removidas": comparison["removed_cols"]}).to_excel(writer, sheet_name="colunas_removidas", index=False)
    buffer.seek(0)
    return buffer.getvalue()


def ppt_add_title(slide, title: str, subtitle: str = "") -> None:
    if slide.shapes.title is not None:
        slide.shapes.title.text = clean_report_text(title, 120)
    else:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9.0), Inches(0.5))
        box.text_frame.text = clean_report_text(title, 120)
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = clean_report_text(subtitle, 220)


def ppt_add_bullets(slide, bullets: List[str], left: float = 0.7, top: float = 1.55, width: float = 8.8, height: float = 4.8) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    for i, item in enumerate(bullets[:8]):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = clean_report_text(item, 240)
        p.level = 0
        p.font.size = Pt(18 if i == 0 else 15)


def ppt_add_table(slide, df: pd.DataFrame, left: float, top: float, width: float, height: float, max_rows: int = 6) -> None:
    if df.empty:
        ppt_add_bullets(slide, ["Sem dados suficientes para esta tabela."], left=left, top=top, width=width, height=height)
        return
    data = df.head(max_rows).copy()
    rows, cols = data.shape[0] + 1, data.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    for j, col in enumerate(data.columns):
        table.cell(0, j).text = clean_report_text(col, 28)
    for i, (_, row) in enumerate(data.iterrows(), start=1):
        for j, col in enumerate(data.columns):
            table.cell(i, j).text = clean_report_text(row[col], 38)
    for i in range(rows):
        for j in range(cols):
            for paragraph in table.cell(i, j).text_frame.paragraphs:
                paragraph.font.size = Pt(9)


def generate_powerpoint_report(
    *,
    source_name: str,
    sheet_name: str,
    quality_report: dict,
    insights: List[str],
    issue_catalog: pd.DataFrame,
    dashboard_bundle: Optional[dict],
    treatment_report: dict,
    ml_sup: Optional[dict],
    version: int,
) -> bytes:
    if Presentation is None:
        raise RuntimeError("Dependencia python-pptx nao instalada.")
    prs = Presentation()
    prs.core_properties.author = CREATOR_SIGNATURE

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    ppt_add_title(slide, "AutoAnalista 2026", f"Relatorio executivo | {source_name} | {sheet_name} | v{version}")

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    ppt_add_title(slide, "Resumo Executivo", f"Score: {quality_report['score']:.1f}/100 ({quality_report['level']})")
    ppt_add_bullets(
        slide,
        [
            f"Linhas finais analisadas: {treatment_report.get('rows_after', 0):,}".replace(",", "."),
            f"Linhas removidas na tratativa: {treatment_report.get('total_removed', 0):,}".replace(",", "."),
            f"Duplicadas removidas: {treatment_report.get('duplicates_removed', 0):,}".replace(",", "."),
            f"Outliers removidos: {treatment_report.get('outliers_removed', 0):,}".replace(",", "."),
        ]
        + insights[:4],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    ppt_add_title(slide, "Dashboard Gerencial", "KPIs e principais categorias")
    if dashboard_bundle:
        kpi_df = dashboard_bundle.get("kpi_df", pd.DataFrame())
        top_cat_df = dashboard_bundle.get("top_categories_df", pd.DataFrame())
        ppt_add_table(slide, kpi_df, 0.4, 1.2, 4.4, 2.2, max_rows=5)
        ppt_add_table(slide, top_cat_df[["categoria", "quantidade", "participacao_pct"]] if len(top_cat_df) > 0 else pd.DataFrame(), 5.0, 1.2, 4.6, 2.8, max_rows=6)
    else:
        ppt_add_bullets(slide, ["Dashboard nao preparado para esta exportacao."])

    if dashboard_bundle:
        profile_df = dashboard_bundle.get("profile_df", pd.DataFrame())
        if len(profile_df) > 0:
            profile_map = dict(zip(profile_df.get("campo", []), profile_df.get("valor", [])))
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            ppt_add_title(slide, "Perfil Inteligente da Planilha", "Tipo detectado e melhor abordagem analitica")
            ppt_add_bullets(
                slide,
                [
                    f"Tipo detectado: {profile_map.get('tipo_detectado', '-')}",
                    f"Confianca: {profile_map.get('confianca_pct', '-')}%",
                    f"Foco recomendado: {profile_map.get('foco', '-')}",
                    f"Melhor analise: {profile_map.get('melhor_analise', '-')}",
                    f"Graficos indicados: {profile_map.get('graficos_recomendados', '-')}",
                ],
            )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    ppt_add_title(slide, "Qualidade e Riscos", "Principais pontos de atencao")
    risk_bullets = []
    if len(issue_catalog) > 0:
        for _, row in issue_catalog.sort_values(["priority_rank", "affected_rows"], ascending=[True, False]).head(6).iterrows():
            risk_bullets.append(f"{row['priority']} | {row['issue']} em {row['column']} | {row['affected_rows']} afetados")
    else:
        risk_bullets.append("Sem riscos relevantes catalogados automaticamente.")
    ppt_add_bullets(slide, risk_bullets)

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    ppt_add_title(slide, "Plano de Acao", "Prioridades sugeridas")
    action_df = build_structured_action_plan("Geral", quality_report, issue_catalog)
    ppt_add_bullets(slide, [f"{r['prioridade']} | {r['acao']} | {r['prazo']}" for _, r in action_df.head(6).iterrows()])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    ppt_add_title(slide, "Machine Learning", "Resultado preditivo quando disponivel")
    if ml_sup:
        bullets = [f"Modelo: {ml_sup.get('model_choice', '-')}", f"Alvo: {ml_sup.get('target_col', '-')}"]
        if ml_sup.get("requested_model_choice") == "Auto (comparar modelos)":
            bullets.append("Selecao automatica: modelos comparados por validacao cruzada.")
        if ml_sup.get("problem_type") == "classification":
            bullets.append(f"F1 holdout: {ml_sup.get('holdout_f1', 0):.3f}")
        else:
            bullets.append(f"R2 holdout: {ml_sup.get('holdout_r2', 0):.3f}")
        bullets.append(ml_sup.get("model_reliability_alert", "Sem alerta de confiabilidade."))
    else:
        bullets = ["Modelo supervisionado nao executado nesta analise."]
    ppt_add_bullets(slide, bullets)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def generate_offline_package() -> bytes:
    buffer = io.BytesIO()
    run_bat = """@echo off
cd /d %~dp0
python -m venv .venv
call .venv\\Scripts\\activate
python -m pip install --upgrade pip
pip install -r requirements.txt
streamlit run app.py
pause
"""
    install_md = f"""# AutoAnalista 2026 - pacote offline

Criador: {CREATOR_SIGNATURE}

## Como instalar no computador
1. Extraia este ZIP em uma pasta local.
2. Instale Python 3.11 ou superior, se ainda nao tiver.
3. Dê duplo clique em `INICIAR_AUTOANALISTA.bat`.
4. Aguarde a instalacao das dependencias.
5. O app abrira em `http://localhost:8501`.

Depois da primeira instalacao, o app funciona localmente. O banco SQLite fica em `app_data/autoanalista.db`.
"""
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.write(APP_DIR / "app.py", "app.py")
        zf.write(APP_DIR / "requirements.txt", "requirements.txt")
        zf.writestr("INICIAR_AUTOANALISTA.bat", run_bat)
        zf.writestr("LEIA-ME_INSTALACAO_OFFLINE.md", install_md)
    buffer.seek(0)
    return buffer.getvalue()


def generate_excel_export(
    df: pd.DataFrame,
    quality_report: dict,
    column_quality: pd.DataFrame,
    issue_catalog: pd.DataFrame,
    violations: pd.DataFrame,
    anomalias_df: pd.DataFrame,
    insights: List[str],
    recommendations: List[str],
    raw_df: Optional[pd.DataFrame] = None,
    treatment_report: Optional[dict] = None,
    dashboard_bundle: Optional[dict] = None,
) -> bytes:
    buffer = io.BytesIO()
    pillars_df = pd.DataFrame({"Pilar": list(quality_report["pillars"].keys()), "Score": list(quality_report["pillars"].values())})
    summary_df = pd.DataFrame(
        [
            {"Metrica": "Score geral", "Valor": quality_report["score"]},
            {"Metrica": "Nivel", "Valor": quality_report["level"]},
            {"Metrica": "Duplicados", "Valor": quality_report["duplicates"]},
            {"Metrica": "Outlier ratio", "Valor": quality_report["outlier_ratio"]},
        ]
    )
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        if raw_df is not None and len(raw_df) > 0:
            raw_df.head(MAX_EXPORT_ROWS).to_excel(writer, sheet_name="dados_brutos_filtrados", index=False)
        df.head(MAX_EXPORT_ROWS).to_excel(writer, sheet_name="dados_tratados", index=False)
        summary_df.to_excel(writer, sheet_name="resumo_qualidade", index=False)
        pillars_df.to_excel(writer, sheet_name="pilares", index=False)
        column_quality.to_excel(writer, sheet_name="qualidade_colunas", index=False)
        issue_catalog.to_excel(writer, sheet_name="catalogo_problemas", index=False)
        if treatment_report:
            treatment_df = pd.DataFrame(
                [
                    {"metrica": "linhas_antes", "valor": treatment_report.get("rows_before", 0)},
                    {"metrica": "linhas_depois", "valor": treatment_report.get("rows_after", 0)},
                    {"metrica": "vazias_removidas", "valor": treatment_report.get("empty_rows_removed", 0)},
                    {"metrica": "duplicadas_removidas", "valor": treatment_report.get("duplicates_removed", 0)},
                    {"metrica": "outliers_removidos", "valor": treatment_report.get("outliers_removed", 0)},
                    {"metrica": "remocao_total", "valor": treatment_report.get("total_removed", 0)},
                    {"metrica": "missing_ratio_antes", "valor": treatment_report.get("missing_ratio_before", 0)},
                    {"metrica": "missing_ratio_depois", "valor": treatment_report.get("missing_ratio_after", 0)},
                ]
            )
            treatment_df.to_excel(writer, sheet_name="tratativa_dados", index=False)
        if len(violations) > 0:
            violations.head(MAX_EXPORT_ROWS).to_excel(writer, sheet_name="violacoes_regras", index=False)
        if len(anomalias_df) > 0:
            anomalias_df.head(MAX_EXPORT_ROWS).to_excel(writer, sheet_name="anomalias_ml", index=False)
        pd.DataFrame({"insight": insights}).to_excel(writer, sheet_name="insights", index=False)
        pd.DataFrame({"acao_recomendada": recommendations}).to_excel(writer, sheet_name="recomendacoes", index=False)
        if dashboard_bundle:
            dashboard_bundle.get("summary_df", pd.DataFrame()).to_excel(writer, sheet_name="dash_resumo", index=False)
            dashboard_bundle.get("profile_df", pd.DataFrame()).to_excel(writer, sheet_name="dash_perfil", index=False)
            dashboard_bundle.get("filters_df", pd.DataFrame()).to_excel(writer, sheet_name="dash_filtros", index=False)
            dashboard_bundle.get("narrative_df", pd.DataFrame()).to_excel(writer, sheet_name="dash_narrativa", index=False)
            dashboard_bundle.get("analytic_reading_df", pd.DataFrame()).to_excel(writer, sheet_name="dash_leitura", index=False)
            dashboard_bundle.get("kpi_df", pd.DataFrame()).to_excel(writer, sheet_name="dash_kpis", index=False)
            top_cat_df = dashboard_bundle.get("top_categories_df", pd.DataFrame())
            if len(top_cat_df) > 0:
                top_cat_df.to_excel(writer, sheet_name="dash_top_categorias", index=False)
            metric_cat_df = dashboard_bundle.get("metric_by_category_df", pd.DataFrame())
            if len(metric_cat_df) > 0:
                metric_cat_df.to_excel(writer, sheet_name="dash_metrica_categoria", index=False)
            dist_df = dashboard_bundle.get("numeric_distribution_df", pd.DataFrame())
            if len(dist_df) > 0:
                dist_df.to_excel(writer, sheet_name="dash_dist_numerica", index=False)
            trend_df = dashboard_bundle.get("trend_df", pd.DataFrame())
            if len(trend_df) > 0:
                trend_df.to_excel(writer, sheet_name="dash_tendencia", index=False)
    buffer.seek(0)
    return buffer.getvalue()


def pdf_table(data: List[List[object]], col_widths: List[int], header_color: str = "#1F4E79", font_size: int = 8) -> Table:
    clean_data = [[clean_report_text(cell, 120) for cell in row] for row in data]
    table = Table(clean_data, colWidths=col_widths, repeatRows=1)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(header_color)),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#CBD5E1")),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), font_size),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F8FAFC")]),
            ]
        )
    )
    return table


def make_horizontal_bar_chart(
    labels: List[object],
    values: List[object],
    title: str,
    width: int = 500,
    height: int = 210,
    color: str = "#4EA3F1",
) -> Optional[Drawing]:
    pairs = []
    for label, value in zip(labels, values):
        try:
            val = float(value)
        except Exception:
            continue
        if pd.notna(val):
            pairs.append((clean_report_text(label, 34), val))
    pairs = pairs[:8]
    if not pairs:
        return None
    pairs = list(reversed(pairs))
    chart = HorizontalBarChart()
    chart.x = 150
    chart.y = 25
    chart.height = height - 62
    chart.width = width - 185
    chart.data = [[p[1] for p in pairs]]
    chart.categoryAxis.categoryNames = [p[0] for p in pairs]
    chart.categoryAxis.labels.fontSize = 6
    chart.valueAxis.labels.fontSize = 7
    chart.valueAxis.valueMin = 0
    chart.valueAxis.valueMax = max(p[1] for p in pairs) * 1.15 if max(p[1] for p in pairs) > 0 else 1
    chart.bars[0].fillColor = colors.HexColor(color)
    drawing = Drawing(width, height)
    drawing.add(String(0, height - 18, clean_report_text(title, 90), fontSize=10, fillColor=colors.HexColor("#0F172A")))
    drawing.add(chart)
    return drawing


def make_vertical_bar_chart(
    labels: List[object],
    values: List[object],
    title: str,
    width: int = 500,
    height: int = 210,
    color: str = "#1B7F5A",
) -> Optional[Drawing]:
    pairs = []
    for label, value in zip(labels, values):
        try:
            val = float(value)
        except Exception:
            continue
        if pd.notna(val):
            pairs.append((clean_report_text(label, 12), val))
    pairs = pairs[-10:]
    if not pairs:
        return None
    chart = VerticalBarChart()
    chart.x = 35
    chart.y = 35
    chart.height = height - 75
    chart.width = width - 60
    chart.data = [[p[1] for p in pairs]]
    chart.categoryAxis.categoryNames = [p[0] for p in pairs]
    chart.categoryAxis.labels.fontSize = 6
    chart.categoryAxis.labels.angle = 30
    chart.valueAxis.labels.fontSize = 7
    chart.valueAxis.valueMin = 0
    chart.valueAxis.valueMax = max(p[1] for p in pairs) * 1.15 if max(p[1] for p in pairs) > 0 else 1
    chart.bars[0].fillColor = colors.HexColor(color)
    drawing = Drawing(width, height)
    drawing.add(String(0, height - 18, clean_report_text(title, 90), fontSize=10, fillColor=colors.HexColor("#0F172A")))
    drawing.add(chart)
    return drawing


def make_pie_chart(labels: List[object], values: List[object], title: str, width: int = 500, height: int = 210) -> Optional[Drawing]:
    pairs = []
    for label, value in zip(labels, values):
        try:
            val = float(value)
        except Exception:
            continue
        if pd.notna(val) and val > 0:
            pairs.append((clean_report_text(label, 24), val))
    if not pairs:
        return None
    top = pairs[:6]
    if len(pairs) > 6:
        top.append(("Outros", sum(v for _, v in pairs[6:])))
    pie = Pie()
    pie.x = 35
    pie.y = 25
    pie.width = 145
    pie.height = 145
    pie.data = [p[1] for p in top]
    pie.labels = [p[0] for p in top]
    pie.sideLabels = 1
    palette = ["#4EA3F1", "#1B7F5A", "#F59E0B", "#EF4444", "#8B5CF6", "#14B8A6", "#94A3B8"]
    for i, color in enumerate(palette[: len(top)]):
        pie.slices[i].fillColor = colors.HexColor(color)
    drawing = Drawing(width, height)
    drawing.add(String(0, height - 18, clean_report_text(title, 90), fontSize=10, fillColor=colors.HexColor("#0F172A")))
    drawing.add(pie)
    return drawing


def append_if_chart(story: List[object], chart: Optional[Drawing]) -> None:
    if chart is not None:
        story.append(chart)
        story.append(Spacer(1, 8))


def build_executive_sections(
    insights: List[str],
    dashboard_bundle: Optional[dict],
    quality_report: dict,
    issue_catalog: pd.DataFrame,
    action_df: pd.DataFrame,
) -> Dict[str, List[str]]:
    analytic_df = dashboard_bundle.get("analytic_reading_df", pd.DataFrame()) if dashboard_bundle else pd.DataFrame()
    analytic_notes = analytic_df["leitura_analitica"].dropna().astype(str).tolist() if len(analytic_df) > 0 else []

    what_happened = (analytic_notes[:2] + insights[:2])[:3]
    if not what_happened:
        what_happened = [f"A base foi analisada com score de qualidade {quality_report.get('score', 0.0):.1f}/100."]

    risks: List[str] = []
    if len(issue_catalog) > 0:
        for _, row in issue_catalog.sort_values(["priority_rank", "affected_rows"], ascending=[True, False]).head(3).iterrows():
            risks.append(
                f"{row['issue']} em {row['column']}: {int(row['affected_rows'])} registro(s) afetado(s)."
            )
    if not risks:
        risks.append("Nao foram encontrados riscos criticos no catalogo automatico, mantendo monitoramento por lote.")

    opportunities: List[str] = []
    if quality_report.get("score", 0) >= 80:
        opportunities.append("A qualidade atual permite uso gerencial mais seguro e evolucao para monitoramento recorrente.")
    else:
        opportunities.append("A maior oportunidade esta em elevar completude, consistencia e padronizacao antes de automatizar decisoes.")
    opportunities.extend([note for note in analytic_notes if "concentracao" in normalize_token_text(note) or "distribuicao" in normalize_token_text(note)][:2])

    actions: List[str] = []
    for _, row in action_df.head(3).iterrows():
        actions.append(
            f"{row['prioridade']} - {row['acao']} | Responsavel: {row.get('responsavel', 'Dono do dado')} | Prazo: {row['prazo']}"
        )
    if not actions:
        actions.append("Manter rotina de validacao e acompanhamento de indicadores por periodo.")

    return {
        "O que aconteceu": what_happened[:3],
        "Principais riscos": risks[:3],
        "Principais oportunidades": opportunities[:3],
        "Acoes recomendadas": actions[:3],
    }


def build_pdf_report(
    *,
    source_name: str,
    sheet_name: str,
    quality_report: dict,
    insights: List[str],
    recommendations: List[str],
    top_corr: pd.DataFrame,
    issue_catalog: pd.DataFrame,
    ml_unsup: Optional[dict],
    ml_sup: Optional[dict],
    version: int,
    dashboard_bundle: Optional[dict] = None,
    treatment_report: Optional[dict] = None,
    active_filters: Optional[List[str]] = None,
    area: str = "Geral",
) -> bytes:
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=36)
    styles = getSampleStyleSheet()
    story = []

    score = float(quality_report.get("score", 0.0))
    level = str(quality_report.get("level", "-"))
    summary_df = dashboard_bundle.get("summary_df", pd.DataFrame()) if dashboard_bundle else pd.DataFrame()
    summary_map = dict(zip(summary_df.get("metrica", []), summary_df.get("valor", []))) if len(summary_df) > 0 else {}
    rows_analyzed = summary_map.get("linhas_analisadas", treatment_report.get("rows_after", "-") if treatment_report else "-")
    cols_analyzed = summary_map.get("colunas_analisadas", "-")
    risk_count = summary_map.get("riscos_criticos_altos", 0)
    active_filters = active_filters or []
    action_df = build_structured_action_plan(area, quality_report, issue_catalog)

    cover = Table(
        [
            [pdf_paragraph("AutoAnalista 2026", styles["Title"])],
            [pdf_paragraph("Relatorio Executivo de Analise de Dados", styles["Heading2"])],
            [pdf_paragraph(f"Arquivo: {source_name} | Tabela: {sheet_name} | Versao: v{version}", styles["Normal"])],
            [pdf_paragraph(f"Criador: {CREATOR_SIGNATURE} | Gerado em: {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}", styles["Normal"])],
        ],
        colWidths=[500],
    )
    cover.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#EAF4FF")),
                ("BOX", (0, 0), (-1, -1), 1.0, colors.HexColor("#1F4E79")),
                ("LEFTPADDING", (0, 0), (-1, -1), 16),
                ("RIGHTPADDING", (0, 0), (-1, -1), 16),
                ("TOPPADDING", (0, 0), (-1, -1), 10),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
            ]
        )
    )
    story.append(cover)
    story.append(Spacer(1, 12))

    story.append(pdf_paragraph("Relatorio Executivo", styles["Heading1"]))
    story.append(Spacer(1, 6))

    executive_metrics = [
        ["Indicador", "Resultado"],
        ["Score de qualidade", f"{score:.1f}/100 ({level})"],
        ["Linhas analisadas", f"{rows_analyzed}"],
        ["Colunas analisadas", f"{cols_analyzed}"],
        ["Riscos criticos/altos", f"{risk_count}"],
    ]
    if treatment_report:
        executive_metrics.append(["Linhas removidas na tratativa", str(treatment_report.get("total_removed", 0))])
    story.append(pdf_table(executive_metrics, [240, 240], header_color="#0F6E9A", font_size=9))
    story.append(Spacer(1, 10))

    filter_text = " | ".join(active_filters[:10]) if active_filters else "Nenhum filtro global ou de dashboard ativo no momento da geracao."
    story.append(pdf_paragraph("Filtros usados no relatorio: " + filter_text, styles["Normal"]))
    story.append(Spacer(1, 6))

    story.append(pdf_paragraph("Resumo executivo", styles["Heading2"]))
    executive_sections = build_executive_sections(insights, dashboard_bundle, quality_report, issue_catalog, action_df)
    for section, bullets in executive_sections.items():
        story.append(pdf_paragraph(section, styles["Heading3"]))
        for item in bullets:
            story.append(pdf_paragraph(f"- {item}", styles["Normal"]))
    story.append(Spacer(1, 8))

    if dashboard_bundle:
        analytic_df = dashboard_bundle.get("analytic_reading_df", pd.DataFrame())
        if len(analytic_df) > 0:
            story.append(pdf_paragraph("Leitura Analitica", styles["Heading2"]))
            for _, row in analytic_df.head(6).iterrows():
                story.append(pdf_paragraph(f"- {row.get('leitura_analitica', '-')}", styles["Normal"]))
            story.append(Spacer(1, 8))

    story.append(pdf_paragraph("Plano de acao gerencial", styles["Heading2"]))
    action_data = [["Prioridade", "Problema", "Impacto", "Acao", "Responsavel", "Prazo"]]
    for _, row in action_df.head(6).iterrows():
        action_data.append(
            [
                row["prioridade"],
                row.get("problema", row.get("foco", "-")),
                row["impacto"],
                row["acao"],
                row.get("responsavel", "Dono do dado"),
                row["prazo"],
            ]
        )
    story.append(pdf_table(action_data, [45, 90, 95, 135, 80, 55], header_color="#B45309", font_size=6))
    story.append(Spacer(1, 8))

    story.append(PageBreak())

    story.append(pdf_paragraph("Relatorio Tecnico", styles["Heading1"]))
    story.append(Spacer(1, 6))
    story.append(pdf_paragraph("Qualidade e Graficos Executivos", styles["Heading2"]))
    quality_data = [["Pilar", "Score"]] + [[k, f"{v:.1f}"] for k, v in quality_report["pillars"].items()]
    story.append(pdf_table(quality_data, [250, 120], header_color="#1F4E79", font_size=9))
    append_if_chart(
        story,
        make_horizontal_bar_chart(
            list(quality_report["pillars"].keys()),
            list(quality_report["pillars"].values()),
            "Score por pilar de qualidade",
            color="#1F4E79",
        ),
    )

    if treatment_report:
        treatment_data = [
            ["Etapa", "Quantidade"],
            ["Linhas antes", treatment_report.get("rows_before", 0)],
            ["Linhas depois", treatment_report.get("rows_after", 0)],
            ["Vazias removidas", treatment_report.get("empty_rows_removed", 0)],
            ["Duplicadas removidas", treatment_report.get("duplicates_removed", 0)],
            ["Outliers removidos", treatment_report.get("outliers_removed", 0)],
            ["Total removido", treatment_report.get("total_removed", 0)],
            ["Ausencia antes", f"{float(treatment_report.get('missing_ratio_before', 0.0)) * 100:.2f}%"],
            ["Ausencia depois", f"{float(treatment_report.get('missing_ratio_after', 0.0)) * 100:.2f}%"],
        ]
        story.append(pdf_paragraph("Tratativa dos dados", styles["Heading3"]))
        story.append(pdf_table(treatment_data, [240, 120], header_color="#334155", font_size=8))
        story.append(Spacer(1, 8))

    story.append(pdf_paragraph("Correlacoes de destaque", styles["Heading2"]))
    if len(top_corr) > 0:
        data = [["Variavel 1", "Variavel 2", "Correlacao"]]
        for _, row in top_corr.head(8).iterrows():
            data.append([row["variavel_1"], row["variavel_2"], f"{row['correlacao']:.3f}"])
        story.append(pdf_table(data, [170, 170, 100], header_color="#3A6EA5", font_size=8))
        append_if_chart(
            story,
            make_horizontal_bar_chart(
                [f"{row['variavel_1']} x {row['variavel_2']}" for _, row in top_corr.head(6).iterrows()],
                [abs(float(row["correlacao"])) for _, row in top_corr.head(6).iterrows()],
                "Forca das principais correlacoes",
                color="#3A6EA5",
            ),
        )
    else:
        story.append(pdf_paragraph("Sem correlacoes calculaveis.", styles["Normal"]))
    story.append(Spacer(1, 8))

    story.append(pdf_paragraph("Catalogo de problemas (top)", styles["Heading2"]))
    if len(issue_catalog) > 0:
        data = [["Prioridade", "Problema", "Coluna", "Afetados"]]
        for _, row in issue_catalog.sort_values(["priority_rank", "affected_rows"], ascending=[True, False]).head(8).iterrows():
            data.append([row["priority"], row["issue"], row["column"], str(row["affected_rows"])])
        story.append(pdf_table(data, [80, 170, 130, 70], header_color="#7A4EAB", font_size=8))
    else:
        story.append(pdf_paragraph("Sem problemas catalogados.", styles["Normal"]))
    story.append(Spacer(1, 8))

    if dashboard_bundle:
        story.append(PageBreak())
        story.append(pdf_paragraph("Analise do Dashboard", styles["Heading2"]))
        profile_df = dashboard_bundle.get("profile_df", pd.DataFrame())
        if len(profile_df) > 0:
            profile_map = dict(zip(profile_df.get("campo", []), profile_df.get("valor", [])))
            story.append(pdf_paragraph("Perfil automatico da planilha", styles["Heading3"]))
            profile_data = [
                ["Campo", "Leitura"],
                ["Tipo detectado", profile_map.get("tipo_detectado", "-")],
                ["Confianca", f"{profile_map.get('confianca_pct', '-')}%"],
                ["Foco recomendado", profile_map.get("foco", "-")],
                ["Melhor analise", profile_map.get("melhor_analise", "-")],
                ["Graficos indicados", profile_map.get("graficos_recomendados", "-")],
            ]
            story.append(pdf_table(profile_data, [130, 350], header_color="#0F6E9A", font_size=7))
            story.append(Spacer(1, 6))

        filters_df = dashboard_bundle.get("filters_df", pd.DataFrame())
        if len(filters_df) > 0:
            story.append(pdf_paragraph("Filtros do Dashboard registrados", styles["Heading3"]))
            filter_data = [["Filtro"]] + [[row.get("filtro", "-")] for _, row in filters_df.head(10).iterrows()]
            story.append(pdf_table(filter_data, [460], header_color="#334155", font_size=8))
            story.append(Spacer(1, 6))

        narrative_df = dashboard_bundle.get("narrative_df", pd.DataFrame())
        if len(narrative_df) > 0:
            for _, row in narrative_df.head(8).iterrows():
                story.append(pdf_paragraph(f"- {row['analise']}", styles["Normal"]))
        else:
            story.append(pdf_paragraph("Sem narrativa analitica disponivel.", styles["Normal"]))
        story.append(Spacer(1, 6))

        kpi_df = dashboard_bundle.get("kpi_df", pd.DataFrame())
        if len(kpi_df) > 0:
            story.append(pdf_paragraph("KPIs do Dashboard", styles["Heading3"]))
            kpi_data = [["KPI", "Valor", "Variacao", "Alerta"]]
            for _, row in kpi_df.head(8).iterrows():
                kpi_data.append([row.get("kpi", "-"), row.get("valor", "-"), row.get("variacao_recente", "-"), row.get("alerta", "-")])
            story.append(pdf_table(kpi_data, [150, 70, 140, 120], header_color="#0F6E9A", font_size=8))
            story.append(Spacer(1, 6))

        top_cat_df = dashboard_bundle.get("top_categories_df", pd.DataFrame())
        if len(top_cat_df) > 0:
            story.append(pdf_paragraph("Top categorias extraidas", styles["Heading3"]))
            append_if_chart(
                story,
                make_pie_chart(
                    top_cat_df["categoria"].head(10).tolist(),
                    top_cat_df["quantidade"].head(10).tolist(),
                    "Participacao das principais categorias",
                ),
            )
            append_if_chart(
                story,
                make_horizontal_bar_chart(
                    top_cat_df["categoria"].head(10).tolist(),
                    top_cat_df["quantidade"].head(10).tolist(),
                    "Top categorias por quantidade",
                    color="#1B7F5A",
                ),
            )
            cat_data = [["Categoria", "Quantidade", "Participacao %"]]
            for _, row in top_cat_df.head(8).iterrows():
                cat_data.append([row.get("categoria", "-"), row.get("quantidade", "-"), f"{float(row.get('participacao_pct', 0.0)):.2f}%"])
            story.append(pdf_table(cat_data, [220, 100, 120], header_color="#1B7F5A", font_size=8))
            story.append(Spacer(1, 6))

        metric_cat_df = dashboard_bundle.get("metric_by_category_df", pd.DataFrame())
        if len(metric_cat_df) > 0:
            metric_name = str(metric_cat_df["coluna_metrica"].iloc[0]) if "coluna_metrica" in metric_cat_df.columns else "metrica"
            append_if_chart(
                story,
                make_horizontal_bar_chart(
                    metric_cat_df["categoria"].head(10).tolist(),
                    metric_cat_df["valor_metrica"].head(10).tolist(),
                    f"Comparativo por categoria - {metric_name}",
                    color="#F59E0B",
                ),
            )

        numeric_dist_df = dashboard_bundle.get("numeric_distribution_df", pd.DataFrame())
        if len(numeric_dist_df) > 0:
            dist_col = str(numeric_dist_df["coluna"].iloc[0]) if "coluna" in numeric_dist_df.columns else "metrica"
            story.append(pdf_paragraph("Distribuicao numerica", styles["Heading3"]))
            append_if_chart(
                story,
                make_vertical_bar_chart(
                    numeric_dist_df["faixa"].head(10).tolist(),
                    numeric_dist_df["quantidade"].head(10).tolist(),
                    f"Distribuicao de {dist_col}",
                    color="#4EA3F1",
                ),
            )
            dist_data = [["Faixa", "Quantidade", "Participacao %"]]
            for _, row in numeric_dist_df.head(8).iterrows():
                dist_data.append([row.get("faixa", "-"), row.get("quantidade", "-"), f"{float(row.get('participacao_pct', 0.0)):.2f}%"])
            story.append(pdf_table(dist_data, [210, 100, 120], header_color="#4EA3F1", font_size=8))
            story.append(Spacer(1, 6))

        trend_df = dashboard_bundle.get("trend_df", pd.DataFrame())
        if len(trend_df) > 0:
            story.append(pdf_paragraph("Tendencia temporal extraida", styles["Heading3"]))
            append_if_chart(
                story,
                make_vertical_bar_chart(
                    trend_df.iloc[:, 0].tail(10).astype(str).tolist(),
                    trend_df["valor"].tail(10).tolist(),
                    dashboard_bundle.get("trend_label", "Tendencia temporal"),
                    color="#8A6D1F",
                ),
            )
            trend_cols = trend_df.columns.tolist()
            trend_data = [[str(trend_cols[0]), "Indicador", "Valor"]]
            for _, row in trend_df.tail(10).iterrows():
                trend_data.append([str(row.iloc[0]), str(row.get("indicador", "-")), f"{float(row.get('valor', 0.0)):.2f}"])
            story.append(pdf_table(trend_data, [150, 190, 100], header_color="#8A6D1F", font_size=8))
            story.append(Spacer(1, 8))

    story.append(pdf_paragraph("Resumo de Machine Learning", styles["Heading2"]))
    if ml_unsup:
        story.append(
            pdf_paragraph(
                f"Sem alvo: k={ml_unsup['best_k']}, silhouette={ml_unsup['silhouette']:.3f}, anomalias={ml_unsup['anomalias']} (taxa={ml_unsup['contamination']:.2f}).",
                styles["Normal"],
            )
        )
    else:
        story.append(pdf_paragraph("Sem resultado de ML sem alvo.", styles["Normal"]))
    if ml_sup:
        if ml_sup["problem_type"] == "classification":
            story.append(
                pdf_paragraph(
                    f"Com alvo ({ml_sup['target_col']}) usando {ml_sup.get('model_choice', 'modelo preditivo')}: CV Acuracia={ml_sup['cv_accuracy']:.3f}, CV F1={ml_sup['cv_f1']:.3f}, Holdout Acuracia={ml_sup['holdout_accuracy']:.3f}.",
                    styles["Normal"],
                )
            )
        else:
            story.append(
                pdf_paragraph(
                    f"Com alvo ({ml_sup['target_col']}) usando {ml_sup.get('model_choice', 'modelo preditivo')}: CV R2={ml_sup['cv_r2']:.3f}, CV RMSE={ml_sup['cv_rmse']:.3f}, Holdout R2={ml_sup['holdout_r2']:.3f}.",
                    styles["Normal"],
                )
            )
        business_notes = ml_business_explanation(ml_sup)
        if business_notes:
            story.append(pdf_paragraph("Leitura gerencial do ML", styles["Heading3"]))
            for note in business_notes[:5]:
                story.append(pdf_paragraph("- " + note, styles["Normal"]))
        model_comparison = ml_sup.get("model_comparison", [])
        if model_comparison:
            comparison_df = pd.DataFrame(model_comparison)
            if len(comparison_df) > 0:
                story.append(pdf_paragraph("Comparacao de modelos testados", styles["Heading3"]))
                if ml_sup.get("problem_type") == "classification":
                    preferred_cols = ["modelo", "f1_cv", "acuracia_cv", "status"]
                    header_map = {"modelo": "Modelo", "f1_cv": "F1 CV", "acuracia_cv": "Acuracia CV", "status": "Status"}
                else:
                    preferred_cols = ["modelo", "r2_cv", "rmse_cv", "mae_cv", "status"]
                    header_map = {"modelo": "Modelo", "r2_cv": "R2 CV", "rmse_cv": "RMSE CV", "mae_cv": "MAE CV", "status": "Status"}
                available_cols = [c for c in preferred_cols if c in comparison_df.columns]
                if available_cols:
                    model_data = [[header_map.get(c, c) for c in available_cols]]
                    for _, row in comparison_df[available_cols].head(6).iterrows():
                        model_data.append(
                            [
                                f"{float(value):.3f}" if isinstance(value, (int, float, np.floating)) and pd.notna(value) else value
                                for value in row.tolist()
                            ]
                        )
                    story.append(pdf_table(model_data, [120, 75, 75, 75, 95][: len(available_cols)], header_color="#334155", font_size=7))
        if ml_sup.get("model_reliability_alert"):
            story.append(pdf_paragraph("Alerta de confiabilidade: " + ml_sup["model_reliability_alert"], styles["Normal"]))
        feature_importance = ml_sup.get("feature_importance", [])
        if feature_importance:
            story.append(pdf_paragraph("Principais variaveis explicativas", styles["Heading3"]))
            feature_data = [["Variavel", "Importancia"]]
            for item in feature_importance[:10]:
                feature_data.append([item.get("feature", "-"), f"{float(item.get('importance', 0.0)):.4f}"])
            story.append(pdf_table(feature_data, [280, 120], header_color="#0F6E9A", font_size=8))
        if ml_sup.get("removed_leakage_features"):
            story.append(
                pdf_paragraph(
                    "Colunas removidas por risco de leakage: " + ", ".join(map(str, ml_sup.get("removed_leakage_features", [])[:8])),
                    styles["Normal"],
                )
            )
        if ml_sup.get("problem_type") == "classification" and ml_sup.get("confusion_matrix"):
            cm_df = pd.DataFrame(ml_sup["confusion_matrix"])
            if len(cm_df) > 0:
                story.append(pdf_paragraph("Matriz de confusao do holdout", styles["Heading3"]))
                cm_cols = cm_df.columns.tolist()[:7]
                cm_data = [cm_cols]
                for _, row in cm_df[cm_cols].head(10).iterrows():
                    cm_data.append(row.tolist())
                story.append(pdf_table(cm_data, [95] + [55] * (len(cm_cols) - 1), header_color="#7A4EAB", font_size=6))
        if ml_sup.get("problem_type") == "regression" and ml_sup.get("prediction_sample"):
            pred_df = pd.DataFrame(ml_sup["prediction_sample"])
            if len(pred_df) > 0:
                story.append(pdf_paragraph("Amostra real vs previsto", styles["Heading3"]))
                pred_cols = [c for c in ["real", "previsto", "erro"] if c in pred_df.columns]
                pred_data = [["Real", "Previsto", "Erro"][: len(pred_cols)]]
                for _, row in pred_df[pred_cols].head(10).iterrows():
                    pred_data.append([f"{float(value):.3f}" if pd.notna(value) else "-" for value in row.tolist()])
                story.append(pdf_table(pred_data, [110, 110, 110][: len(pred_cols)], header_color="#0F6E9A", font_size=8))
        story.append(
            pdf_paragraph(
                "Nota: os resultados de ML devem ser usados como apoio a decisao. Para uso operacional critico, valide o modelo com dados historicos, regras de negocio e acompanhamento recorrente.",
                styles["Normal"],
            )
        )
    else:
        story.append(pdf_paragraph("Sem resultado supervisionado.", styles["Normal"]))

    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()


def main() -> None:
    ensure_data_storage()
    st.sidebar.title("AutoAnalista 2026")
    turbo_config = render_turbo_ai_controls()
    st.sidebar.caption(f"Criador do app: {CREATOR_SIGNATURE}")
    theme = st.sidebar.selectbox("Tema visual", options=["Escuro", "Claro"], index=0, key="theme_selector")
    apply_theme(theme)

    auth = login_panel()
    if auth is None:
        st.title("AutoAnalista de Dados 2026")
        st.info("Faca login para iniciar as analises.")
        st.caption(f"Assinatura do criador: {CREATOR_SIGNATURE}")
        st.stop()

    role = auth["role"]
    perms = ROLE_PERMISSIONS[role]
    area = st.sidebar.selectbox(
        "Area de negocio",
        options=["Geral", "Financeiro", "Vendas", "Operacoes", "RH"],
        index=0,
        key="business_area_selector",
    )
    render_database_controls(role)
    render_offline_install_controls()

    st.title("AutoAnalista de Dados 2026")
    st.caption("Analista profissional com qualidade, insights executivos, ML explicavel e governanca de historico.")
    st.caption(f"Assinatura do criador: {CREATOR_SIGNATURE}")

    uploaded = st.file_uploader("Envie sua planilha", type=SUPPORTED_TYPES, key="uploaded_sheet")
    if uploaded is None:
        st.subheader("Instalacao local/offline")
        st.info("Para rodar o AutoAnalista no proprio computador, use o botao `Baixar app offline` no menu lateral.")
        st.caption("O pacote inclui app, dependencias, inicializador Windows e instrucoes. O banco SQLite fica local no computador instalado.")
        hist = load_history()
        if hist:
            st.subheader("Historico recente")
            st.dataframe(pd.DataFrame(hist).tail(20).sort_values("timestamp", ascending=False), use_container_width=True)
        else:
            st.info("Aguardando arquivo para iniciar.")
        st.stop()

    prog = st.progress(0, text="Iniciando leitura do arquivo...")
    try:
        file_bytes = uploaded.getvalue()
        uploaded_hash_short = hashlib.sha256(file_bytes).hexdigest()[:12]
        workbook = read_workbook(file_bytes, uploaded.name)
    except Exception as exc:
        st.error("Nao foi possivel abrir o arquivo.")
        st.markdown("**Como corrigir:**")
        st.markdown("- Abra o arquivo no Excel/LibreOffice e salve como `.xlsx`.")
        st.markdown("- Verifique se o arquivo nao esta protegido por senha.")
        st.markdown("- Se for CSV, confirme separador e codificacao UTF-8.")
        st.code(str(exc))
        st.stop()

    prog.progress(20, text="Arquivo lido. Preparando abas...")
    if not workbook:
        st.warning("O arquivo nao contem abas/tabelas legiveis.")
        st.stop()

    if len(workbook) == 1:
        sheet = list(workbook.keys())[0]
        st.caption(f"Tabela detectada automaticamente: `{sheet}`")
    else:
        sheet = st.selectbox("Selecione a aba/tabela", options=list(workbook.keys()), key=f"sheet_selector_{uploaded_hash_short}")

    raw_df = workbook[sheet].copy()
    if raw_df.empty:
        st.warning("A aba selecionada esta vazia.")
        st.stop()

    prog.progress(35, text="Convertendo tipos e preparando perfil...")
    df_coerced, notes = coerce_data_types(raw_df)
    groups_coerced = get_column_groups(df_coerced)

    custom_rules = render_custom_rules_editor(
        df_coerced.columns.tolist(),
        perms["rules"],
        username=auth["username"],
        area=area,
    )
    filtered_df, active_filters = render_global_filters(df_coerced, groups_coerced)
    if filtered_df.empty:
        st.warning("Os filtros removeram todos os registros. Ajuste os filtros.")
        st.stop()
    treatment_options = render_data_treatment_controls()

    if "analysis_ready" not in st.session_state:
        st.session_state.analysis_ready = False
    if "analysis_key" not in st.session_state:
        st.session_state.analysis_key = None
    if "analysis_bundle" not in st.session_state:
        st.session_state.analysis_bundle = None
    if "ml_cache" not in st.session_state:
        st.session_state.ml_cache = {}
    if "export_cache" not in st.session_state:
        st.session_state.export_cache = {}

    file_hash = hashlib.sha256(file_bytes).hexdigest()
    analysis_payload = {
        "file_hash": file_hash,
        "file_name": uploaded.name,
        "sheet": sheet,
        "area": area,
        "filters": active_filters,
        "rules": custom_rules,
        "treatment": treatment_options,
        "rows": len(filtered_df),
    }
    analysis_key = hash_text(json.dumps(analysis_payload, ensure_ascii=False, sort_keys=True))

    if st.session_state.analysis_key != analysis_key:
        st.session_state.analysis_ready = False

    run_analysis = st.sidebar.button("Executar Analise", type="primary", use_container_width=True)
    st.sidebar.caption("A analise so recalcula quando voce clica no botao acima.")

    if run_analysis:
        prog.progress(50, text="Executando diagnostico de qualidade...")
        treated_df_full, treatment_report = apply_data_treatment(filtered_df, treatment_options)
        if treated_df_full.empty:
            st.error("A tratativa removeu todos os registros. Ajuste filtros ou parametros de outlier.")
            st.stop()

        analysis_df = treated_df_full
        sampled = False
        if len(analysis_df) > MAX_ANALYSIS_ROWS:
            analysis_df = analysis_df.sample(MAX_ANALYSIS_ROWS, random_state=SEED)
            sampled = True

        groups = get_column_groups(analysis_df)
        v_auto = run_auto_validations(analysis_df, groups)
        v_custom = run_custom_validations(analysis_df, custom_rules)
        violations = pd.concat([v_auto, v_custom], ignore_index=True) if len(v_auto) + len(v_custom) > 0 else pd.DataFrame()
        quality = compute_quality_report(analysis_df, groups, violations)
        top_corr = strongest_correlations(analysis_df, top_n=10)
        kpis = detect_kpis(analysis_df, groups)
        insights = generate_professional_insights(analysis_df, groups, quality, top_corr, kpis, treatment_report=treatment_report)
        missing_by_col = (analysis_df.isna().mean() * 100).sort_values(ascending=False)
        issue_catalog = build_issue_catalog(analysis_df, quality, violations, missing_by_col)
        column_quality = compute_column_quality(analysis_df, violations, quality["mixed_cols"])
        recommendations = build_action_plan(area, quality, issue_catalog)

        prog.progress(70, text="Registrando historico/versionamento...")
        history = load_history()
        fingerprint = f"{auth['username']}|{analysis_key}"
        if st.session_state.get("last_logged_fingerprint") != fingerprint:
            hist_entry = register_history(
                history,
                user=auth["username"],
                role=auth["role"],
                file_name=uploaded.name,
                sheet_name=sheet,
                rows=treatment_report["rows_after"],
                cols=analysis_df.shape[1],
                score=quality["score"],
            )
            st.session_state.last_logged_fingerprint = fingerprint
        else:
            same = [h for h in history if h.get("file_name") == uploaded.name and h.get("sheet_name") == sheet]
            hist_entry = same[-1] if same else {"version": 1}

        st.session_state.analysis_bundle = {
            "analysis_df": analysis_df,
            "treated_df_export": treated_df_full.head(MAX_EXPORT_ROWS).copy(),
            "filtered_df_export": filtered_df.head(MAX_EXPORT_ROWS).copy(),
            "sampled": sampled,
            "groups": groups,
            "violations": violations,
            "quality": quality,
            "top_corr": top_corr,
            "kpis": kpis,
            "insights": insights,
            "missing_by_col": missing_by_col,
            "issue_catalog": issue_catalog,
            "column_quality": column_quality,
            "recommendations": recommendations,
            "treatment_report": treatment_report,
            "version": int(hist_entry.get("version", 1)),
        }
        st.session_state.analysis_key = analysis_key
        st.session_state.analysis_ready = True
        prog.progress(100, text="Analise concluida.")

    if not st.session_state.analysis_ready or st.session_state.analysis_bundle is None or st.session_state.analysis_key != analysis_key:
        prog.empty()
        st.info("Configure filtros/regras e clique em `Executar Analise` para processar a base.")
        st.dataframe(filtered_df.head(40), use_container_width=True)
        st.stop()

    bundle = st.session_state.analysis_bundle
    analysis_df = bundle["analysis_df"]
    treated_df_export = bundle.get("treated_df_export", analysis_df.head(MAX_EXPORT_ROWS).copy())
    filtered_df_export = bundle.get("filtered_df_export", analysis_df.head(MAX_EXPORT_ROWS).copy())
    sampled = bundle["sampled"]
    groups = bundle["groups"]
    violations = bundle["violations"]
    quality = bundle["quality"]
    top_corr = bundle["top_corr"]
    kpis = bundle["kpis"]
    insights = bundle["insights"]
    missing_by_col = bundle["missing_by_col"]
    issue_catalog = bundle["issue_catalog"]
    column_quality = bundle["column_quality"]
    recommendations = bundle["recommendations"]
    treatment_report = bundle.get(
        "treatment_report",
        {
            "rows_before": len(analysis_df),
            "rows_after": len(analysis_df),
            "empty_rows_removed": 0,
            "duplicates_removed": 0,
            "outliers_removed": 0,
            "total_removed": 0,
            "warnings": [],
            "outlier_columns": [],
        },
    )
    version = bundle["version"]

    prog.empty()
    if sampled:
        st.info(
            f"Base grande detectada: analises avancadas executadas em amostra de {MAX_ANALYSIS_ROWS} linhas para performance."
        )
    if active_filters:
        st.caption("Filtros ativos: " + " | ".join(active_filters))
    if notes:
        st.caption("Conversoes automaticas: " + " | ".join(notes[:8]))
    st.caption(
        "Tratativa aplicada: "
        f"vazias={treatment_report['empty_rows_removed']}, "
        f"duplicadas={treatment_report['duplicates_removed']}, "
        f"outliers={treatment_report['outliers_removed']}."
    )

    target_options = ["(nao usar)"] + analysis_df.columns.tolist()
    model_options = ["Auto (comparar modelos)", "Random Forest", "Regressao Linear / Logistica"]
    ml_target_suggestions = suggest_ml_targets(analysis_df, groups)
    best_ml_target = ml_target_suggestions[0]["coluna"] if ml_target_suggestions else None
    target_state_key = f"predictive_target_value_{analysis_key}"
    target_sidebar_key = f"predictive_target_sidebar_{analysis_key}"
    target_tab_key = f"predictive_target_tab_{analysis_key}"
    model_state_key = f"predictive_model_value_{analysis_key}"
    model_sidebar_key = f"predictive_model_sidebar_{analysis_key}"
    model_tab_key = f"predictive_model_tab_{analysis_key}"

    if st.session_state.get(target_state_key) not in target_options:
        st.session_state[target_state_key] = "(nao usar)"
    if st.session_state.get(model_state_key) not in model_options:
        st.session_state[model_state_key] = "Auto (comparar modelos)"
    for widget_key in [target_sidebar_key, target_tab_key]:
        if st.session_state.get(widget_key) not in target_options:
            st.session_state[widget_key] = st.session_state[target_state_key]
        elif st.session_state.get(widget_key) != st.session_state[target_state_key]:
            st.session_state[widget_key] = st.session_state[target_state_key]
    for widget_key in [model_sidebar_key, model_tab_key]:
        if st.session_state.get(widget_key) not in model_options:
            st.session_state[widget_key] = st.session_state[model_state_key]
        elif st.session_state.get(widget_key) != st.session_state[model_state_key]:
            st.session_state[widget_key] = st.session_state[model_state_key]

    predictive_target_col = st.session_state[target_state_key]
    predictive_model_choice = st.session_state[model_state_key]
    train_predictive_model = False
    if perms["ml"]:
        st.sidebar.markdown("---")
        st.sidebar.subheader("Modelagem Preditiva")
        st.sidebar.caption("Usa os dados ja tratados, filtrados e saneados pela analise atual.")
        if best_ml_target:
            st.sidebar.caption(f"Sugestao: `{best_ml_target}` ({ml_target_suggestions[0]['tipo']}).")
            if st.sidebar.button("Usar alvo sugerido", use_container_width=True, key=f"use_suggested_target_sidebar_{analysis_key}"):
                st.session_state[target_state_key] = best_ml_target
                st.rerun()
        st.sidebar.selectbox(
            "Variavel alvo para prever",
            options=target_options,
            key=target_sidebar_key,
            on_change=sync_session_value,
            args=(target_sidebar_key, target_state_key, target_tab_key),
        )
        st.sidebar.selectbox(
            "Tipo de modelo",
            options=model_options,
            key=model_sidebar_key,
            on_change=sync_session_value,
            args=(model_sidebar_key, model_state_key, model_tab_key),
        )
        predictive_target_col = st.session_state[target_state_key]
        predictive_model_choice = st.session_state[model_state_key]
        train_predictive_model = st.sidebar.button(
            "Treinar Modelo",
            use_container_width=True,
            key=f"predictive_train_{analysis_key}",
            disabled=(predictive_target_col == "(nao usar)"),
        )

    tabs = st.tabs(["Resumo", "Dashboard", "Qualidade", "Insights", "ML", "Comparacao", "Relatorio"])

    with tabs[0]:
        st.subheader("Resumo Executivo")
        c1, c2, c3, c4, c5, c6 = st.columns(6)
        c1.metric("Linhas tratadas", f"{treatment_report['rows_after']:,}".replace(",", "."))
        c2.metric("Colunas", f"{analysis_df.shape[1]:,}".replace(",", "."))
        c3.metric("Ausencia", f"{quality['missing_ratio'] * 100:.2f}%")
        c4.metric("Score", f"{quality['score']:.1f}")
        c5.metric("Semaforo", quality_status(quality["score"]))
        c6.metric("Linhas removidas", int(treatment_report["total_removed"]))
        render_quality_overview(quality, groups, key_prefix=f"summary_quality_{analysis_key}")
        st.subheader("Preview")
        st.dataframe(analysis_df.head(80), use_container_width=True)

    with tabs[1]:
        render_management_dashboard(
            analysis_df,
            groups,
            quality,
            issue_catalog,
            kpis,
            treatment_report,
            violations,
            key_prefix=f"dash_{analysis_key}",
            turbo_config=turbo_config,
        )

    with tabs[2]:
        st.subheader("Qualidade de Dados Profissional")
        st.markdown("**Catalogo de problemas (prioridade)**")
        st.dataframe(issue_catalog.sort_values(["priority_rank", "affected_rows"], ascending=[True, False]), use_container_width=True)
        st.markdown("**Semaforo por coluna**")
        st.dataframe(column_quality, use_container_width=True)
        if (missing_by_col > 0).any():
            miss_df = missing_by_col.reset_index()
            miss_df.columns = ["coluna", "percentual_ausente"]
            st.plotly_chart(
                px.bar(miss_df.head(20), x="coluna", y="percentual_ausente", title="Percentual de ausencia por coluna (Top 20)"),
                use_container_width=True,
                key=f"quality_missing_{analysis_key}",
            )

    with tabs[3]:
        st.subheader("Insights e Visualizacoes")
        st.caption("Nesta aba, a leitura prioriza contagem de registros/IDs e segmentacoes. Somas aparecem apenas como analise complementar quando a coluna for metrica real.")
        insights_df, insight_filters = apply_inline_filters(
            analysis_df,
            groups,
            key_prefix=f"insights_filters_{analysis_key}",
            title="Filtros da aba Insights",
            expanded=True,
            max_categories=80,
        )
        if insight_filters:
            st.caption("Filtros ativos em Insights: " + " | ".join(insight_filters))
        if insights_df.empty:
            st.warning("Os filtros da aba Insights removeram todos os registros. Ajuste o recorte para recalcular os insights.")
        else:
            insights_groups = get_column_groups(insights_df)
            insights_quality = compute_quality_report(insights_df, insights_groups, pd.DataFrame())
            insights_missing_by_col = insights_df.isna().mean().mul(100).sort_values(ascending=False)
            insights_issue_catalog = build_issue_catalog(insights_df, insights_quality, pd.DataFrame(), insights_missing_by_col)
            insights_corr = strongest_correlations(insights_df, top_n=10)
            insights_kpis = detect_kpis(insights_df, insights_groups)
            insights_profile = detect_spreadsheet_profile(insights_df, insights_groups, insights_quality)
            insight_metric_numeric, insight_id_like = select_numeric_metric_columns(insights_df, insights_groups)
            insight_volume_col = pick_volume_id_column(insights_df, insight_id_like)
            s1, s2, s3, s4, s5 = st.columns(5)
            s1.metric("Registros no recorte", f"{len(insights_df):,}".replace(",", "."))
            s2.metric("Tipo detectado", clean_report_text(insights_profile.get("tipo", "-"), 26))
            s3.metric("Confianca", f"{float(insights_profile.get('confianca', 0.0)):.1f}%")
            s4.metric("IDs unicos", f"{insights_df[insight_volume_col].nunique(dropna=True):,}".replace(",", ".") if insight_volume_col else "-")
            s5.metric("Score qualidade", f"{insights_quality['score']:.1f}")
            st.info(clean_report_text(insights_profile.get("melhor_analise", "Analise recomendada indisponivel."), 420))
            insights_notes = generate_professional_insights(
                insights_df,
                insights_groups,
                insights_quality,
                insights_corr,
                insights_kpis,
                treatment_report=treatment_report,
            )
            st.markdown("**Insights claros do recorte filtrado**")
            for i in insights_notes:
                st.markdown(f"- {i}")

            render_turbo_ai_box(
                turbo_config=turbo_config,
                section="Insights",
                button_label="Gerar Insights Turbo com IA",
                context=build_ai_dataset_context(
                    insights_df,
                    insights_groups,
                    quality_report=insights_quality,
                    profile=insights_profile,
                    kpis=insights_kpis,
                    issue_catalog=insights_issue_catalog,
                    active_filters=insight_filters,
                ),
                instruction=(
                    "Transforme os dados filtrados em uma leitura analitica clara. "
                    "Separe: resumo do recorte, principais concentracoes, alertas, hipoteses de causa, "
                    "o que investigar primeiro e sugestao de graficos/indicadores. "
                    "Evite somar IDs/codigos; use contagem quando o campo representar unidade, chamado, ticket ou protocolo."
                ),
                cache_key=f"insights_{analysis_key}_{hash_text('|'.join(insight_filters))[:10]}",
            )

            st.markdown("**KPIs detectados automaticamente**")
            if insights_kpis:
                cols = st.columns(len(insights_kpis))
                for i, item in enumerate(insights_kpis):
                    delta = item["delta_text"] if item["delta_text"] else None
                    cols[i].metric(item["name"], format_kpi_metric(item), delta=delta)
                    if item["alert"]:
                        cols[i].caption(item["alert"])
            else:
                st.info("Sem KPIs suficientes para deteccao automatica.")

            render_count_based_insight_charts(
                insights_df,
                insights_groups,
                insights_profile,
                key_prefix=f"insights_count_{analysis_key}",
            )

            st.markdown("**Plano de acao recomendado para este recorte**")
            scoped_action_df = build_structured_action_plan(area, insights_quality, insights_issue_catalog)
            st.dataframe(scoped_action_df.head(6), use_container_width=True)

            with st.expander("Analises avancadas complementares", expanded=False):
                st.caption("Esta area pode usar distribuicao numerica, correlacao e graficos estatisticos. Para IDs/codigos, prefira sempre contagem.")
                render_auto_charts(insights_df, insights_groups, insights_corr, key_prefix=f"insights_{analysis_key}")
            if len(insights_corr) > 0:
                st.markdown("**Correlacoes de destaque**")
                st.dataframe(insights_corr, use_container_width=True)
                insight_metric_cols = insight_metric_numeric
                if len(insight_metric_cols) >= 2:
                    corr_matrix = insights_df[insight_metric_cols].corr(numeric_only=True)
                    st.plotly_chart(
                        px.imshow(corr_matrix, text_auto=True, aspect="auto", title="Matriz de correlacao"),
                        use_container_width=True,
                        key=f"insights_corr_{analysis_key}_{hash_text('|'.join(insight_filters))[:8]}",
                    )

    ml_unsup = None
    ml_sup = None
    anomalias_df = pd.DataFrame()
    with tabs[4]:
        if perms["ml"]:
            st.caption("Execucao sob demanda: rode ML apenas quando necessario.")
            st.markdown("**Modelagem Preditiva**")
            st.info(
                "Escolha abaixo a variavel alvo que deseja prever. O treinamento usa o dataframe principal ja tratado, "
                "com limpeza, duplicidades e outliers conforme a configuracao atual."
            )
            if ml_target_suggestions:
                st.markdown("**Sugestoes automaticas de alvo**")
                sug_df = pd.DataFrame(ml_target_suggestions)
                st.dataframe(sug_df, use_container_width=True, height=180)
                if best_ml_target and st.button("Aplicar melhor alvo sugerido", use_container_width=True, key=f"use_suggested_target_tab_{analysis_key}"):
                    st.session_state[target_state_key] = best_ml_target
                    st.rerun()
            else:
                st.caption("Nao encontrei uma variavel alvo claramente recomendada. Voce ainda pode selecionar manualmente.")
            ml_cfg1, ml_cfg2 = st.columns(2)
            ml_cfg1.selectbox(
                "Variavel alvo para prever",
                options=target_options,
                key=target_tab_key,
                on_change=sync_session_value,
                args=(target_tab_key, target_state_key, target_sidebar_key),
            )
            ml_cfg2.selectbox(
                "Tipo de modelo",
                options=model_options,
                key=model_tab_key,
                on_change=sync_session_value,
                args=(model_tab_key, model_state_key, model_sidebar_key),
            )
            predictive_target_col = st.session_state[target_state_key]
            predictive_model_choice = st.session_state[model_state_key]
            if predictive_target_col == "(nao usar)":
                st.warning("Selecione uma variavel alvo para liberar o treinamento supervisionado.")
            else:
                st.success(f"Pronto para treinar `{predictive_model_choice}` prevendo `{predictive_target_col}`.")
                target_warnings, target_strengths = diagnose_ml_target(analysis_df, predictive_target_col)
                d1, d2 = st.columns(2)
                with d1:
                    st.markdown("**Pontos fortes do alvo**")
                    if target_strengths:
                        for item in target_strengths:
                            st.caption(f"- {item}")
                    else:
                        st.caption("Sem ponto forte automatico destacado.")
                with d2:
                    st.markdown("**Alertas antes do treino**")
                    if target_warnings:
                        for item in target_warnings:
                            st.warning(item)
                    else:
                        st.caption("Nenhum alerta relevante detectado.")
            contamination = st.slider(
                "Taxa de anomalias (ML sem alvo)",
                min_value=0.01,
                max_value=0.20,
                value=0.05,
                step=0.01,
                key=f"contamination_{analysis_key}",
            )

            c1, c2 = st.columns(2)
            run_unsup = c1.button("Executar ML Sem Alvo", use_container_width=True, key=f"run_unsup_{analysis_key}")
            run_sup_tab = c2.button(
                "Treinar Modelo Agora",
                use_container_width=True,
                key=f"run_sup_{analysis_key}",
                disabled=(predictive_target_col == "(nao usar)"),
            )
            run_sup = train_predictive_model or run_sup_tab

            unsup_cache_key = f"{analysis_key}::unsup::{contamination:.2f}"
            sup_cache_key = f"{analysis_key}::sup::{predictive_target_col}::{predictive_model_choice}"
            model_key_suffix = re.sub(r"[^A-Za-z0-9_]+", "_", predictive_model_choice).strip("_").lower()

            if run_unsup:
                ml_unsup, anomalias_df = run_unsupervised_ml(
                    analysis_df,
                    groups["numeric"],
                    contamination=contamination,
                    key_prefix=f"unsup_{analysis_key}",
                )
                st.session_state.ml_cache[unsup_cache_key] = {"summary": ml_unsup, "anomalias": anomalias_df}
            elif unsup_cache_key in st.session_state.ml_cache:
                cached = st.session_state.ml_cache[unsup_cache_key]
                ml_unsup = cached["summary"]
                anomalias_df = cached["anomalias"]
                st.info("Resultado de ML sem alvo carregado do cache desta configuracao.")
                if ml_unsup:
                    st.metric("Anomalias detectadas (cache)", ml_unsup["anomalias"])

            if run_sup:
                ml_sup = run_supervised_ml(
                    analysis_df,
                    target_col=predictive_target_col,
                    model_choice=predictive_model_choice,
                    key_prefix=f"sup_{analysis_key}_{model_key_suffix}",
                )
                st.session_state.ml_cache[sup_cache_key] = {"summary": ml_sup}
            elif sup_cache_key in st.session_state.ml_cache:
                cached_sup = st.session_state.ml_cache[sup_cache_key]
                ml_sup = cached_sup["summary"]
                if ml_sup:
                    st.info("Resultado de ML supervisionado carregado do cache desta configuracao.")
                    st.json(public_ml_summary(ml_sup))

            if ml_sup:
                st.markdown("**Explicacao gerencial do modelo**")
                for item in ml_business_explanation(ml_sup):
                    st.markdown(f"- {item}")

                if ml_sup.get("model_comparison"):
                    st.markdown("**Comparacao dos modelos avaliados**")
                    st.dataframe(pd.DataFrame(ml_sup["model_comparison"]), use_container_width=True)

                ml_context = build_ai_dataset_context(
                    analysis_df,
                    groups,
                    quality_report=quality,
                    profile=detect_spreadsheet_profile(analysis_df, groups, quality),
                    kpis=kpis,
                    issue_catalog=issue_catalog,
                    active_filters=active_filters,
                )
                ml_context["resultado_ml"] = public_ml_summary(ml_sup)
                ml_context["resultado_ml_sem_alvo"] = ml_unsup or {}
                render_turbo_ai_box(
                    turbo_config=turbo_config,
                    section="Machine Learning",
                    button_label="Gerar explicacao Turbo do ML",
                    context=ml_context,
                    instruction=(
                        "Explique o resultado do Machine Learning para um gestor. "
                        "Inclua: objetivo do modelo, metrica principal, confiabilidade, variaveis mais relevantes, "
                        "riscos de uso, proximos testes recomendados e como aplicar o resultado na operacao. "
                        "Se a metrica for fraca, deixe isso claro sem mascarar o risco."
                    ),
                    cache_key=f"ml_{analysis_key}_{predictive_target_col}_{model_key_suffix}",
                )

                if ml_sup.get("problem_type") == "classification" and ml_sup.get("confusion_matrix"):
                    st.markdown("**Matriz de confusao salva do holdout**")
                    st.dataframe(pd.DataFrame(ml_sup["confusion_matrix"]), use_container_width=True)
                if ml_sup.get("problem_type") == "regression" and ml_sup.get("prediction_sample"):
                    st.markdown("**Amostra real vs previsto**")
                    sample_pred = pd.DataFrame(ml_sup["prediction_sample"])
                    st.dataframe(sample_pred.head(30), use_container_width=True)

                st.markdown("**Prever nova planilha com este modelo**")
                pred_file = st.file_uploader(
                    "Envie uma nova planilha com as mesmas colunas de entrada",
                    type=SUPPORTED_TYPES,
                    key=f"predict_new_file_{analysis_key}_{predictive_target_col}",
                )
                if pred_file is not None:
                    try:
                        pred_workbook = read_workbook(pred_file.getvalue(), pred_file.name)
                        pred_sheet = (
                            list(pred_workbook.keys())[0]
                            if len(pred_workbook) == 1
                            else st.selectbox(
                                "Selecione a aba para previsao",
                                options=list(pred_workbook.keys()),
                                key=f"predict_new_sheet_{analysis_key}_{predictive_target_col}",
                            )
                        )
                        pred_raw = pred_workbook[pred_sheet].copy()
                        pred_df, _ = coerce_data_types(pred_raw)
                        pipe = ml_sup.get("pipeline")
                        training_features = ml_sup.get("training_features", [])
                        if pipe is None or not training_features:
                            st.warning("Modelo carregado sem pipeline ativo. Treine novamente para prever uma nova planilha.")
                        else:
                            x_new = prepare_prediction_input(pred_df, training_features)
                            preds = pipe.predict(x_new)
                            result_pred = pred_df.copy()
                            result_pred[f"previsao_{ml_sup.get('target_col', 'alvo')}"] = preds
                            st.dataframe(result_pred.head(80), use_container_width=True)
                            pred_bytes = result_pred.head(MAX_EXPORT_ROWS).to_csv(index=False).encode("utf-8-sig")
                            st.download_button(
                                "Baixar previsoes CSV",
                                data=pred_bytes,
                                file_name=f"previsoes_{Path(pred_file.name).stem}.csv",
                                mime="text/csv",
                                use_container_width=True,
                                key=f"download_predictions_{analysis_key}_{predictive_target_col}",
                            )
                    except Exception as exc:
                        st.error("Nao foi possivel gerar previsoes para a nova planilha.")
                        st.caption(str(exc)[:500])
        else:
            st.info("Perfil viewer sem permissao para executar ML.")

    with tabs[5]:
        st.subheader("Comparacao entre Planilhas")
        st.caption("Compare a planilha atual com outro arquivo ou com uma planilha salva no banco SQLite local.")
        compare_source = st.radio(
            "Origem da base para comparacao",
            options=["Enviar outro arquivo", "Usar planilha salva no SQLite"],
            horizontal=True,
            key=f"compare_source_{analysis_key}",
        )
        compare_df = None
        compare_label = "Comparacao"

        if compare_source == "Enviar outro arquivo":
            compare_file = st.file_uploader(
                "Envie a planilha de comparacao",
                type=SUPPORTED_TYPES,
                key=f"compare_upload_{analysis_key}",
            )
            if compare_file is not None:
                try:
                    compare_workbook = read_workbook(compare_file.getvalue(), compare_file.name)
                    compare_sheet = (
                        list(compare_workbook.keys())[0]
                        if len(compare_workbook) == 1
                        else st.selectbox(
                            "Selecione a aba da planilha de comparacao",
                            options=list(compare_workbook.keys()),
                            key=f"compare_sheet_{analysis_key}",
                        )
                    )
                    compare_raw = compare_workbook[compare_sheet].copy()
                    compare_df, _ = coerce_data_types(compare_raw)
                    compare_label = Path(compare_file.name).stem
                except Exception as exc:
                    st.error("Nao foi possivel ler a planilha de comparacao.")
                    st.caption(str(exc)[:500])
        else:
            stored_df = list_stored_sheets()
            if len(stored_df) == 0:
                st.info("Nenhuma planilha salva no SQLite ainda. Salve uma planilha na aba Relatorio primeiro.")
            else:
                stored_df["rotulo"] = stored_df.apply(
                    lambda r: f"{r['timestamp']} | {r['file_name']} | {r['sheet_name']} | {r['rows']} linhas",
                    axis=1,
                )
                selected_stored = st.selectbox(
                    "Planilha salva para comparar",
                    options=stored_df["rotulo"].tolist(),
                    key=f"compare_stored_{analysis_key}",
                )
                selected_id = stored_df.loc[stored_df["rotulo"] == selected_stored, "id"].iloc[0]
                loaded = load_stored_sheet(selected_id)
                if loaded is not None:
                    compare_df, _ = coerce_data_types(loaded)
                    compare_label = selected_stored.split("|")[1].strip() if "|" in selected_stored else "SQLite"

        if compare_df is not None:
            comp = compare_dataframes(analysis_df, compare_df, current_name="Atual", other_name=compare_label)
            st.markdown("### Resumo Executivo")
            vol = comp["volume_df"].iloc[0] if len(comp["volume_df"]) else {}
            score_delta = float(comp["summary_df"].loc[comp["summary_df"]["indicador"] == "Score qualidade", "delta"].iloc[0])
            c1, c2, c3, c4 = st.columns(4)
            c1.metric("Volume atual", f"{int(vol.get('Atual', len(analysis_df))):,}".replace(",", "."), delta=f"{int(vol.get('delta', 0)):,}".replace(",", "."))
            c2.metric("Delta volume %", f"{float(vol.get('delta_pct', 0) or 0):.1f}%")
            c3.metric("Delta qualidade", f"{score_delta:.1f} pts")
            c4.metric("Alertas gerenciais", int(len(comp["action_df"])))

            for item in comp["insights"]:
                st.markdown(f"- {item}")

            st.markdown("### Mudancas Relevantes")
            delta_df = comp.get("category_delta_df", pd.DataFrame())
            ranking_df = comp.get("ranking_df", pd.DataFrame())
            m1, m2 = st.columns(2)
            with m1:
                st.markdown("**Maiores altas e quedas por categoria**")
                if len(delta_df) > 0:
                    st.dataframe(delta_df.head(15), use_container_width=True)
                else:
                    st.caption("Sem categorias comuns suficientes para detectar mudancas de perfil.")
            with m2:
                st.markdown("**Entradas, saidas e mudancas no ranking**")
                if len(ranking_df) > 0:
                    st.dataframe(ranking_df.head(15), use_container_width=True)
                else:
                    st.caption("Sem ranking categorico comparavel.")

            st.markdown("### Graficos Comparativos")
            chart_df = comp["summary_df"].copy()
            chart_df["delta_abs"] = pd.to_numeric(chart_df["delta"], errors="coerce").abs()
            g1, g2 = st.columns(2)
            fig_comp = px.bar(chart_df, x="indicador", y="delta", text="delta", title="Deltas executivos")
            fig_comp = apply_analytics_chart_layout(fig_comp, height=360)
            g1.plotly_chart(fig_comp, use_container_width=True, key=f"comparison_delta_{analysis_key}")

            if len(delta_df) > 0:
                top_delta = delta_df.head(12).copy()
                top_delta["rotulo"] = top_delta["coluna"].astype(str) + " | " + top_delta["categoria"].astype(str)
                fig_cat_delta = px.bar(
                    top_delta,
                    x="delta_qtd",
                    y="rotulo",
                    orientation="h",
                    text="delta_qtd",
                    title="Maiores mudancas de perfil",
                    color="delta_participacao_pct",
                    color_continuous_scale=["#DC2626", "#F8FAFC", "#16A34A"],
                )
                fig_cat_delta = apply_analytics_chart_layout(fig_cat_delta, height=420)
                g2.plotly_chart(fig_cat_delta, use_container_width=True, key=f"comparison_cat_delta_{analysis_key}")
            else:
                g2.info("Sem deltas categoricos para graficos.")

            if len(comp["numeric_df"]) > 0:
                st.markdown("**Variacao de metricas numericas reais**")
                num_chart = comp["numeric_df"].head(12).copy()
                fig_num = px.bar(
                    num_chart,
                    x="coluna",
                    y="delta_soma",
                    text="delta_soma",
                    title="Delta de soma por metrica",
                    color="delta_soma",
                    color_continuous_scale=["#DC2626", "#F8FAFC", "#16A34A"],
                )
                fig_num = apply_analytics_chart_layout(fig_num, height=380)
                st.plotly_chart(fig_num, use_container_width=True, key=f"comparison_numeric_delta_{analysis_key}")
                st.dataframe(comp["numeric_df"].head(15), use_container_width=True)

            st.markdown("### Plano de Acao")
            st.dataframe(comp["action_df"], use_container_width=True)

            comp_excel = generate_comparison_excel(comp)
            st.download_button(
                "Baixar Excel da Comparacao",
                data=comp_excel,
                file_name=f"comparacao_{Path(uploaded.name).stem}_v{version}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True,
            )

    with tabs[6]:
        st.subheader("Relatorios e Historico")
        st.markdown(f"Versao atual do relatorio: **v{version}**")
        if perms["export"]:
            export_key = f"{analysis_key}::exports"
            build_exports = st.button("Preparar Exportacoes", type="primary", use_container_width=True, key=f"prep_export_{analysis_key}")
            if build_exports:
                with st.spinner("Gerando pacotes de exportacao..."):
                    dashboard_key_prefix = f"dash_{analysis_key}"
                    dashboard_filters_for_report = st.session_state.get(f"{dashboard_key_prefix}_active_filters", [])
                    dashboard_config_for_report = st.session_state.get(f"{dashboard_key_prefix}_config", {})
                    dashboard_report_df = st.session_state.get(f"{dashboard_key_prefix}_filtered_df")
                    if not isinstance(dashboard_report_df, pd.DataFrame) or dashboard_report_df.empty:
                        dashboard_report_df = analysis_df
                    dashboard_report_groups = get_column_groups(dashboard_report_df)
                    dashboard_report_kpis = detect_kpis(dashboard_report_df, dashboard_report_groups)
                    dashboard_top_corr = strongest_correlations(dashboard_report_df, top_n=10)
                    report_filters = [f"Global - {item}" for item in active_filters]
                    report_filters.extend([f"Dashboard - {item}" for item in dashboard_filters_for_report])
                    dashboard_bundle = build_dashboard_export_bundle(
                        dashboard_report_df,
                        dashboard_report_groups,
                        quality,
                        issue_catalog,
                        dashboard_report_kpis,
                        active_filters=dashboard_filters_for_report,
                        dashboard_config=dashboard_config_for_report,
                    )
                    excel_bytes = generate_excel_export(
                        treated_df_export,
                        quality,
                        column_quality,
                        issue_catalog,
                        violations,
                        anomalias_df,
                        insights,
                        recommendations,
                        raw_df=filtered_df_export,
                        treatment_report=treatment_report,
                        dashboard_bundle=dashboard_bundle,
                    )
                    csv_bytes = (
                        violations.head(MAX_EXPORT_ROWS).to_csv(index=False).encode("utf-8")
                        if len(violations) > 0
                        else "".encode("utf-8")
                    )
                    pdf_bytes = build_pdf_report(
                        source_name=uploaded.name,
                        sheet_name=sheet,
                        quality_report=quality,
                        insights=insights,
                        recommendations=recommendations,
                        top_corr=dashboard_top_corr if len(dashboard_top_corr) > 0 else top_corr,
                        issue_catalog=issue_catalog,
                        ml_unsup=ml_unsup,
                        ml_sup=ml_sup,
                        version=version,
                        dashboard_bundle=dashboard_bundle,
                        treatment_report=treatment_report,
                        active_filters=report_filters,
                        area=area,
                    )
                    try:
                        pptx_bytes = generate_powerpoint_report(
                            source_name=uploaded.name,
                            sheet_name=sheet,
                            quality_report=quality,
                            insights=insights,
                            issue_catalog=issue_catalog,
                            dashboard_bundle=dashboard_bundle,
                            treatment_report=treatment_report,
                            ml_sup=ml_sup,
                            version=version,
                        )
                    except Exception as exc:
                        pptx_bytes = b""
                        st.warning(f"PowerPoint nao gerado nesta execucao: {str(exc)[:180]}")
                    offline_bytes = generate_offline_package()
                    st.session_state.export_cache[export_key] = {
                        "excel_bytes": excel_bytes,
                        "csv_bytes": csv_bytes,
                        "pdf_bytes": pdf_bytes,
                        "pptx_bytes": pptx_bytes,
                        "offline_bytes": offline_bytes,
                    }
                    st.success("Exportacoes prontas para download.")

            export_bundle = st.session_state.export_cache.get(export_key)
            safe_sheet = "".join(ch if ch.isalnum() or ch in ["_", "-"] else "_" for ch in sheet)[:60]
            base_name = Path(uploaded.name).stem
            st.download_button(
                "Baixar Relatorio PDF",
                data=export_bundle["pdf_bytes"] if export_bundle else b"",
                file_name=f"relatorio_{base_name}_{safe_sheet}_v{version}.pdf",
                mime="application/pdf",
                use_container_width=True,
                disabled=(export_bundle is None),
            )
            st.download_button(
                "Baixar Pacote Excel de Evidencias",
                data=export_bundle["excel_bytes"] if export_bundle else b"",
                file_name=f"evidencias_{base_name}_{safe_sheet}_v{version}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True,
                disabled=(export_bundle is None),
            )
            st.download_button(
                "Baixar Apresentacao PowerPoint Executiva",
                data=export_bundle.get("pptx_bytes", b"") if export_bundle else b"",
                file_name=f"apresentacao_{base_name}_{safe_sheet}_v{version}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                disabled=(export_bundle is None or not export_bundle.get("pptx_bytes")),
            )
            st.download_button(
                "Baixar CSV de Violacoes",
                data=export_bundle["csv_bytes"] if export_bundle else b"",
                file_name=f"violacoes_{base_name}_{safe_sheet}_v{version}.csv",
                mime="text/csv",
                use_container_width=True,
                disabled=(export_bundle is None or len(violations) == 0),
            )
            st.download_button(
                "Baixar Pacote para Instalar Offline",
                data=export_bundle.get("offline_bytes", b"") if export_bundle else b"",
                file_name="autoanalista_2026_offline.zip",
                mime="application/zip",
                use_container_width=True,
                disabled=(export_bundle is None),
            )
        else:
            st.info("Perfil viewer sem permissao de exportacao.")

        st.markdown("**Banco de dados local**")
        db_c1, db_c2, db_c3 = st.columns([1.2, 1, 1])
        db_c1.info(f"SQLite: {database_status_label()}")
        if db_c2.button("Instalar/Inicializar SQLite", use_container_width=True, key=f"install_sqlite_tab_{analysis_key}"):
            migrated = init_local_database(migrate_json=True)
            st.success(f"Banco SQLite pronto. Historico migrado: {migrated} registro(s).")
        if db_c1.button("Salvar planilha tratada no SQLite", use_container_width=True, key=f"save_sheet_sqlite_{analysis_key}"):
            if not local_database_installed():
                init_local_database(migrate_json=True)
            stored_id = save_sheet_to_local_db(
                treated_df_export,
                user=auth["username"],
                file_name=uploaded.name,
                sheet_name=sheet,
                quality_score=quality["score"],
                version=version,
                metadata={
                    "area": area,
                    "filters": active_filters,
                    "rows_full_analysis": len(analysis_df),
                    "saved_limit": MAX_EXPORT_ROWS,
                },
            )
            st.success(f"Planilha salva no banco local SQLite. ID: {stored_id}")
        if role == "admin":
            clear_tab_confirm = db_c3.checkbox("Confirmar limpeza", key=f"confirm_clear_history_tab_{analysis_key}")
            if st.button(
                "Limpar historico de analises",
                use_container_width=True,
                key=f"clear_history_tab_{analysis_key}",
                disabled=not clear_tab_confirm,
            ):
                clear_analysis_history()
                st.session_state.analysis_ready = False
                st.session_state.analysis_bundle = None
                st.session_state.analysis_key = None
                st.session_state.export_cache = {}
                st.success("Historico de analises limpo.")
        else:
            db_c3.caption("Limpeza apenas para admin.")

        st.markdown("**Historico de analises**")
        hist_df = pd.DataFrame(load_history())
        if len(hist_df) > 0:
            hist_view = hist_df.sort_values("timestamp", ascending=False).head(50)
            st.dataframe(hist_view, use_container_width=True)
            hist_plot = hist_df.copy()
            hist_plot["timestamp"] = pd.to_datetime(hist_plot["timestamp"], errors="coerce")
            hist_plot = hist_plot.dropna(subset=["timestamp"]).sort_values("timestamp")
            if len(hist_plot) >= 2:
                h1, h2 = st.columns(2)
                h1.plotly_chart(
                    apply_analytics_chart_layout(
                        px.line(hist_plot, x="timestamp", y="quality_score", markers=True, title="Evolucao do score de qualidade"),
                        height=340,
                    ),
                    use_container_width=True,
                    key=f"history_quality_{analysis_key}",
                )
                h2.plotly_chart(
                    apply_analytics_chart_layout(
                        px.bar(hist_plot.tail(20), x="timestamp", y="rows", title="Volume analisado por execucao"),
                        height=340,
                    ),
                    use_container_width=True,
                    key=f"history_rows_{analysis_key}",
                )
                latest_score = float(hist_plot["quality_score"].iloc[-1])
                prev_score = float(hist_plot["quality_score"].iloc[-2])
                delta_score = latest_score - prev_score
                st.caption(f"Historico inteligente: variacao do ultimo score vs execucao anterior = {delta_score:.2f} ponto(s).")
        else:
            st.caption("Sem registros no historico ainda.")

        st.markdown("**Planilhas salvas no SQLite**")
        stored_df = list_stored_sheets()
        if len(stored_df) > 0:
            st.dataframe(stored_df.head(50), use_container_width=True)
            if role == "admin":
                confirm_clear_sheets = st.checkbox("Confirmar limpeza das planilhas salvas", key=f"confirm_clear_stored_sheets_{analysis_key}")
                if st.button(
                    "Limpar planilhas salvas no SQLite",
                    use_container_width=True,
                    key=f"clear_stored_sheets_{analysis_key}",
                    disabled=not confirm_clear_sheets,
                ):
                    clear_stored_sheets()
                    st.success("Planilhas salvas removidas do SQLite.")
        else:
            st.caption("Nenhuma planilha salva no banco local ainda.")


if __name__ == "__main__":
    main()
